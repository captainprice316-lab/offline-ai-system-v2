"""
generate_finetune_pptx.py  --  VANI Fine-Tuning 22-slide PPTX
"""
import sys
import matplotlib
matplotlib.use("Agg")
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Import chart generators from the PDF report script
sys.path.insert(0, str(Path(__file__).parent))
from generate_report_pdf import (
    chart_dataset_sizes, chart_summary_bar, chart_wer_all,
    chart_wer_per_lang, LANG_ORDER, LANG_META,
)
# Shared hero charts (identical figures in the PDF + PPTX)
import report_charts

# ── Colour palette — light professional theme ──────────────────────────────────
C_BG     = RGBColor(0xFF, 0xFF, 0xFF)   # white background
C_CARD   = RGBColor(0xF0, 0xF4, 0xFF)   # very light indigo card
C_CARD2  = RGBColor(0xF8, 0xF9, 0xFA)   # off-white alternate card
C_NAVY   = RGBColor(0x1A, 0x23, 0x7E)   # deep navy (headings, top band)
C_TEXT   = RGBColor(0x21, 0x21, 0x21)   # near-black body text
C_SUB    = RGBColor(0x54, 0x6E, 0x7A)   # blue-grey secondary text
C_TEAL   = RGBColor(0x00, 0x79, 0x6B)   # professional teal accent
C_TEAL_L = RGBColor(0xE0, 0xF2, 0xF1)   # very light teal fill
C_BLUE   = RGBColor(0x15, 0x65, 0xC0)   # medium blue accent
C_GOLD   = RGBColor(0xF5, 0x7F, 0x17)   # amber (darker for light bg)
C_GREEN  = RGBColor(0x2E, 0x7D, 0x32)   # deep green for good results
C_RED    = RGBColor(0xC6, 0x28, 0x28)   # deep red for warnings
C_PURPLE = RGBColor(0x6A, 0x1B, 0x9A)   # purple accent
C_BORDER = RGBColor(0xC5, 0xCA, 0xE9)   # light indigo border

OUT_PATH = Path(__file__).parent / "docs" / "VANI_Finetune_Presentation_v6.pptx"

FONT = "Times New Roman"   # global presentation typeface

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0.75)):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, left, top, width, height,
          font_size=18, bold=False, color=C_TEXT,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name    = FONT
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return tb

def _multiline(slide, text, left, top, width, height,
               font_size=9, bold=False, color=C_TEXT, align=PP_ALIGN.CENTER):
    """Textbox that renders each '\\n'-separated line as its own paragraph."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def top_band(slide, color=C_NAVY):
    """Thin decorative band at very top — does NOT overlap headings."""
    box(slide, 0, 0, 13.33, 0.10, fill_color=color)

def hline(slide, top, color=C_TEAL, width_in=12.5, left=0.4):
    box(slide, left, top, width_in, 0.03, fill_color=color)

def slide_num(slide, n):
    txbox(slide, str(n), 12.6, 7.1, 0.6, 0.32,
          font_size=11, color=C_SUB, align=PP_ALIGN.RIGHT)

def section_tag(slide, text, left=0.4, top=0.15):
    box(slide, left, top, 2.5, 0.26, fill_color=C_NAVY)
    txbox(slide, text.upper(), left+0.08, top+0.03, 2.4, 0.22,
          font_size=8.5, bold=True, color=RGBColor(0xFF,0xFF,0xFF))

def slide_header(slide, title, n, tag=None, top=0.55):
    """Standard slide header: navy top band + title + teal underline."""
    top_band(slide)
    if tag:
        section_tag(slide, tag)
    txbox(slide, title, 0.4, top, 12.5, 0.65,
          font_size=27, bold=True, color=C_NAVY)
    hline(slide, top + 0.70)
    slide_num(slide, n)

def card(slide, left, top, width, height, accent_color=None):
    box(slide, left, top, width, height, fill_color=C_CARD,
        line_color=C_BORDER, line_width=Pt(0.75))
    if accent_color:
        box(slide, left, top, 0.06, height, fill_color=accent_color)

def bullet_block(slide, items, left, top, width, font_size=13,
                 color=C_SUB, spacing=0.38, bullet="▸"):
    y = top
    for item in items:
        txbox(slide, f"{bullet}  {item}", left, y, width, 0.36,
              font_size=font_size, color=color)
        y += spacing

# ── Slide 1 — Title ───────────────────────────────────────────────────────────

def slide_01_title(prs):
    s = blank_slide(prs)
    bg(s)
    # Navy header block
    box(s, 0, 0, 13.33, 3.2, fill_color=C_NAVY)
    box(s, 0, 3.18, 13.33, 0.07, fill_color=C_TEAL)

    txbox(s, "VANI", 1.0, 0.10, 11.33, 1.55,
          font_size=90, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    txbox(s, "Voice Analysis & Neural Intelligence System",
          1.0, 1.68, 11.33, 0.5,
          font_size=17, color=RGBColor(0xC5, 0xCA, 0xE9), align=PP_ALIGN.CENTER)
    txbox(s, "Whisper ASR Fine-Tuning Report",
          1.0, 2.28, 11.33, 0.55,
          font_size=20, bold=True, color=RGBColor(0xFF, 0xD5, 0x4F), align=PP_ALIGN.CENTER)

    # Light content area
    txbox(s, "LoRA Domain Adaptation for 7 Border-Region Languages",
          1.0, 3.45, 11.33, 0.48,
          font_size=16, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    hline(s, 4.03, width_in=5.5, left=3.92, color=C_TEAL)
    txbox(s, "M.Tech Research Project  ·  IIT Indore  ·  2026",
          1.0, 4.12, 11.33, 0.38,
          font_size=13, italic=True, color=C_SUB, align=PP_ALIGN.CENTER)
    txbox(s, "Hardware: NVIDIA RTX 5060 8 GB  ·  Windows 11  ·  CUDA  ·  faster-whisper CT2 int8",
          1.0, 4.58, 11.33, 0.32,
          font_size=10, color=C_SUB, align=PP_ALIGN.CENTER)

    # Bottom info strip
    box(s, 0, 7.1, 13.33, 0.40, fill_color=C_CARD)
    txbox(s, "7 Languages  ·  LoRA r=8/16  ·  CTranslate2 int8  ·  Fully Offline",
          0.5, 7.14, 12.33, 0.30,
          font_size=11, color=C_NAVY, align=PP_ALIGN.CENTER)
    slide_num(s, 1)

# ── Slide 2 — Motivation ──────────────────────────────────────────────────────

def slide_02_motivation(prs):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Why Fine-Tune Whisper?", 2, "Background")

    # Left panel — problem
    card(s, 0.4, 1.42, 5.85, 5.0, accent_color=C_RED)
    box(s, 0.4, 1.42, 5.85, 0.38, fill_color=C_RED)
    txbox(s, "⚠  The Problem", 0.55, 1.47, 5.5, 0.32,
          font_size=13, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    issues = [
        "Baseline Whisper large-v3 WER on\nPunjabi (pa): ~78%",
        "Pashto baseline WER: ~90%",
        "Nepali baseline WER: ~89%",
        "Kashmiri: no Whisper support at all\n(baseline ~97% WER)",
        "Radio intercept audio is noisy\n& domain-specific",
        "Generic model lacks border-language\nacoustic knowledge",
    ]
    bullet_block(s, issues, 0.58, 1.95, 5.5, font_size=12,
                 color=C_TEXT, spacing=0.68, bullet="✗")

    # Right panel — solution
    card(s, 6.75, 1.42, 5.85, 5.0, accent_color=C_TEAL)
    box(s, 6.75, 1.42, 5.85, 0.38, fill_color=C_TEAL)
    txbox(s, "✔  The Solution — LoRA Fine-Tuning", 6.90, 1.47, 5.5, 0.32,
          font_size=13, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    solutions = [
        "Train only 0.25% of parameters\n(LoRA r=8/16, α=16/32)",
        "Use FLEURS speech corpus + AI4Bharat\nIndicVoices-R for Punjabi & Nepali",
        "Base model weights stay frozen\n— no catastrophic forgetting",
        "Export to CTranslate2 int8\nfor fast offline inference",
        "Plug directly into VANI\n10-stage pipeline",
        "Fully offline — no internet after\ninitial model download",
    ]
    bullet_block(s, solutions, 6.90, 1.95, 5.5, font_size=12,
                 color=C_TEXT, spacing=0.68, bullet="✔")

# ── Slide 3 — Methodology ─────────────────────────────────────────────────────

def slide_03_methodology(prs):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "LoRA Fine-Tuning Methodology", 3, "Methodology")

    # Pipeline flow
    steps = [
        ("1", "FLEURS /\nIndicVoices-R", C_TEAL),
        ("2", "Whisper\nProcessor", C_TEAL),
        ("3", "LoRA Adapter\nr=8, q+v proj", C_GOLD),
        ("4", "Seq2Seq\nTrainer", C_GOLD),
        ("5", "Merge &\nCT2 Convert", C_GREEN),
        ("6", "VANI\nDeploy", C_GREEN),
    ]
    x = 0.3
    for i, (num, label, col) in enumerate(steps):
        box(s, x, 1.55, 1.9, 1.25,
            fill_color=C_CARD, line_color=col, line_width=Pt(1.2))
        box(s, x, 1.55, 1.9, 0.06, fill_color=col)
        txbox(s, num, x, 1.62, 1.9, 0.45,
              font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
        txbox(s, label, x, 2.1, 1.9, 0.65,
              font_size=10, color=C_SUB, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            txbox(s, "→", x+1.9, 2.05, 0.35, 0.4,
                  font_size=20, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
        x += 2.22

    # Left config table
    txbox(s, "LoRA Configuration", 0.4, 3.05, 6.0, 0.38,
          font_size=14, bold=True, color=C_NAVY)
    params = [
        ("LoRA Rank (r)",     "8 (v1/v2)  /  16 (PA v3)"),
        ("LoRA Alpha",        "16 (v1/v2)  /  32 (PA v3)"),
        ("Target Modules",    "q_proj, v_proj"),
        ("Trainable Params",  "~3.9M (r=8)  /  ~7.9M (r=16)"),
        ("Learning Rate",     "5×10⁻⁵ with linear warmup"),
        ("Grad Clip",         "max_grad_norm = 0.5"),
        ("Quantization",      "int8 via CTranslate2"),
    ]
    y = 3.48
    for k, v in params:
        bg_col = C_CARD if (params.index((k,v))) % 2 == 0 else C_CARD2
        box(s, 0.4, y, 6.0, 0.36, fill_color=bg_col)
        txbox(s, k, 0.52, y+0.06, 2.5, 0.26, font_size=11, color=C_SUB)
        txbox(s, v, 3.0, y+0.06, 3.3, 0.26, font_size=11, bold=True, color=C_TEXT)
        y += 0.38

    # Right training notes
    txbox(s, "Training Notes", 6.9, 3.05, 6.0, 0.38,
          font_size=14, bold=True, color=C_NAVY)
    notes = [
        "Batch size 2, effective batch = 2",
        "Eval metric: WER (jiwer)",
        "Best checkpoint: load_best_model_at_end",
        "fp16 mixed precision training",
        "PA v3: 20k IndicVoices-R + FLEURS",
        "CT2 tokenizer.json fix (large-v3 vocab shift)",
        "preprocessor_config.json: feature_size=128",
        "faster-whisper patched for ks language code",
    ]
    y2 = 3.48
    for note in notes:
        box(s, 6.9, y2, 6.0, 0.36,
            fill_color=C_CARD if notes.index(note) % 2 == 0 else C_CARD2)
        txbox(s, f"▸  {note}", 7.02, y2+0.06, 5.8, 0.26, font_size=11, color=C_SUB)
        y2 += 0.38

# ── Slide 4 — Results Overview ────────────────────────────────────────────────

def slide_04_results_overview(prs):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Fine-Tuning Results — All 7 Languages", 4, "Results")

    cols_x = [0.35, 2.05, 3.55, 5.55, 7.15, 8.75, 9.9, 11.2]
    headers = ["Language", "Script", "Base\nModel", "Training Data\n(samples)",
               "Train\nSteps", "Baseline\nWER (%)", "Best\nWER (%)", "WER Reduction\n(pp)"]
    y = 1.45
    box(s, 0.3, y, 12.75, 0.48, fill_color=C_NAVY)
    for x, h in zip(cols_x, headers):
        _multiline(s, h, x, y+0.04, 1.65, 0.40,
                   font_size=9.5, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.LEFT)

    # Baseline = true openai/whisper-large-v3 (not turbo); Best WER = held-out n=100
    # FLEURS test (corrected 2026-07-11). Mandarin regressed vs baseline (red).
    rows = [
        ("Punjabi  (pa) v3", "Gurmukhi",   "large-v3",       "FLEURS+IV-R (21,923)", "4000",  "77.6%",  "57.39%", "−20 pp", C_GREEN),
        ("Pashto   (ps)",    "Nastaliq",   "medium-pashto",  "FLEURS ps_af (2,082)", "2000",  "89.8%",  "38.55%", "−51 pp", C_GREEN),
        ("Urdu     (ur)",    "Nastaliq",   "large-v3",       "FLEURS ur_pk (2,109)", "1000",  "21.2%",  "19.82%", "−1.4pp", C_GREEN),
        ("Nepali   (ne) v2", "Devanagari", "large-v3",       "FLEURS+IV-R (13,332)", "3000",  "88.9%",  "50.92%", "−38 pp", C_GOLD),
        ("Mandarin (zh)",    "Simplified", "large-v3",       "FLEURS cmn (3,246)",   "400†",  "11.0%",  "14.22%‡","+3.2 pp",C_RED),
        ("Hindi    (hi)",    "Devanagari", "large-v3",       "FLEURS hi_in (2,120)", "600",   "26.3%",  "19.78%", "−6.6 pp",C_GREEN),
        ("Kashmiri (ks)",    "Nastaliq",   "large-v3+<|ks|>","IndicVoices-R (20k)",  "2400§", "96.9%",  "74.02%", "−23 pp", C_GOLD),
    ]
    y += 0.52
    for i, row in enumerate(rows):
        lang, script, base, data, steps, bwer, fwer, imp, col = row
        bg_col = C_CARD if i % 2 == 0 else C_CARD2
        box(s, 0.3, y, 12.75, 0.36, fill_color=bg_col)
        vals = [lang, script, base, data, steps, bwer, fwer, imp]
        for j, (x, v) in enumerate(zip(cols_x, vals)):
            c = col if j in (6, 7) else (C_TEXT if j == 0 else C_SUB)
            txbox(s, v, x, y+0.05, 1.65, 0.26,
                  font_size=9.5, bold=(j in (0, 6, 7)), color=c)
        y += 0.37

    txbox(s, "Baseline = true openai/whisper-large-v3 · Best WER = held-out n=100 FLEURS test  "
             "† ZH diverged at step ~820; ckpt-400 best.  ‡ ZH fine-tuning REGRESSED vs baseline "
             "(prior 100%→16% was a whitespace-scoring artefact); ZH is served by SeamlessM4T.  "
             "§ KS custom <|ks|> token ID 51866; SeamlessM4T has no Kashmiri.",
          0.35, 7.05, 12.5, 0.32, font_size=8.5, italic=True, color=C_SUB)

# ── Slide 5 — WER Comparison ──────────────────────────────────────────────────

def slide_05_wer_chart(prs, n=8):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Fine-Tuned Whisper vs SeamlessM4T — ASR WER", n, "Backend Selection")

    # Held-out n=100 FLEURS test WER (CER in parentheses on the labels). The operational
    # question is FT Whisper vs zero-shot SeamlessM4T; the lower bar (green) is deployed.
    langs  = ["Mandarin (zh)", "Urdu (ur)",   "Hindi (hi)",   "Nepali (ne)",  "Punjabi (pa)", "Pashto (ps)"]
    ftwer  = [14.22,           19.82,          19.78,          50.92,           57.39,           38.55]
    smwer  = [11.69,           16.90,          15.44,          28.46,           19.77,           44.40]
    ftcer  = [14.22,           7.29,           7.46,           18.83,           32.52,           17.65]
    smcer  = [11.69,           7.00,           9.12,           11.22,           9.97,            22.92]

    scale  = 4.8 / 60.0
    bar_left = 2.65

    txbox(s, "■ FT Whisper", bar_left + 0.05, 1.43, 2.0, 0.28,
          font_size=10, color=RGBColor(0x78, 0x90, 0xA0), italic=True)
    txbox(s, "■ SeamlessM4T (zero-shot)", bar_left + 2.2, 1.43, 2.6, 0.28,
          font_size=10, color=C_TEAL, italic=True)

    y = 1.82
    for lang, ft, sm, fc, sc in zip(langs, ftwer, smwer, ftcer, smcer):
        txbox(s, lang, 0.3, y+0.02, 2.3, 0.28,
              font_size=12, color=C_TEXT, align=PP_ALIGN.RIGHT)
        # FT Whisper bar (green only if FT is the winner, i.e. ps)
        fw = ft * scale
        col_ft = C_GREEN if ft < sm else RGBColor(0x37, 0x52, 0x65)
        box(s, bar_left, y, fw, 0.145, fill_color=col_ft)
        # SeamlessM4T bar (green if SM is the winner)
        sw = sm * scale
        col_sm = C_GREEN if sm < ft else RGBColor(0x37, 0x52, 0x65)
        box(s, bar_left, y+0.155, sw, 0.145, fill_color=col_sm)
        txbox(s, f"{ft:.2f}% ({fc:.1f})", bar_left+fw+0.06, y, 1.45, 0.16,
              font_size=9, bold=(ft < sm), color=col_ft)
        txbox(s, f"{sm:.2f}% ({sc:.1f})", bar_left+sw+0.06, y+0.155, 1.45, 0.16,
              font_size=9.5, bold=(sm < ft), color=col_sm)
        y += 0.81

    txbox(s, "Labels: WER% (CER%). CER shown because Han and Perso-Arabic orthographies make "
             "whitespace WER misleading; zh is char-segmented so WER = CER.",
          bar_left, y + 0.05, 5.3, 0.5, font_size=8.5, italic=True, color=C_SUB)

    # Summary panel — deployed backend per language
    card(s, 8.15, 1.55, 4.8, 5.55, accent_color=C_GOLD)
    txbox(s, "Deployed ASR Backend", 8.35, 1.65, 4.4, 0.36,
          font_size=13, bold=True, color=C_NAVY)
    selection = [
        ("Punjabi",  "57.4 → 19.8%", "SeamlessM4T"),
        ("Nepali",   "50.9 → 28.5%", "SeamlessM4T"),
        ("Hindi",    "19.8 → 13.9%", "SM4T + LoRA"),
        ("Urdu",     "19.8 → 16.9%", "SeamlessM4T"),
        ("Mandarin", "14.2 → 11.7%", "SeamlessM4T"),
        ("Pashto",   "38.6 → 36.2%", "SM4T + LoRA"),
        ("Kashmiri", "74.0 → 50.3%*","SM4T + LoRA"),
    ]
    y2 = 2.12
    for lang, prog, backend in selection:
        box(s, 8.25, y2, 4.6, 0.56, fill_color=C_CARD if selection.index((lang,prog,backend))%2==0 else C_CARD2)
        txbox(s, lang,  8.35, y2+0.05, 1.2, 0.26, font_size=11, bold=True, color=C_TEXT)
        txbox(s, prog,  9.55, y2+0.05, 1.6, 0.26, font_size=9.5, color=C_SUB)
        txbox(s, backend, 11.15, y2+0.05, 1.65, 0.26, font_size=10, bold=True,
              color=(C_TEAL if "SM4T" in backend or backend=="SeamlessM4T" else C_GREEN), align=PP_ALIGN.RIGHT)
        y2 += 0.61

# ── Slide 6 — PA v3 Training Progress ────────────────────────────────────────

def slide_06_pa_v3_progress(prs, n=9):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Punjabi v3 Training Progress (LoRA r=16)", n, "PA v3 Update")

    # Config comparison
    txbox(s, "Config Comparison: v2 vs v3", 0.4, 1.42, 5.8, 0.36,
          font_size=13, bold=True, color=C_NAVY)
    comp = [
        ("Parameter",          "v2 (superseded)",   "v3 (deployed)"),
        ("LoRA Rank",          "r = 8",             "r = 16"),
        ("LoRA Alpha",         "α = 16",            "α = 32"),
        ("Trainable Params",   "~3.9M (0.25%)",     "~7.9M (0.51%)"),
        ("IndicVoices-R",      "9,407 samples",     "20,000 samples"),
        ("Total Train",        "11,923 samples",    "21,923 samples"),
        ("Steps",              "3,000",             "4,000"),
        ("Best WER",           "52.55% (step 3000)","49.31% (step 4000)"),
    ]
    y = 1.83
    for i, (k, v2, v3) in enumerate(comp):
        bg_col = C_TEAL if i == 0 else (C_CARD if i%2==1 else C_CARD2)
        box(s, 0.4, y, 6.0, 0.37, fill_color=bg_col)
        fc = C_TEXT if i==0 else C_SUB
        txbox(s, k,  0.52, y+0.06, 2.0, 0.27, font_size=10, bold=(i==0), color=fc)
        txbox(s, v2, 2.52, y+0.06, 1.9, 0.27, font_size=10, bold=(i==0),
              color=C_TEXT if i>0 else C_TEXT)
        txbox(s, v3, 4.42, y+0.06, 2.0, 0.27, font_size=10, bold=(i==0),
              color=C_GOLD if i>0 else C_TEXT)
        y += 0.38

    # WER progression table
    txbox(s, "Eval WER Progression (v3)", 6.7, 1.42, 6.2, 0.36,
          font_size=13, bold=True, color=C_NAVY)
    wer_data = [
        ("Train Step", "v3 WER (%)",  "v2 WER (%)", "v3 − v2 (pp)"),
        ("200",  "69.97%",  "70.75%", "−0.78 pp"),
        ("1000", "56.96%",  "59.09%", "−2.13 pp ✓"),
        ("1800", "52.06%",  "54.65%", "−2.59 pp ✓"),
        ("2200", "50.62%",  "53.88%", "−3.26 pp ✓"),
        ("2800", "50.51%",  "52.61%", "−2.10 pp ✓"),
        ("3000", "50.05%",  "52.55%", "−2.50 pp ✓"),
        ("3400", "49.40%",  "  —  ",  "sub-50%"),
        ("3800", "49.49%",  "  —  ",  "  —  "),
        ("4000", "49.31%★", "  —  ",  "best ✓"),
    ]
    y = 1.83
    for i, (st, v3w, v2w, delta) in enumerate(wer_data):
        bg_col = C_TEAL if i==0 else (C_CARD if i%2==1 else C_CARD2)
        box(s, 6.7, y, 6.2, 0.37, fill_color=bg_col)
        fc = C_TEXT
        is_best = "★" in v3w
        txbox(s, st,  6.82, y+0.06, 1.1, 0.27, font_size=10, bold=(i==0 or is_best), color=fc)
        txbox(s, v3w, 7.92, y+0.06, 1.5, 0.27, font_size=10,
              bold=(i==0 or is_best), color=C_GOLD if is_best else fc)
        txbox(s, v2w, 9.42, y+0.06, 1.5, 0.27, font_size=10, color=C_SUB)
        col = C_GREEN if "✓" in delta else (C_RED if delta.startswith("+") else C_SUB)
        txbox(s, delta, 10.9, y+0.06, 1.9, 0.27, font_size=10, color=col, align=PP_ALIGN.RIGHT)
        y += 0.38

    # Status banner — DEPLOYED
    box(s, 0.4, 6.5, 12.55, 0.52, fill_color=C_GREEN)
    txbox(s, "✔  DEPLOYED: step-4000 @ 49.31% (v2 was 52.55%) — v3 wins by 3.24 pp, verified transcribing Gurmukhi  "
             "|  Completed 4000 steps (OOM at 2400 recovered via greedy eval)",
          0.55, 6.55, 12.3, 0.4, font_size=10.5, bold=True, color=C_TEXT)

# ── Slide 7 — Per-Language Deep Dive ─────────────────────────────────────────

def slide_07_language_deep(prs, n=12):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Per-Language Training Details", n, "Language Deep-Dive")

    cards = [
        ("Punjabi (pa) v3", "→ SM4T", C_TEAL,
         ["Base: whisper-large-v3  |  LoRA r=16, α=32",
          "Train: FLEURS pa_in + IV-R 20k = 21,923 samples",
          "Baseline 77.6%  →  FT test 57.39%  (−20 pp)",
          "SeamlessM4T 19.77% wins → DEPLOYED backend",
          "FT model retained but not served"]),
        ("Urdu (ur)", "→ SM4T", C_TEAL,
         ["Base: whisper-large-v3  |  LoRA r=8",
          "Train: FLEURS ur_pk  (2,109 samples)",
          "Baseline 21.23%  →  FT test 19.82%  (−1.4 pp)",
          "SeamlessM4T 16.9% wins → DEPLOYED backend",
          "CT2 tokenizer fix earlier: 100.66% → 19.52%"]),
        ("Hindi (hi)", "→ SM4T", C_TEAL,
         ["Base: whisper-large-v3  |  LoRA r=8",
          "Train: FLEURS hi_in  (2,120 samples)",
          "Baseline 26.34%  →  FT test 19.78%  (−6.6 pp)",
          "SM4T + IV-R LoRA adapter 12.91% → DEPLOYED",
          "max_grad_norm=0.5 (post-Mandarin fix)"]),
        ("Pashto (ps)", "→ SM4T", C_TEAL,
         ["Base: pashto-ghag-whisper-medium-asr",
          "Train: FLEURS ps_af  (2,082 samples)",
          "Baseline 89.76%  →  FT test 38.55%  (−51 pp)",
          "Beat zero-shot SM4T for a year — finally lost to a",
          "noise-aug SM4T LoRA (36.91%), then its r=128 cloud",
          "retrain ps_cloud (36.16%) → DEPLOYED"]),
        ("Mandarin (zh)", "→ SM4T", C_RED,
         ["Base: whisper-large-v3  |  LoRA r=8",
          "Baseline 10.99%  →  FT test 14.22%  (+3.2 pp ✗)",
          "Fine-tuning REGRESSED — prior 100→16% was a",
          "whitespace-scoring artefact (see correction slide)",
          "SeamlessM4T 11.69% wins → DEPLOYED backend"]),
        ("Nepali (ne) v2", "→ SM4T", C_TEAL,
         ["Base: whisper-large-v3  |  LoRA r=8",
          "Train: FLEURS ne_np + IV-R = 13,332 samples",
          "Baseline 88.85%  →  FT test 50.92%  (−38 pp)",
          "SM4T + IV-R LoRA adapter 24.34% → DEPLOYED",
          "Nepali's first working adapter (post label fix)"]),
    ]

    positions = [(0.3, 1.42), (4.55, 1.42), (8.8, 1.42),
                 (0.3,  4.38), (4.55, 4.38), (8.8, 4.38)]
    for (left, top), (lang, wer, col, bullets) in zip(positions, cards):
        card(s, left, top, 4.1, 2.8, accent_color=col)
        box(s, left, top, 4.1, 0.36, fill_color=col)
        txbox(s, lang, left+0.12, top+0.06, 2.6, 0.28,
              font_size=12, bold=True, color=C_TEXT)
        txbox(s, wer, left+2.65, top+0.06, 1.3, 0.28,
              font_size=16, bold=True, color=C_TEXT, align=PP_ALIGN.RIGHT)
        y = top + 0.44
        for b in bullets:
            txbox(s, f"• {b}", left+0.12, y, 3.85, 0.36, font_size=9.5, color=C_SUB)
            y += 0.37

# ── Slide 8 — Kashmiri ────────────────────────────────────────────────────────

def slide_08_kashmiri(prs, n=20):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Kashmiri (ks) — Custom Token Fine-Tuning (Whisper, superseded)", n, "Kashmiri")

    # Challenge panel
    card(s, 0.4, 1.42, 5.85, 5.35, accent_color=C_RED)
    box(s, 0.4, 1.42, 5.85, 0.35, fill_color=C_RED)
    txbox(s, "The Challenge", 0.55, 1.47, 5.5, 0.28,
          font_size=12, bold=True, color=C_TEXT)
    chs = [
        "Whisper has no native <|ks|> language token",
        "No FLEURS or Common Voice Kashmiri corpus",
        "Only IndicVoices-R available (AI4Bharat)",
        "Baseline WER: 96.87% (turbo model)",
        "Nastaliq script — same as Urdu & Pashto",
        "Requires patching faster-whisper source",
    ]
    bullet_block(s, chs, 0.55, 1.9, 5.5, font_size=12, color=C_SUB, spacing=0.62)

    # Solution panel
    card(s, 6.75, 1.42, 5.85, 5.35, accent_color=C_TEAL)
    box(s, 6.75, 1.42, 5.85, 0.35, fill_color=C_TEAL)
    txbox(s, "Engineering Solution", 6.9, 1.47, 5.5, 0.28,
          font_size=12, bold=True, color=C_TEXT)
    sols = [
        "Added <|ks|> token (ID 51866) to vocab & embedding",
        "Embedding initialised from <|ur|> (Urdu — same script)",
        "TemplateProcessing prefix forced at inference",
        "Trained on 20,000 IndicVoices-R KS samples",
        "3,000 steps  →  best WER 74.02% at step 2400",
        "Improvement: −22.85 pp from 96.87% baseline",
        "faster-whisper patched to accept language='ks'",
    ]
    bullet_block(s, sols, 6.9, 1.9, 5.5, font_size=12, color=C_SUB, spacing=0.62)

    # Stats bar
    box(s, 0.4, 6.88, 12.55, 0.42, fill_color=C_CARD)
    stats = [
        ("Base", "whisper-large-v3"),
        ("Token ID", "51866  (<|ks|>)"),
        ("Train Data", "20,000 IV-R samples"),
        ("Best WER", "74.02% @ step 2400"),
        ("CT2 Model", "whisper-large-v3-ks-ct2"),
    ]
    x = 0.55
    for k, v in stats:
        txbox(s, k, x, 6.92, 1.5, 0.18, font_size=9, color=C_SUB)
        txbox(s, v, x, 7.08, 1.9, 0.18, font_size=9.5, bold=True, color=C_GOLD)
        x += 2.48
    txbox(s, "This model documents HOW the Whisper <|ks|> token was engineered. It was "
             "superseded 2026-07-20 by ks_max — the analogous SeamlessM4T trick (custom "
             "__kas__ token, r=32 LoRA + a trainable embedding row) — and the line kept "
             "improving: ks_max2 (4x combined corpus, 61.88%), ks_cloud (r=128 on a rented "
             "A6000, 56.44%), and ks_cloud2 — the same run allowed to converge instead of "
             "early-stopping at 0.8 epochs — at 52.60%, and finally ks_cloud3, which repaired 20 "
             "Kashmiri characters that had no token in the model vocabulary, at 50.26% "
             "diacritic-normalised — now DEPLOYED. Rollback only.",
          0.4, 6.55, 12.5, 0.32, font_size=9, italic=True, color=C_GREEN)

# ── Slide 9 — VANI Pipeline ───────────────────────────────────────────────────

def slide_09_pipeline(prs, n=21):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "VANI 10-Stage Intelligence Pipeline", n, "Pipeline")

    stages = [
        ("1",  "VAD",            "Silero — detect speech, discard silence",            C_TEAL),
        ("2",  "Preprocessing",  "Bandpass 300–3400 Hz + noise reduction + normalise",  C_TEAL),
        ("3",  "MMS-LID",        "Facebook MMS 256-language ID — route to correct model",C_GOLD),
        ("4",  "ASR",            "SeamlessM4T v2, per-language backend (LoRA / zero-shot)", C_GOLD),
        ("5",  "Script Cascade", "Arabic-script override: >20% Nastaliq → Urdu routing",C_GOLD),
        ("6",  "Translation",    "NLLB-200 distilled 600M → English transcription",     C_GREEN),
        ("7",  "Diarization",    "Speaker separation — up to 4 speakers identified",    C_GREEN),
        ("8",  "Keywords",       "Multilingual keyword + threat-entity detection",       C_GREEN),
        ("9",  "ISUM",           "Gemma 3:12B (Ollama) → 4-sentence Intel Summary",    C_TEAL),
        ("10", "Export",         "SQLite + JSON output + Streamlit web UI (offline)",   C_TEAL),
    ]

    y = 1.42
    for i, (num, name, desc, col) in enumerate(stages):
        bg_col = C_CARD if i % 2 == 0 else C_CARD2
        box(s, 0.3, y, 12.75, 0.44, fill_color=bg_col)
        box(s, 0.3, y, 0.055, 0.44, fill_color=col)
        txbox(s, num,  0.43, y+0.08, 0.6,  0.28, font_size=11, bold=True, color=col)
        txbox(s, name, 1.05, y+0.08, 2.2,  0.28, font_size=12, bold=True, color=C_TEXT)
        txbox(s, desc, 3.3,  y+0.08, 9.5,  0.28, font_size=11, color=C_SUB)
        y += 0.46

# ── Slide 10 — Key Findings ───────────────────────────────────────────────────

def slide_10_key_findings(prs, n=22):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Key Findings & Technical Insights", n, "Key Findings")

    findings = [
        ("01", C_GOLD,   "LoRA Effective at 0.25% Params",
         "WER drops 1.4–51 pp where fine-tuning helps (zh regressed). r=8/r=16 sufficient. "
         "Full fine-tune (40+ GB VRAM) unnecessary."),
        ("02", C_GREEN,  "FLEURS + IndicVoices-R Transfer to Radio Domain",
         "Clean read-speech training improves military-domain accuracy. PA v2: +9,407 IV-R samples → −4.1 pp."),
        ("03", C_TEAL,   "MMS-LID Critical for Pashto & Robustness",
         "Whisper-only scores 0% on Pashto LangID. MMS-LID lifts it to 53–87% across 5 degradation conditions."),
        ("04", C_RED,    "fp16 Gradient Explosion Risk at Low LR",
         "Mandarin diverged at step ~820 (grad_norm=12.9, loss 0.15→0.77). Fix: max_grad_norm=0.5 for all subsequent."),
        ("05", C_GOLD,   "CT2 Tokenizer Bug — All large-v3 Models Translated Not Transcribed",
         "ct2-transformers-converter omits tokenizer.json. large-v3: translate=50359; tiny: transcribe=50359. One-token shift. Now automated."),
        ("06", C_TEAL,   "Script Cascade Prevents Nastaliq Misidentification",
         ">20% Nastaliq chars in transcript → Urdu override. Critical for noisy radio audio."),
        ("07", C_PURPLE, "Full VANI +14 pp over Whisper-Only (Robustness Eval)",
         "5 conditions × 7 langs × 4 configs. ZH: 97–100%. PS: MMS-LID essential. Mandarin most robust."),
    ]

    y = 1.42
    for num, col, title, detail in findings:
        box(s, 0.3,  y, 0.55, 0.82, fill_color=col)
        txbox(s, num, 0.3, y+0.18, 0.55, 0.4,
              font_size=15, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
        box(s, 0.9, y, 12.1, 0.82, fill_color=C_CARD)
        txbox(s, title,  1.05, y+0.05, 11.8, 0.32, font_size=12, bold=True, color=col)
        txbox(s, detail, 1.05, y+0.38, 11.8, 0.42, font_size=10.5, color=C_SUB)
        y += 0.94

# ── Slide 11 — Robustness Eval ────────────────────────────────────────────────

def slide_11_robustness(prs, n=23):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Radio-Channel Robustness — LangID Accuracy (%)", n, "Robustness Eval")

    txbox(s, "5 degradation conditions × 7 languages × 4 pipeline configs  |  30 samples/lang  |  1-Jul-2026  |  "
             "cells = MMS language-ID accuracy (%), higher is better",
          0.4, 1.33, 12.5, 0.28, font_size=9.5, italic=True, color=C_SUB)

    cols_x = [0.3, 2.05, 3.15, 4.05, 4.95, 5.85, 6.8, 7.7, 8.6, 9.55]
    hdrs   = ["Channel\nCondition", "Pipeline\nConfig", "PA", "HI", "UR", "NE", "ZH", "PS", "KS", "Overall\nAvg"]

    box(s, 0.3, 1.63, 12.75, 0.46, fill_color=C_NAVY)
    for x, h in zip(cols_x, hdrs):
        _multiline(s, h, x, 1.66, 1.05, 0.40,
                   font_size=9, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.LEFT)

    rows_data = [
        ("Clean",     "Whisper",   "90%","87%","73%","27%","100%", "0%", "0%","54%", C_SUB),
        ("",          "Full VANI", "90%","83%","70%","33%","100%","80%","17%","68%", C_GREEN),
        ("Bandpass",  "Whisper",   "70%","83%","33%","27%","100%", "0%", "0%","45%", C_SUB),
        ("",          "Full VANI", "73%","83%","33%","27%","100%","83%","23%","61%", C_GREEN),
        ("AWGN 10dB", "Whisper",   "97%","63%","80%","43%","100%", "0%", "0%","55%", C_SUB),
        ("",          "Full VANI", "97%","63%","80%","43%","100%","83%", "3%","67%", C_GREEN),
        ("AWGN 0dB",  "Whisper",   "57%","63%","27%","13%","100%", "0%", "0%","37%", C_SUB),
        ("",          "Full VANI", "60%","67%","27%","17%", "97%","53%", "0%","46%", C_RED),
        ("MP3 16kbps","Whisper",   "77%","47%","63%","23%","100%", "0%", "0%","44%", C_SUB),
        ("",          "Full VANI", "77%","47%","63%","23%", "97%","87%","20%","59%", C_GREEN),
    ]

    y = 2.14
    for i, row in enumerate(rows_data):
        cond, cfg = row[0], row[1]
        vals, color = row[2:10], row[10]
        # group separator
        if cond and i > 0:
            box(s, 0.3, y, 12.75, 0.025, fill_color=C_BORDER)
        bg_col = C_CARD if (i // 2) % 2 == 0 else C_CARD2
        box(s, 0.3, y, 12.75, 0.33, fill_color=bg_col)
        if cond:
            txbox(s, cond, cols_x[0], y+0.03, 1.7, 0.27,
                  font_size=10, bold=True, color=C_TEXT)
        is_vani = cfg == "Full VANI"
        txbox(s, cfg, cols_x[1], y+0.03, 1.05, 0.27,
              font_size=10, bold=is_vani, color=color)
        for xi, v in zip(cols_x[2:], vals):
            txbox(s, v, xi, y+0.03, 0.88, 0.27,
                  font_size=10, bold=is_vani, color=color, align=PP_ALIGN.CENTER)
        y += 0.33

    txbox(s, "ZH: most robust (97–100% all conditions)  ·  "
             "PS: MMS-LID essential, Whisper-only = 0%  ·  "
             "Full VANI avg +14 pp over Whisper-only",
          0.3, 5.45, 12.5, 0.38, font_size=10.5, italic=True, color=C_SUB)

# ── Slide — Next Steps to Improve Accuracy ───────────────────────────────────

def slide_next_steps_accuracy(prs, n):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Next Steps to Improve Accuracy", n, "Roadmap")

    txbox(s, "Prioritised levers to push WER lower — ordered by expected impact for the high-WER languages (PA, NE, KS).",
          0.4, 1.30, 12.5, 0.28, font_size=10, italic=True, color=C_SUB)

    cards = [
        ("1  ·  Scale Training Data", C_GOLD,
         ["Biggest lever for PA / NE / KS (highest WER)",
          "Add IndicVoices-R full split + Shrutilipi",
          "KS starved at 20k → target 50k+ samples",
          "Common Voice + FLEURS-R for extra hours"]),
        ("2  ·  Increase Model Capacity", C_TEAL,
         ["PA v3 r=16 already beats r=8 by 2–3 pp",
          "Try r=32 / α=64 for PA, NE, KS",
          "Expand targets: add k_proj, o_proj",
          "Optionally add fc1/fc2 (MLP) adapters"]),
        ("3  ·  Longer / Smarter Training", C_BLUE,
         ["NE loss still falling at step 3000 → extend",
          "Cosine LR schedule + warm restarts",
          "Higher LR (7e-5) with max_grad_norm=0.5",
          "More frequent eval near convergence"]),
        ("4  ·  Radio-Domain Augmentation", C_PURPLE,
         ["SpecAugment (time/freq masking)",
          "RIR reverb + additive noise injection",
          "MP3 / codec + bandpass augmentation",
          "Pseudo-label real radio intercepts"]),
        ("5  ·  Decoding-Time Gains", C_GREEN,
         ["Tune beam width + temperature fallback",
          "KenLM shallow fusion for low-resource langs",
          "Repetition / no-speech threshold tuning",
          "Script-cascade confidence calibration"]),
        ("6  ·  Recover Mandarin Headroom", C_RED,
         ["Val WER hit 8.97% before divergence",
          "Resume past ckpt-400 at lower LR (2e-5)",
          "Stronger clipping to avoid fp16 overflow",
          "Target ZH test WER < 12%"]),
    ]

    positions = [(0.3, 1.60), (4.55, 1.60), (8.8, 1.60),
                 (0.3, 4.55), (4.55, 4.55), (8.8, 4.55)]
    for (left, top), (title, col, bullets) in zip(positions, cards):
        card(s, left, top, 4.1, 2.78, accent_color=col)
        box(s, left, top, 4.1, 0.40, fill_color=col)
        txbox(s, title, left+0.12, top+0.07, 3.85, 0.30,
              font_size=12, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
        y = top + 0.52
        for b in bullets:
            txbox(s, f"▸  {b}", left+0.12, y, 3.85, 0.4, font_size=10, color=C_SUB)
            y += 0.53

    box(s, 0.3, 7.14, 12.75, 0.30, fill_color=C_TEAL_L)
    txbox(s, "Delivered: PA v3 (r=16, 21,923 samples) completed 4000 steps → 49.31% deployed, −3.24 pp below v2. Next levers below.",
          0.45, 7.17, 12.5, 0.24, font_size=9.5, bold=True, italic=True, color=C_NAVY)

# ── Slide 12 — Deployment & Conclusions ──────────────────────────────────────

def slide_12_conclusion(prs, n=24):
    s = blank_slide(prs)
    bg(s)
    top_band(s)
    txbox(s, "Deployment & Conclusions", 0.5, 0.12, 12.0, 0.62,
          font_size=27, bold=True, color=C_TEXT)
    hline(s, 0.78, width_in=12.33, left=0.5)
    slide_num(s, n)

    # Summary boxes
    summ = [
        (C_TEAL,   "7 Models Fine-Tuned\npa · ps · ur · ne · zh · hi · ks",
                   "All CT2 int8, fully offline"),
        (C_GOLD,   "Backend Selected\nper Language",
                   "SeamlessM4T ×7  ·  FT Whisper ×0"),
        (C_GREEN,  "All 7 Languages Now\nRun on SeamlessM4T",
                   "Whisper models kept only for rollback"),
        (C_PURPLE, "0.25–0.51% Params\nLoRA r=8 / r=16",
                   "No full fine-tune needed"),
    ]
    x = 0.4
    for col, title, detail in summ:
        card(s, x, 0.93, 3.08, 1.55, accent_color=col)
        box(s, x, 0.93, 3.08, 0.055, fill_color=col)
        txbox(s, title,  x+0.15, 1.0,  2.8, 0.6,  font_size=11.5, bold=True, color=col)
        txbox(s, detail, x+0.15, 1.62, 2.8, 0.36, font_size=10, color=C_SUB)
        x += 3.18

    # Production ASR backend per language (held-out n=100 test WER)
    txbox(s, "Production ASR Backend per Language", 0.4, 2.65, 7.0, 0.35,
          font_size=13, bold=True, color=C_NAVY)
    models = [
        ("SeamlessM4T (zero-shot)",   "Punjabi",  "19.77%", C_TEAL),
        ("SeamlessM4T + ps LoRA",     "Pashto",   "36.16%", C_TEAL),
        ("SeamlessM4T (zero-shot)",   "Urdu",     "16.90%", C_TEAL),
        ("SeamlessM4T + ne LoRA",     "Nepali",   "24.34%", C_TEAL),
        ("SeamlessM4T (zero-shot)",   "Mandarin", "11.69%", C_TEAL),
        ("SeamlessM4T + hi LoRA",     "Hindi",    "12.91%", C_TEAL),
        ("SeamlessM4T + ks LoRA",     "Kashmiri", "50.26%*",C_TEAL),
    ]
    y = 3.04
    for i, (mname, lang, wer, col) in enumerate(models):
        bg_col = C_CARD if i % 2 == 0 else C_CARD2
        box(s, 0.4, y, 7.1, 0.36, fill_color=bg_col)
        box(s, 0.4, y, 0.055, 0.36, fill_color=col)
        txbox(s, mname, 0.55, y+0.06, 4.0, 0.25, font_size=9.5, color=C_SUB)
        txbox(s, lang,  4.6,  y+0.06, 1.1, 0.25, font_size=10, bold=True, color=C_TEXT)
        txbox(s, wer,   5.8,  y+0.06, 1.6, 0.25, font_size=10.5, bold=True, color=col, align=PP_ALIGN.RIGHT)
        y += 0.37
    txbox(s, "* ks: diacritic-normalised WER (raw is inflated by dense Perso-Arabic marks for BOTH "
             "systems). Deployed ks_cloud3 (r=128 + 20 repaired vocab chars) vs Whisper-ks 65.19 — see notes.",
          0.4, y+0.02, 7.1, 0.3, font_size=8, italic=True, color=C_SUB)

    # Next steps
    txbox(s, "Next Steps", 7.8, 2.65, 5.4, 0.35,
          font_size=13, bold=True, color=C_NAVY)
    nexts = [
        (C_GREEN,  "Backend routing", "DEPLOYED — SeamlessM4T for all 7 languages (hi/ne/ps/ks w/ LoRA)"),
        (C_GREEN,  "Robustness",    "COMPLETE — 5 cond × 7 langs; SM4T lead holds under noise"),
        (C_GOLD,   "Seamless timing", "Per-utterance segments added; per-VAD-segment ASR live"),
        (C_PURPLE, "Reports",       "Regenerated from corrected n=100 scoring (this deck)"),
        (C_TEAL,   "Paper",         "VANI → IJAINN / SLT 2026 submission"),
    ]
    y2 = 3.04
    for col, title, detail in nexts:
        box(s, 7.8, y2, 5.4, 0.52, fill_color=C_CARD)
        box(s, 7.8, y2, 0.055, 0.52, fill_color=col)
        txbox(s, title,  7.92, y2+0.05, 1.35, 0.38, font_size=11, bold=True, color=col)
        txbox(s, detail, 9.28, y2+0.05, 3.85, 0.38, font_size=10, color=C_SUB)
        y2 += 0.57

    box(s, 0, 7.1, 13.33, 0.40, fill_color=C_NAVY)
    txbox(s, "VANI  ·  M.Tech Research  ·  IIT Indore  ·  2026",
          0.5, 7.14, 12.33, 0.28,
          font_size=11, italic=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)


# ── Chart slide helper ────────────────────────────────────────────────────────

def chart_slide(prs, title, n, tag, buf, caption="",
                img_top=1.42, img_width=11.0, aspect=2.0):
    """Embed a matplotlib PNG BytesIO into a new slide."""
    s = blank_slide(prs)
    bg(s)
    slide_header(s, title, n, tag)
    img_h = img_width / aspect
    img_left = (13.33 - img_width) / 2
    s.shapes.add_picture(buf, Inches(img_left), Inches(img_top),
                         Inches(img_width), Inches(img_h))
    if caption:
        txbox(s, caption, 0.4, img_top + img_h + 0.10, 12.5, 0.32,
              font_size=9.5, italic=True, color=C_SUB, align=PP_ALIGN.CENTER)

# ── Chart slides — overview ────────────────────────────────────────────────────

def slide_chartA_dataset(prs, n):
    chart_slide(
        prs, "Training Dataset Sizes per Language", n, "Data",
        chart_dataset_sizes(),
        caption="Figure 1: Training vs validation samples per language. PA and NE bars reflect v2 totals (FLEURS + IndicVoices-R).",
        img_width=10.0, aspect=2.0,
    )

def slide_chartB_summary_bar(prs, n):
    chart_slide(
        prs, "Baseline vs. Fine-Tuned WER — All Languages", n, "Results",
        chart_summary_bar(),
        caption="Figure 2: True large-v3 baseline WER (no fine-tuning) vs best fine-tuned WER, held-out n=100 test. "
                "Mandarin shows a small regression (baseline 10.99% < fine-tuned 14.22%); the prior 100% baseline "
                "was a whitespace-scoring artefact, now corrected.",
        img_width=11.0, aspect=2.0,
    )

def slide_chartC_wer_all(prs, n):
    chart_slide(
        prs, "WER Progression During LoRA Training", n, "Training Curves",
        chart_wer_all(),
        caption="Figure 3: Eval WER at each checkpoint for all 7 languages. "
                "× marks Mandarin divergence at step ~820 (252% — clipped to axis). "
                "PA v3 dotted line shown separately.",
        img_width=11.5, aspect=2.0,
    )

# ── Chart slides — per language ────────────────────────────────────────────────

def slide_chartD_lang(prs, lang, n):
    m = LANG_META[lang]
    title = f"{m['name']} ({lang.upper()}) — WER & Training Loss Curves"
    if lang == "zh":
        # ZH: training-val WER 8.97% at step 400 did NOT generalise; held-out test 14.22%,
        # which is WORSE than the 10.99% large-v3 baseline. Not deployed; routes to Seamless.
        caption = (
            "Left: Eval WER on the 409-sample training-validation set — best 8.97% at step 400 (★). "
            "★ This did not generalise: held-out FLEURS test WER is 14.22%, above the 10.99% large-v3 "
            "baseline, so fine-tuning regressed Mandarin. Not deployed — routes to SeamlessM4T (11.69%). "
            "Training diverged at step ~820 (grad_norm=12.9). Right: Training loss (healthy until divergence)."
        )
    else:
        caption = (
            f"Left: Training-val WER over training steps (★ = best checkpoint, {m['best_wer']:.2f}%). "
            f"Right: Training loss (monotonic descent confirms healthy convergence). "
            f"Held-out test: baseline {m['baseline_wer']:.2f}%  →  fine-tuned {m['eval_wer']:.2f}%  "
            f"({m['eval_wer'] - m['baseline_wer']:+.1f} pp)."
        )
    chart_slide(
        prs, title, n, "Training Curves",
        chart_wer_per_lang(lang),
        caption=caption,
        img_width=12.2, aspect=2.5,
    )

# ── Table slide — training version history ────────────────────────────────────

def slide_table_version_history(prs, n):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Training Version History — All Reruns", n, "Version History")

    txbox(s, "Documents every training run across all 7 languages, including re-runs triggered by data additions and architecture changes.",
          0.4, 1.37, 12.5, 0.28, font_size=10, italic=True, color=C_SUB)

    cols_x  = [0.30, 1.55, 2.65, 4.55, 6.25, 7.95, 9.30, 10.50, 11.65]
    headers = ["Language", "Run\nVersion", "Base\nModel", "Training\nDataset",
               "Train\nSamples", "Train\nSteps", "Best Val\nWER (%)", "Test\nWER (%)", "Deploy\nStatus"]
    y = 1.72
    box(s, 0.28, y, 12.77, 0.52, fill_color=C_NAVY)
    for x, h in zip(cols_x, headers):
        _multiline(s, h, x, y+0.05, 1.25, 0.44,
                   font_size=8.5, bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    rows = [
        ("PA",  "v1",         "large-v3",       "FLEURS pa_in",          "2,516",  "3,000", "56.67%", "56.67%", "Superseded", C_SUB),
        ("PA",  "v2",         "large-v3",       "FLEURS+IV-R",           "11,923", "3,000", "52.55%", "52.55%", "Superseded", C_SUB),
        ("PA",  "v3",         "large-v3",       "FLEURS+IV-R (20k)",     "21,923", "4,000", "49.31%", "57.39%", "SM4T serves", C_SUB),
        ("NE",  "v1",         "large-v3",       "FLEURS ne_np",          "3,332",  "2,000", "52.14%", "52.14%", "Superseded", C_SUB),
        ("NE",  "v2",         "large-v3",       "FLEURS+IV-R",           "13,332", "3,000", "50.82%", "50.92%", "SM4T serves", C_SUB),
        ("ZH",  "v1 (div.)", "large-v3",       "FLEURS cmn_hans_cn",    "3,246",  "400†",  "8.97%",  "14.22%", "SM4T serves", C_RED),
        ("PS",  "v1",         "medium-pashto",  "FLEURS ps_af",          "2,082",  "2,000", "38.55%", "38.55%", "SM4T serves", C_SUB),
        ("UR",  "v1",         "large-v3",       "FLEURS ur_pk",          "2,109",  "1,000", "22.27%", "19.82%", "SM4T serves", C_SUB),
        ("HI",  "v1",         "large-v3",       "FLEURS hi_in",          "2,120",  "600",   "23.13%", "19.78%", "SM4T serves", C_SUB),
        ("KS",  "v1",         "large-v3+<|ks|>","IndicVoices-R KS",     "20,000", "3,000", "74.02%", "74.02%", "SM4T serves", C_SUB),
    ]

    y += 0.54
    for i, row in enumerate(rows):
        lang, ver, base, data, samp, steps, bwer, ewer, status, col = row
        bg_col = C_CARD if i % 2 == 0 else C_CARD2
        box(s, 0.28, y, 12.77, 0.34, fill_color=bg_col)
        vals = [lang, ver, base, data, samp, steps, bwer, ewer, status]
        for j, (x, v) in enumerate(zip(cols_x, vals)):
            c = col if j in (6, 8) else (C_TEXT if j <= 1 else C_SUB)
            txbox(s, v, x, y+0.05, 1.25, 0.25,
                  font_size=8.5, bold=(j <= 1 or j == 8), color=c)
        y += 0.35

    txbox(s, "† ZH diverged step ~820, ckpt-400 kept  ·  IV-R = IndicVoices-R  ·  "
             "\"SM4T serves\" = this fine-tuned Whisper model was trained and evaluated, but SeamlessM4T "
             "(zero-shot or + LoRA) is the deployed backend for this language; the Whisper model is "
             "retained on disk for rollback only. As of 2026-07-20 this holds for ALL 7 languages.",
          0.35, y+0.08, 12.5, 0.28, font_size=9, italic=True, color=C_SUB)

# ── Table slide — FT Whisper vs SeamlessM4T ───────────────────────────────────

def slide_table_sm4t(prs, n):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Cross-Model Benchmark — FT Whisper · NLLB · SeamlessM4T", n, "Comparison")

    txbox(s, "n=100 FLEURS test / IndicVoices-R  ·  re-run 11 Jul 2026, corrected scoring  ·  "
             "ASR: WER ↓ better  ·  Translation: chrF ↑ better  ·  "
             "Baseline = true large-v3 (not turbo)  ·  SM4T = SeamlessM4T-large-v2",
          0.4, 1.37, 12.5, 0.28, font_size=9.5, italic=True, color=C_SUB)

    cols_x  = [0.30, 1.55, 2.80, 4.05, 5.40, 6.75, 8.20, 9.60, 11.05]
    headers = ["Language", "Script", "Baseline\nWER (%)", "FT Whisper\nWER (%)", "SeamlessM4T\nWER (%)",
               "Deployed\nBackend", "NLLB-200\nchrF ↑", "SeamlessM4T\nchrF ↑", "Best Transl.\nModel"]
    y = 1.72
    box(s, 0.28, y, 12.77, 0.52, fill_color=C_NAVY)
    for x, h in zip(cols_x, headers):
        _multiline(s, h, x, y+0.05, 1.28, 0.44,
                   font_size=8, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.LEFT)

    sm4t_data = [
        ("Punjabi (pa)",  "Gurmukhi",  " 77.60", "57.39", "19.77", "SM4T",     "40.15", "54.53", "SM4T"),
        ("Pashto (ps)",   "Nastaliq",  " 89.76", "38.55", "44.40", "SM4T+LoRA","44.48", "40.15", "NLLB"),
        ("Urdu (ur)",     "Nastaliq",  " 21.23", "19.82", "16.90", "SM4T",     "51.34", "50.73", "NLLB"),
        ("Nepali (ne)",   "Devanagari"," 88.85", "50.92", "28.46", "SM4T",     "45.55", "51.67", "SM4T"),
        ("Mandarin (zh)", "Han Simpl.", " 10.99", "14.22", "11.69", "SM4T",    "42.00", "49.15", "SM4T"),
        ("Hindi (hi)",    "Devanagari"," 26.34", "19.78", "15.44", "SM4T",     "53.71", "51.54", "NLLB"),
        # ks is NOT on FLEURS and NOT on this ruler: 372-clip IndicVoices-R, L2
        # diacritic-normalised. It previously showed 96.87 / 74.02, which are
        # training-split validation figures sitting in a table of held-out ones
        # — the same ruler mix the report warns about. Held-out L2 values now.
        ("Kashmiri (ks)‡", "Nastaliq", "     —", "65.19", "50.26", "SM4T+tok", "—", "—", "N/A"),
    ]

    y += 0.54
    for i, row in enumerate(sm4t_data):
        lang, script, bwer, ftwer, smwer, winner_asr, nllb, smchrf, winner_t = row
        bg_col = C_CARD if i % 2 == 0 else C_CARD2
        box(s, 0.28, y, 12.77, 0.44, fill_color=bg_col)
        # winner colour
        w_col = C_TEAL if "FT" in winner_asr else C_SUB
        vals_cols = [
            (lang,  C_TEXT), (script, C_SUB), (bwer, C_TEXT),
            (ftwer, C_BLUE), (smwer, C_SUB),
            (winner_asr, w_col),
            (nllb, C_SUB), (smchrf, C_SUB),
            (winner_t, C_SUB),
        ]
        for x, (v, c) in zip(cols_x, vals_cols):
            txbox(s, v, x, y+0.08, 1.25, 0.32,
                  font_size=8.5, bold=(c in (C_TEXT, C_BLUE, C_TEAL)), color=c)
        y += 0.45

    txbox(s, "Zero-shot SeamlessM4T wins ASR for pa / ne / hi / ur / zh. Pashto and Kashmiri were "
             "fine-tuned Whisper's last strongholds; both fell in 2026-07 to SeamlessM4T LoRA "
             "adapters (ps: noise-augmented training; ks: a trainable __kas__ embedding row, "
             "rank 128 and a 20-character vocabulary repair, plus a scoring-ruler correction — see "
             "notes). VANI now routes all 7 languages to SeamlessM4T; FT Whisper models are "
             "retained on disk for rollback only.   "
             "‡ Kashmiri is absent from FLEURS: its two figures are the 372-clip IndicVoices-R "
             "test split at L2, not this table's FLEURS n=100 ruler, and are held-out rather than "
             "training-split. The two rulers are not comparable across a row.",
          0.35, y+0.06, 12.5, 0.32, font_size=10, italic=True, color=C_SUB)

# ══════════════════════════════════════════════════════════════════════════════
#  Professor-presentation slides (7-section flow, added 2026-07-20)
# ══════════════════════════════════════════════════════════════════════════════

def slide_section_divider(prs, sec_no, title, subtitle, n):
    s = blank_slide(prs)
    bg(s, C_NAVY)
    box(s, 0, 3.55, 13.33, 0.07, fill_color=C_TEAL)
    txbox(s, f"{sec_no:02d}", 1.0, 2.05, 3.0, 1.5,
          font_size=96, bold=True, color=C_TEAL, align=PP_ALIGN.LEFT)
    txbox(s, title, 3.6, 2.55, 9.0, 1.0,
          font_size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    txbox(s, subtitle, 3.65, 3.75, 9.0, 0.8,
          font_size=15, italic=True, color=RGBColor(0xB0, 0xC4, 0xDE))
    slide_num(s, n)


def slide_intro_problem(prs, n):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Introduction & Problem Statement", n, "Introduction")

    txbox(s, "VANI — an offline AI system that turns foreign-language radio intercepts into "
             "English intelligence summaries, on a single air-gapped workstation.",
          0.4, 1.42, 12.5, 0.6, font_size=14, italic=True, color=C_NAVY)

    card(s, 0.4, 2.15, 6.0, 4.55, accent_color=C_RED)
    box(s, 0.4, 2.15, 6.0, 0.4, fill_color=C_RED)
    txbox(s, "The Operational Problem", 0.55, 2.20, 5.7, 0.32,
          font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    probs = [
        "High-volume radio traffic in many low-resource\nborder languages (Punjabi, Pashto, Urdu,\nNepali, Mandarin, Hindi, Kashmiri, Dogri).\nAll eight now have fine-tuned ASR - Dogri was\nthe last, added 2026-07-28 (102.25 -> 46.73).",
        "Must run FULLY OFFLINE / air-gapped —\nno cloud, no internet at runtime",
        "Generic ASR fails on these languages AND\non the noisy, band-limited radio channel",
        "Manual transcription + translation does not\nscale to an operational intercept volume",
        "Analysts need actionable output (who / what /\nwhere / threat), not just a raw transcript",
    ]
    bullet_block(s, probs, 0.58, 2.72, 5.7, font_size=11.5,
                 color=C_TEXT, spacing=0.75, bullet="✗")

    card(s, 6.7, 2.15, 6.2, 4.55, accent_color=C_TEAL)
    box(s, 6.7, 2.15, 6.2, 0.4, fill_color=C_TEAL)
    txbox(s, "The Objective", 6.85, 2.20, 5.9, 0.32,
          font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    objs = [
        "One end-to-end offline pipeline:\nspeech → language ID → transcription →\nEnglish translation → threat detection →\nstructured intelligence summary",
        "Best-in-class ASR per language, auto-selected\nby audio language identification",
        "Rigorous, reproducible evaluation to choose the\nright ASR backend — under realistic radio noise",
        "A usable analyst GUI: map, dashboard, search,\nlive-mic capture, coded-terminology alerts",
    ]
    bullet_block(s, objs, 6.88, 2.72, 5.9, font_size=11.5,
                 color=C_TEXT, spacing=0.92, bullet="▸")


def slide_deliverables(prs, n):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Deliverables", n, "Deliverables")
    items = [
        ("Working VANI system", "Streamlit GUI processing radio intercepts end-to-end, fully offline on one RTX 5060 workstation", C_TEAL),
        ("7-language ASR backend", "SeamlessM4T v2, auto-selected by MMS-LID; 4 languages on LoRA adapters, 3 zero-shot", C_TEAL),
        ("Full intelligence pipeline", "VAD → LID → ASR → NLLB translation → keyword/threat detection → ISUM (Gemma 3) → SQLite + reports", C_BLUE),
        ("Backend-selection research", "Head-to-head WER + 5-condition robustness study; 5 scoring-methodology defects found and corrected", C_BLUE),
        ("Analyst features", "Map + dashboard, transcript search, live-mic capture, coded-terminology alerts, speaker diarization", C_GOLD),
        ("3-node LAN integration", "Optional distributed mode: denoise+diarize (A) + LID/dialect (B) + VANI orchestrator (C)", C_GOLD),
        ("Documentation", "This fine-tuning report (PDF), presentation deck, correction memo, reproducible eval scripts", C_PURPLE),
    ]
    y = 1.55
    for title, desc, col in items:
        box(s, 0.4, y, 12.5, 0.72, fill_color=C_CARD, line_color=C_BORDER)
        box(s, 0.4, y, 0.08, 0.72, fill_color=col)
        txbox(s, title, 0.62, y + 0.06, 3.5, 0.6, font_size=13, bold=True, color=C_NAVY)
        txbox(s, desc, 4.2, y + 0.10, 8.6, 0.55, font_size=11, color=C_SUB)
        y += 0.76


def slide_architecture(prs, n):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "System Architecture", n, "Architecture")
    txbox(s, "Single offline workstation · all models local (HF_HUB_OFFLINE) · GPU-resident with CPU parking on an 8 GB card",
          0.4, 1.42, 12.5, 0.4, font_size=12, italic=True, color=C_SUB)

    # three horizontal model-stack layers: input → models → output
    layers = [
        ("INPUT", ["Radio intercept WAV", "Live microphone", "16 kHz mono"], C_SUB),
        ("PERCEPTION", ["Silero VAD", "MMS-LID-256 (language ID)", "SeamlessM4T v2 ASR\n(4 LoRA + 3 zero-shot)"], C_TEAL),
        ("UNDERSTANDING", ["NLLB-200 translation", "Keyword / threat detection", "Gemma 3 ISUM (Ollama)"], C_BLUE),
        ("OUTPUT", ["Intelligence summary (5W)", "Map · dashboard · search", "SQLite transcript store"], C_GREEN),
    ]
    x = 0.4
    w = 3.05
    for name, boxes, col in layers:
        box(s, x, 1.95, w, 0.42, fill_color=col)
        txbox(s, name, x, 2.0, w, 0.32, font_size=12, bold=True,
              color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        yy = 2.5
        for b in boxes:
            box(s, x, yy, w, 0.92, fill_color=C_CARD, line_color=col, line_width=Pt(1))
            _multiline(s, b, x + 0.05, yy + 0.14, w - 0.1, 0.7,
                       font_size=10.5, color=C_TEXT, align=PP_ALIGN.CENTER)
            yy += 1.05
        if name != "OUTPUT":
            txbox(s, "→", x + w, 3.1, 0.18, 0.5, font_size=22, bold=True,
                  color=C_TEAL, align=PP_ALIGN.CENTER)
        x += w + 0.22

    box(s, 0.4, 6.35, 12.5, 0.75, fill_color=C_TEAL_L, line_color=C_TEAL)
    txbox(s, "Backend selection (Stage 3.5): MMS-LID routes each language to its best ASR backend. "
             "All 7 languages now run on SeamlessM4T v2 — Hindi/Nepali/Pashto/Kashmiri with per-language "
             "LoRA adapters, Punjabi/Urdu/Mandarin zero-shot. Fine-tuned Whisper models are retained on disk "
             "for rollback only.",
          0.55, 6.42, 12.2, 0.62, font_size=10.5, color=C_NAVY)


def slide_features(prs, n):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Features", n, "Features")
    feats = [
        ("Multilingual ASR", "7 border languages auto-routed by LID; turbo fallback for others", C_TEAL),
        ("Audio Language ID", "MMS-LID 256-language model + Arabic-script cascade", C_TEAL),
        ("English Translation", "NLLB-200-distilled; IndicTrans2 for Dogri", C_BLUE),
        ("Threat Detection", "Keyword + entity spotting with a coded-terminology lexicon", C_GOLD),
        ("Intelligence Summary", "5W structured ISUM via Gemma 3 (Ollama), offline", C_PURPLE),
        ("Speaker Diarization", "Who-spoke-when; per-speaker tracks in 3-node mode", C_TEAL),
        ("Fully Offline", "Air-gapped runtime; no internet after model download", C_GREEN),
        ("GPU-Accelerated", "int8 CT2 + CPU parking to fit an 8 GB card", C_BLUE),
        ("Searchable Database", "SQLite transcript + summary store with full-text search", C_GOLD),
        ("Analyst GUI", "Map, dashboard, live-mic capture, network view", C_PURPLE),
    ]
    cols, cw, ch = 2, 6.15, 1.0
    x0, y0, gx, gy = 0.4, 1.55, 0.2, 0.1
    for i, (title, desc, col) in enumerate(feats):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        y = y0 + r * (ch + gy)
        box(s, x, y, cw, ch, fill_color=C_CARD, line_color=C_BORDER)
        box(s, x, y, 0.08, ch, fill_color=col)
        txbox(s, title, x + 0.25, y + 0.10, cw - 0.4, 0.4, font_size=13, bold=True, color=C_NAVY)
        txbox(s, desc, x + 0.25, y + 0.52, cw - 0.4, 0.42, font_size=10.5, color=C_SUB)


def slide_results_hero(prs, n):
    chart_slide(
        prs, "Result: SeamlessM4T surpasses fine-tuned Whisper on WER", n, "Results",
        report_charts.hero_backend_dumbbell(),
        caption="Deployed backend vs the fine-tuned Whisper model, per language (n=100 FLEURS held-out, same scorer). "
                "Every language improves; the arrow points to the deployed SeamlessM4T backend. Kashmiri is shown separately "
                "(different corpus + scoring ruler).",
        img_width=10.5, aspect=1.923,
    )


def slide_results_robustness_chart(prs, n):
    chart_slide(
        prs, "Result: the advantage holds — and widens — under radio degradation", n, "Results",
        report_charts.robustness_heatmap(),
        caption="SeamlessM4T's WER advantage over fine-tuned Whisper (percentage points), across five channel conditions on "
                "30 clips/language. Positive everywhere; Hindi (+5→+19.5) and Mandarin (+3.6→+17.2) widen most at 0 dB SNR. "
                "Conservative zero-shot comparison — the deployed hi/ne LoRA adapters improve on this further.",
        img_width=8.8, aspect=1.955,
    )


def slide_results_ks_ruler(prs, n):
    chart_slide(
        prs, "Result: Kashmiri — the WER gap was the scoring ruler, not the model", n, "Results",
        report_charts.ks_ruler_bars(),
        caption="Raw word-error rate once made Whisper look ahead of the first adapter, but Perso-Arabic references are densely "
                "diacritised and BOTH systems drop the marks — the gap was the ruler. Corrected, every successive adapter widens "
                "the win: the deployed r=128 ks_cloud3 (rented cloud GPU, plus a repaired vocabulary) leads on all three measures and "
                "won the 5-condition degradation sweep 5/5.",
        img_width=8.8, aspect=1.87,
    )


def slide_campaign_hours(prs, n):
    """Training hours + dataset sizes for the closing SM4T campaign (§4.4 of the PDF)."""
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "Training hours & dataset sizes — the closing adapter campaign", n, "Results")
    txbox(s, "The Whisper phase cost ~100 laptop GPU-hours. The adapter campaign that replaced it was cheaper per win — "
             "the two production adapters that closed it cost ~$6 of rented cloud compute, total.",
          0.4, 1.42, 12.5, 0.55, font_size=12.5, italic=True, color=C_NAVY)

    cols_x  = [0.30, 1.60, 3.05, 7.15, 8.05, 9.75, 10.95, 12.05]
    headers = ["Run", "Language", "Training data (clips / hours)", "LoRA r",
               "Hardware", "Wall time", "Result", "Status"]
    y = 2.15
    box(s, 0.28, y, 12.77, 0.5, fill_color=C_NAVY)
    for x, h in zip(cols_x, headers):
        txbox(s, h, x, y+0.12, 1.4, 0.3, font_size=9, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    rows = [
        ("ks_max2",  "Kashmiri", "97,456 / 240 h — humair025 + IV-R + OpenSLR-122",   "32",  "RTX 5060 8 GB",  "~19 h",    "61.88%", "Rollback",  C_SUB),
        ("ps_aug2",  "Pashto",   "32,656 / ~55 h — CV 30k + FLEURS ps_af x8",          "32",  "RTX 5060 8 GB",  "~5 h",     "37.46%", "Negative",  C_RED),
        ("ks_cloud", "Kashmiri", "144,749 / 335 h — same 3 sources, rebuilt (grown)",  "128", "A6000 48 GB ☁",  "6 h 35 m", "56.44%", "Rollback",  C_SUB),
        ("ks_cloud2","Kashmiri", "144,749 / 335 h — as ks_cloud, 2 epochs (patience 5)", "128", "A6000 48 GB ☁",  "~9 h",     "52.60%", "Rollback",  C_SUB),
        ("ks_cloud3","Kashmiri", "144,942 / 336 h — as ks_cloud2 + 20 vocab chars repaired","128", "A6000 48 GB ☁",  "~8 h",     "50.26%", "DEPLOYED",  C_GREEN),
        ("ks_cloud4","Kashmiri", "as ks_cloud3, warm start, token rows at 5x LR",         "128", "A6000 48 GB ☁",  "~1 h",     "50.69%", "Rejected",  C_RED),
        ("ps_cloud", "Pashto",   "18,656 / ~30 h — CV 10k + FLEURS ps_af x8",          "128", "A6000 48 GB ☁",  "1 h 32 m", "36.16%", "DEPLOYED",  C_GREEN),
        ("doi_iv",   "Dogri",    "IndicVoices-R Dogri — 97 shards, 43.8 GB",           "128", "A6000 48 GB ☁",  "~6 h",     "50.07%", "cut short", C_SUB),
        ("doi_iv2",  "Dogri",    "as doi_iv, fresh LR schedule, 9,000 steps",          "128", "A6000 48 GB ☁",  "~9 h",     "46.73%", "BEST",      C_GREEN),
    ]
    y += 0.52
    for i, (run, lang, data, r, hw, wt, res, status, col) in enumerate(rows):
        bg_col = C_CARD if i % 2 == 0 else C_CARD2
        box(s, 0.28, y, 12.77, 0.44, fill_color=bg_col)
        box(s, 0.28, y, 0.055, 0.44, fill_color=col)
        for x, v, bold in zip(cols_x, (run, lang, data, r, hw, wt, res, status),
                              (True, False, False, False, False, False, True, True)):
            txbox(s, v, x, y+0.10, 4.0 if x == 3.05 else 1.5, 0.28, font_size=9.5,
                  bold=bold, color=col if bold and x >= 10.9 else (C_TEXT if bold else C_SUB))
        y += 0.46
    y += 0.15
    for line in [
        "Kashmiri result = diacritic-normalised WER, 372-clip IndicVoices-R test  ·  Pashto result = clean FLEURS ps_af, n=100.",
        "Cloud runs: rented RTX A6000s at $0.53/hr — ks_cloud ≈ $4, ks_cloud2 ≈ $5, ks_cloud3 ≈ $4, ps_cloud ≈ $2. Corpus rebuilt from source on the box;",
        "humair025 had grown upstream → 144,749 clips / 335.4 h (131,868 unique sentences), 2–20 s filter, eval-leak blocklist applied.",
        "The decisive lever was CAPACITY (r=128 = 6.6% trainable vs r=32 = 1.75%) — more data at unchanged capacity regressed (ps_aug2).",
    ]:
        txbox(s, line, 0.4, y, 12.5, 0.3, font_size=9.5, italic=True, color=C_SUB)
        y += 0.28


def slide_finetune_summary(prs, n):
    s = blank_slide(prs)
    bg(s)
    slide_header(s, "From Whisper fine-tuning to SeamlessM4T — how the backend was chosen", n, "Results")
    txbox(s, "Seven Whisper LoRA models were fine-tuned first; a corrected evaluation then showed SeamlessM4T wins, "
             "and a targeted adapter campaign brought all seven languages onto it.",
          0.4, 1.42, 12.5, 0.55, font_size=12.5, italic=True, color=C_NAVY)

    stages = [
        ("1 · Fine-tune Whisper", "7 languages, LoRA r=8/16 on FLEURS + IndicVoices-R. Kashmiri needed a custom <|ks|> vocab token.", C_GOLD),
        ("2 · Correct the scoring", "Found the baseline was mislabelled turbo, and Mandarin WER was a whitespace artefact. Re-scored with one CJK-aware normaliser.", C_RED),
        ("3 · Head-to-head + robustness", "Zero-shot SeamlessM4T beat fine-tuned Whisper on 5/6 languages and held under bandpass/noise/codec.", C_TEAL),
        ("4 · Adapter campaign", "hi/ne gained from IndicVoices data; Pashto fell to noise-augmented training; Kashmiri to a trainable __kas__ token + a ruler correction.", C_BLUE),
        ("5 · Cloud capacity push", "Rented A6000s (~$11) retrained ps+ks at r=128 — beyond the 8 GB laptop: Pashto 36.91→36.16, Kashmiri 61.88→50.26 (3.84 pp from simply letting the run converge, then 2.34 pp from repairing the vocabulary). Both passed their gates and deployed.", C_TEAL),
        ("6 · Exhausting the levers", "Two honest negatives closed the campaign: a warm start improved validation loss but LOST 0.43 pp WER (5th loss/WER divergence), and beam search gains ~3 pp on clean speech while losing ~3 pp at 0 dB SNR. Production decoding unchanged.", C_RED),
        ("7 · Outcome", "All 7 languages route to SeamlessM4T. Kashmiri 65.19 → 50.26 on one held-out L2 ruler. What remains is a data/acoustic limit, not an optimisation one. Five scoring defects caught in total.", C_GREEN),
    ]
    y = 2.02
    for title, desc, col in stages:
        box(s, 0.4, y, 12.5, 0.68, fill_color=C_CARD, line_color=C_BORDER)
        box(s, 0.4, y, 0.08, 0.68, fill_color=col)
        txbox(s, title, 0.62, y + 0.06, 3.7, 0.56, font_size=11.5, bold=True, color=col)
        txbox(s, desc, 4.35, y + 0.04, 8.4, 0.60, font_size=9, color=C_SUB)
        y += 0.75


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_01_title(prs)                                                   # 1

    slide_section_divider(prs, 1, "Introduction &\nProblem Statement",
                          "What VANI is, and the operational problem it solves", 2)
    slide_intro_problem(prs, 3)                                           # 3

    slide_section_divider(prs, 2, "Deliverables",
                          "What was built and shipped", 4)
    slide_deliverables(prs, 5)                                            # 5

    slide_section_divider(prs, 3, "Architecture",
                          "The offline model stack, end to end", 6)
    slide_architecture(prs, 7)                                            # 7

    slide_section_divider(prs, 4, "Pipeline",
                          "How one intercept flows through ten stages", 8)
    slide_09_pipeline(prs, n=9)                                           # 9

    slide_section_divider(prs, 5, "Features",
                          "Analyst-facing capabilities", 10)
    slide_features(prs, 11)                                               # 11

    slide_section_divider(prs, 6, "Results",
                          "Backend selection, robustness, and the headline WER win", 12)
    slide_results_hero(prs, 13)                                          # 13  HERO
    slide_table_sm4t(prs, 14)                                            # 14  cross-model comparison
    slide_results_robustness_chart(prs, 15)                              # 15  robustness heatmap
    slide_11_robustness(prs, n=16)                                       # 16  robustness (LangID) table
    slide_results_ks_ruler(prs, 17)                                      # 17  Kashmiri ruler
    slide_finetune_summary(prs, 18)                                     # 18  training campaign (one slide)
    slide_campaign_hours(prs, 19)                                       # 19  training hours + dataset sizes
    slide_10_key_findings(prs, n=20)                                     # 20  key findings

    slide_section_divider(prs, 7, "Future Work",
                          "Where VANI goes next", 21)
    slide_next_steps_accuracy(prs, 22)                                   # 22
    slide_12_conclusion(prs, n=23)                                       # 23

    OUT_PATH.parent.mkdir(exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved -> {OUT_PATH}  ({OUT_PATH.stat().st_size // 1024} KB)")

if __name__ == "__main__":
    main()
