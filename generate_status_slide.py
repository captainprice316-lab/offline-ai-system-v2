#!/usr/bin/env python3
"""Generate a single project-status slide (name / deliverables / completed / pending+timeline)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY   = RGBColor(0x0E, 0x2A, 0x47)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)
GREEN  = RGBColor(0x1E, 0x7A, 0x46)
AMBER  = RGBColor(0xB5, 0x6A, 0x00)
GREY   = RGBColor(0x3A, 0x3A, 0x3A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def box(l, t, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill if fill else WHITE
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1.25)
    else:    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def add_text(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (txt, sz, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tb

# ── Title band ────────────────────────────────────────────────────────────────
box(0, 0, 13.333, 1.30, fill=NAVY)
add_text(0.35, 0.12, 12.6, 0.6,
         [("VANI — Voice Analysis & Neural Intelligence", 30, WHITE, True)])
add_text(0.35, 0.74, 12.6, 0.45,
         [("Offline AI-based Radio Intercept Analysis System   ·   M.Tech Research Project, IIT Indore   ·   Status as of 06 Jul 2026",
           13, RGBColor(0xC7, 0xD8, 0xEA), False)])

def header(l, t, w, text, color):
    box(l, t, w, 0.45, fill=color)
    add_text(l + 0.12, t + 0.04, w - 0.2, 0.4, [(text, 15, WHITE, True)],
             anchor=MSO_ANCHOR.MIDDLE)

col_w = 6.35
top_t = 1.55
body_h = 3.05

# ── Deliverables ──────────────────────────────────────────────────────────────
header(0.30, top_t, col_w, "DELIVERABLES", BLUE)
box(0.30, top_t + 0.45, col_w, body_h, fill=LIGHT)
deliv = [
    "Real-time multilingual ASR (fine-tuned Whisper)",
    "English translation (NLLB-200 / IndicTrans2)",
    "Sentiment analysis and reports",
    "Fully offline — no internet after setup",
]
add_text(0.50, top_t + 0.58, col_w - 0.35, body_h,
         [("• " + d, 13, GREY, False) for d in deliv], space=7)

# ── Completed work ────────────────────────────────────────────────────────────
header(6.85, top_t, col_w, "COMPLETED WORK", GREEN)
box(6.85, top_t + 0.45, col_w, body_h, fill=LIGHT)
done = [
    "7 languages fine-tuned via LoRA (pa, ps, ur, ne, zh, hi, ks)",
    "Full 10-stage pipeline verified end-to-end (06 Jul)",
    "Cross-model eval: Whisper vs SeamlessM4T (7 langs)",
    "Robustness eval: 7 langs × 5 noise/distortion conditions",
]
add_text(7.05, top_t + 0.58, col_w - 0.35, body_h,
         [("✓ " + d, 13, GREY, False) for d in done], space=7)

# ── Pending work + timeline (full width) ──────────────────────────────────────
pend_t = 5.30
header(0.30, pend_t, 12.73, "PENDING WORK  —  TIMELINE", AMBER)
box(0.30, pend_t + 0.45, 12.73, 1.45, fill=LIGHT)
pend = [
    "ASR improvement  →  Reduce WER for languages still above acceptable (~25%):",
    "                          Punjabi 55.7%,   Nepali 49.2%,   Pashto 38.6%",
    "Backlog          →  Kashmiri WER not meaningful; SeamlessM4T S2TT (ASR-only LoRA) to revisit",
]
runs = [(line, 13, GREY, False) for line in pend]
add_text(0.50, pend_t + 0.56, 12.4, 1.3, runs, space=8)

out = "docs/VANI_Status_Slide.pptx"
prs.save(out)
print("Saved:", out)
