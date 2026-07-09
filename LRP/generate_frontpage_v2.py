"""
Professional frontpage for VANI LRP — clean navy + gold on white.
Military/academic standard. No image centre piece.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x21, 0x37)   # deep navy
GOLD       = RGBColor(0xB8, 0x96, 0x0C)   # warm gold
GOLD_LIGHT = RGBColor(0xF5, 0xD9, 0x6E)   # pale gold (accent lines)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE   = RGBColor(0xF7, 0xF7, 0xF5)   # very light cream for main bg
DARKGREY   = RGBColor(0x2C, 0x2C, 0x2C)
MIDGREY    = RGBColor(0x55, 0x55, 0x55)
LIGHTGREY  = RGBColor(0xE8, 0xE8, 0xE8)

# ── Slide dimensions: A4 portrait ─────────────────────────────────────────────
W  = Inches(7.5)
H  = Inches(10.83)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(blank_layout)

# ── Helper: solid-fill shape ──────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    return shape

def add_textbox(slide, l, t, w, h, text, size, bold=False, colour=WHITE,
                align=PP_ALIGN.CENTER, italic=False, space_before=0, space_after=0):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return txb

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 0: Off-white full background
# ─────────────────────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 7.5, 10.83, OFFWHITE)

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 1: Outer border — thin navy rectangle (2-line frame feel)
#          We draw two nested rect borders via shape.line
# ─────────────────────────────────────────────────────────────────────────────
def add_border_rect(slide, l, t, w, h, line_rgb, line_width_pt, fill=None):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.background() if fill is None else (shape.fill.solid(), setattr(shape.fill.fore_color, 'rgb', fill))
    shape.line.color.rgb = line_rgb
    shape.line.width = _Pt(line_width_pt)
    return shape

add_border_rect(slide, 0.15, 0.15, 7.20, 10.53, NAVY,  2.0)   # outer frame
add_border_rect(slide, 0.22, 0.22, 7.06, 10.39, GOLD,  0.75)  # inner gold trim

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 2: Navy header bar (top)
# ─────────────────────────────────────────────────────────────────────────────
add_rect(slide, 0.22, 0.22, 7.06, 2.05, NAVY)

# Thin gold accent line below header
add_rect(slide, 0.22, 2.27, 7.06, 0.05, GOLD)

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 3: Logos inside header
# ─────────────────────────────────────────────────────────────────────────────
LOGO_H = 0.80
LOGO_TOP = 0.35

# Left logo (Picture_31 — small ~14KB, likely IIT or unit badge)
try:
    slide.shapes.add_picture(str(Path(__file__).parent / 'Picture_31.png'),
                              Inches(0.40), Inches(LOGO_TOP),
                              height=Inches(LOGO_H))
except Exception:
    pass

# Right logo (Picture_27 — larger ~423KB)
try:
    slide.shapes.add_picture(str(Path(__file__).parent / 'Picture_27.png'),
                              Inches(5.90), Inches(LOGO_TOP),
                              height=Inches(LOGO_H))
except Exception:
    pass

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 4: Title text inside header
# ─────────────────────────────────────────────────────────────────────────────
# "VANI" large
txb = slide.shapes.add_textbox(Inches(1.50), Inches(0.28), Inches(4.50), Inches(0.70))
tf = txb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "VANI"
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = "Tahoma"

# "VOICE ANALYSIS AND NEURAL INTELLIGENCE" in gold
txb2 = slide.shapes.add_textbox(Inches(0.50), Inches(0.98), Inches(6.50), Inches(0.45))
tf2 = txb2.text_frame
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "VOICE ANALYSIS AND NEURAL INTELLIGENCE"
run2.font.size = Pt(11)
run2.font.bold = True
run2.font.color.rgb = GOLD_LIGHT
run2.font.name = "Tahoma"

# Subtitle line
txb3 = slide.shapes.add_textbox(Inches(0.50), Inches(1.44), Inches(6.50), Inches(0.38))
tf3 = txb3.text_frame
p3 = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "Military-Grade Offline Radio Intercept Analysis System"
run3.font.size = Pt(9)
run3.font.italic = True
run3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
run3.font.name = "Calibri"

# Divider rule inside header
add_rect(slide, 0.50, 1.83, 6.50, 0.025, GOLD_LIGHT)

# LRP label
txb4 = slide.shapes.add_textbox(Inches(0.50), Inches(1.86), Inches(6.50), Inches(0.33))
tf4 = txb4.text_frame
p4 = tf4.paragraphs[0]
p4.alignment = PP_ALIGN.CENTER
run4 = p4.add_run()
run4.text = "LITERATURE REVIEW PROPOSAL (LRP)"
run4.font.size = Pt(9)
run4.font.bold = True
run4.font.color.rgb = RGBColor(0xA0, 0xB4, 0xC8)
run4.font.name = "Calibri"

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 5: Topic box
# ─────────────────────────────────────────────────────────────────────────────
# Gold accent strip on left edge of topic area
add_rect(slide, 0.32, 2.38, 0.07, 1.50, GOLD)

# Topic label
txb5 = slide.shapes.add_textbox(Inches(0.52), Inches(2.38), Inches(6.60), Inches(0.28))
tf5 = txb5.text_frame
p5 = tf5.paragraphs[0]
p5.alignment = PP_ALIGN.LEFT
run5 = p5.add_run()
run5.text = "RESEARCH TOPIC"
run5.font.size = Pt(7.5)
run5.font.bold = True
run5.font.color.rgb = GOLD
run5.font.name = "Calibri"

# Topic text
txb6 = slide.shapes.add_textbox(Inches(0.52), Inches(2.65), Inches(6.60), Inches(1.10))
tf6 = txb6.text_frame
tf6.word_wrap = True
p6 = tf6.paragraphs[0]
p6.alignment = PP_ALIGN.LEFT
run6 = p6.add_run()
run6.text = "A Comparative Study of Speech-to-Text Models\nfor Noisy Radio Transmission"
run6.font.size = Pt(17)
run6.font.bold = True
run6.font.color.rgb = NAVY
run6.font.name = "Tahoma"

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 6: Thin gold divider
# ─────────────────────────────────────────────────────────────────────────────
add_rect(slide, 0.32, 3.90, 6.86, 0.035, GOLD)

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 7: Three-column info cards (Course | Category | Date)
# ─────────────────────────────────────────────────────────────────────────────
cards = [
    ("COURSE",    "SOATE-44",          0.32),
    ("CATEGORY",  "SIGINT / ASR / NLP", 2.82),
    ("DATE",      "28 March 2026",      5.32),
]
CARD_W = 2.32
CARD_TOP = 4.00

for label, value, left in cards:
    # card background (very light navy tint)
    box = slide.shapes.add_shape(1, Inches(left), Inches(CARD_TOP),
                                  Inches(CARD_W), Inches(0.80))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEA, 0xEF, 0xF5)
    box.line.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
    box.line.width = Pt(0.5)

    # label
    lbl = slide.shapes.add_textbox(Inches(left+0.08), Inches(CARD_TOP+0.06),
                                    Inches(CARD_W-0.16), Inches(0.24))
    lp = lbl.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = label
    lr.font.size = Pt(6.5)
    lr.font.bold = True
    lr.font.color.rgb = GOLD
    lr.font.name = "Calibri"

    # value
    val = slide.shapes.add_textbox(Inches(left+0.08), Inches(CARD_TOP+0.30),
                                    Inches(CARD_W-0.16), Inches(0.42))
    vp = val.text_frame.paragraphs[0]
    vp.alignment = PP_ALIGN.CENTER
    vr = vp.add_run()
    vr.text = value
    vr.font.size = Pt(10)
    vr.font.bold = True
    vr.font.color.rgb = NAVY
    vr.font.name = "Tahoma"

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 8: Submission details table (python-pptx native table)
# ─────────────────────────────────────────────────────────────────────────────
from pptx.util import Inches as _I

TABLE_TOP    = 4.95
TABLE_LEFT   = 0.32
TABLE_WIDTH  = 6.86
TABLE_HEIGHT = 2.60

rows_data = [
    ("Submitted By",   "Lt Col Vishal Sharma",                          True),
    ("Course",         "SOATE-44, IIT Indore",                          False),
    ("Guide (MCTE)",   "Dr. Krishan Berwal, MCTE Mhow",                 False),
    ("Guide (IIT)",    "Dr. Chandresh Maurya, IIT Indore (CSE Dept.)",  False),
    ("Institution",    "Military College of Telecommunication Engineering (MCTE), Mhow", False),
    ("Date",           "28 March 2026",                                  False),
]

table = slide.shapes.add_table(
    len(rows_data), 2,
    _I(TABLE_LEFT), _I(TABLE_TOP),
    _I(TABLE_WIDTH), _I(TABLE_HEIGHT)
).table

# Column widths
table.columns[0].width = _I(1.80)
table.columns[1].width = _I(5.06)

def style_cell(cell, text, font_size, bold, colour, bg_rgb, align=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = "Calibri"
    # fill
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _et
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for s in tcPr.findall(_qn('w:shd')):
        tcPr.remove(s)
    # use a:solidFill approach via XML
    for existing in tcPr.findall(_qn('a:solidFill')):
        tcPr.remove(existing)
    sf = _et.SubElement(tcPr, _qn('a:solidFill'))
    srgb = _et.SubElement(sf, _qn('a:srgbClr'))
    srgb.set('val', f'{bg_rgb[0]:02X}{bg_rgb[1]:02X}{bg_rgb[2]:02X}')

# Row heights — distribute evenly
row_h = _I(TABLE_HEIGHT / len(rows_data))
for i, row in enumerate(table.rows):
    row.height = row_h
    label, value, highlight = rows_data[i]
    # label cell
    if highlight:
        lbl_bg = (0x0D, 0x21, 0x37)   # navy for "Submitted By"
        lbl_fg = WHITE
    else:
        lbl_bg = (0xD8, 0xE4, 0xEF)
        lbl_fg = NAVY
    style_cell(row.cells[0], label, 8.5, True,  lbl_fg, lbl_bg)

    if highlight:
        val_bg = (0x1A, 0x3A, 0x5C)
        val_fg = GOLD_LIGHT
        val_bold = True
        val_size = 10
    else:
        val_bg = (0xF7, 0xF7, 0xF5)
        val_fg = DARKGREY
        val_bold = False
        val_size = 9
    style_cell(row.cells[1], value, val_size, val_bold, val_fg, val_bg)

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 9: Second gold divider
# ─────────────────────────────────────────────────────────────────────────────
add_rect(slide, 0.32, 7.62, 6.86, 0.035, GOLD)

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 10: Abstract / scope blurb
# ─────────────────────────────────────────────────────────────────────────────
add_rect(slide, 0.32, 7.68, 0.07, 1.88, NAVY)   # left navy strip

txb7 = slide.shapes.add_textbox(Inches(0.52), Inches(7.70), Inches(6.60), Inches(0.25))
p7 = txb7.text_frame.paragraphs[0]
run7 = p7.add_run()
run7.text = "ABSTRACT"
run7.font.size = Pt(7.5)
run7.font.bold = True
run7.font.color.rgb = NAVY
run7.font.name = "Calibri"

txb8 = slide.shapes.add_textbox(Inches(0.52), Inches(7.95), Inches(6.60), Inches(1.55))
tf8 = txb8.text_frame
tf8.word_wrap = True
p8 = tf8.paragraphs[0]
p8.alignment = PP_ALIGN.JUSTIFY
run8 = p8.add_run()
run8.text = (
    "This Literature Review Proposal examines state-of-the-art Speech-to-Text (STT) models "
    "for deployment in degraded military radio communication environments. The study "
    "evaluates sixteen seminal works spanning transformer-based architectures, self-supervised "
    "pre-training, noise-robust augmentation, multilingual Indic language recognition, and "
    "downstream NLP pipelines. Findings directly inform the design of VANI — an offline, "
    "CPU-optimised SIGINT intercept analysis system targeting low-resource multilingual "
    "scenarios prevalent in border-region operations."
)
run8.font.size = Pt(8.5)
run8.font.color.rgb = MIDGREY
run8.font.name = "Calibri"

# ─────────────────────────────────────────────────────────────────────────────
# LAYER 11: Navy footer bar
# ─────────────────────────────────────────────────────────────────────────────
FOOTER_TOP = 9.62
add_rect(slide, 0.22, FOOTER_TOP, 7.06, 0.68, NAVY)

# Gold accent line above footer
add_rect(slide, 0.22, FOOTER_TOP - 0.04, 7.06, 0.04, GOLD)

# Footer text left
txb9 = slide.shapes.add_textbox(Inches(0.40), Inches(FOOTER_TOP + 0.08),
                                  Inches(3.5), Inches(0.50))
p9 = txb9.text_frame.paragraphs[0]
p9.alignment = PP_ALIGN.LEFT
r9 = p9.add_run()
r9.text = "SOATE-44  |  IIT Indore  |  MCTE Mhow"
r9.font.size = Pt(7.5)
r9.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
r9.font.name = "Calibri"

# Footer text right — classification marker
txb10 = slide.shapes.add_textbox(Inches(4.10), Inches(FOOTER_TOP + 0.08),
                                   Inches(3.00), Inches(0.50))
p10 = txb10.text_frame.paragraphs[0]
p10.alignment = PP_ALIGN.RIGHT
r10 = p10.add_run()
r10.text = "FOR ACADEMIC USE  |  UNCLASSIFIED"
r10.font.size = Pt(7.5)
r10.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
r10.font.name = "Calibri"

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
OUT = Path(__file__).parent / "frontpage_new.pptx"
prs.save(str(OUT))
print(f"Saved -> {OUT}")
