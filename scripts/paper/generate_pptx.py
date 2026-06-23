"""
VANI LRP PowerPoint Generator
Generates VANI_LRP_Comparative_STT.pptx — 20 slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt
import copy
from lxml import etree

# ── Colour constants ───────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1F, 0x38, 0x64)
STEEL_BLUE  = RGBColor(0x2E, 0x74, 0xB5)
GOLD        = RGBColor(0xC9, 0xA8, 0x4C)
LIGHT_BLUE  = RGBColor(0xDE, 0xEA, 0xF1)
LIGHT_RED   = RGBColor(0xFC, 0xE4, 0xD6)
LIGHT_GREEN = RGBColor(0xE2, 0xEF, 0xDA)
LIGHT_GOLD  = RGBColor(0xFF, 0xF2, 0xCC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0x80, 0x80, 0x80)
RED         = RGBColor(0xC0, 0x00, 0x00)
AMBER       = RGBColor(0xFF, 0xC0, 0x00)
GREEN_DARK  = RGBColor(0x37, 0x5C, 0x23)
BLACK       = RGBColor(0x00, 0x00, 0x00)

# ── Slide dimensions ───────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ══════════════════════════════════════════════════════════════════════════════
# Helper functions
# ══════════════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def set_shape_bg(shape, fill_color, line_color=None, line_width_pt=0.75):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    if line_color:
        ln = shape.line
        ln.color.rgb = line_color
        ln.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()  # no line

def add_textbox(slide, text, left, top, width, height,
                font_name="Calibri", font_size=14, bold=False, italic=False,
                color=BLACK, align=PP_ALIGN.LEFT,
                fill_color=None, line_color=None,
                word_wrap=True, margin_pt=5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Pt(margin_pt)
    tf.margin_right  = Pt(margin_pt)
    tf.margin_top    = Pt(margin_pt)
    tf.margin_bottom = Pt(margin_pt)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if fill_color:
        set_shape_bg(txBox, fill_color, line_color)
    elif line_color:
        txBox.line.color.rgb = line_color
        txBox.line.width = Pt(0.75)
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width_pt=0.75):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    set_shape_bg(shape, fill_color, line_color, line_width_pt)
    shape.line.fill.background() if not line_color else None
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color,
                     line_color=None, line_width_pt=0.75, radius_pct=10):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    # Use rounded rectangle shape (freeform shape type 5)
    shape = slide.shapes.add_shape(
        5,  # msoShapeRoundedRectangle
        left, top, width, height
    )
    # Set corner radius
    shape.adjustments[0] = radius_pct / 100.0
    set_shape_bg(shape, fill_color, line_color, line_width_pt)
    if not line_color:
        shape.line.fill.background()
    return shape

def add_title_line(slide, title_text, font_size=28):
    """Add standard slide title + horizontal rule."""
    # Title text box
    tb = add_textbox(slide, title_text,
                     left=Inches(0.4), top=Inches(0.15),
                     width=Inches(12.5), height=Inches(0.55),
                     font_size=font_size, bold=True, color=NAVY,
                     align=PP_ALIGN.LEFT)
    # Horizontal line shape (thin navy rectangle)
    line = add_rect(slide,
                    left=Inches(0.4), top=Inches(0.75),
                    width=Inches(12.5), height=Inches(0.03),
                    fill_color=NAVY)
    return tb

def add_slide_number(slide, num):
    add_textbox(slide, str(num),
                left=Inches(12.7), top=Inches(7.1),
                width=Inches(0.5), height=Inches(0.3),
                font_size=9, color=GREY, align=PP_ALIGN.RIGHT)

def add_bullet_para(tf, text, font_size=13, bold=False, color=BLACK,
                    level=0, bullet_char="•"):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.level = level
    p.alignment = PP_ALIGN.LEFT
    indent = "  " * level
    run = p.add_run()
    run.text = f"{indent}{bullet_char} {text}"
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def add_bullets_box(slide, bullets, left, top, width, height,
                    font_size=13, fill_color=None, line_color=None,
                    title=None, title_size=14, title_color=NAVY,
                    bullet_color=BLACK, bold_first=False):
    """Add a text box with multiple bullet lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Pt(6)
    tf.margin_right  = Pt(6)
    tf.margin_top    = Pt(4)
    tf.margin_bottom = Pt(4)
    if fill_color:
        set_shape_bg(txBox, fill_color, line_color)
    elif line_color:
        txBox.line.color.rgb = line_color
        txBox.line.width = Pt(0.75)
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.name = "Calibri"
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        first = False
    for i, bullet in enumerate(bullets):
        if first and i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.bold = (bold_first and i == 0)
        run.font.color.rgb = bullet_color
    return txBox

def set_cell_bg(cell, rgb_color):
    """Set table cell background color."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove existing fills first
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr  = etree.SubElement(solidFill, qn('a:srgbClr'))
    hex_val = f'{rgb_color[0]:02X}{rgb_color[1]:02X}{rgb_color[2]:02X}'
    srgbClr.set('val', hex_val)

def set_cell_text(cell, text, font_size=10, bold=False,
                  color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # clear existing runs
    for run in p.runs:
        run.text = ""
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()

# Navy background
bg = add_rect(slide, 0, 0, W, H, NAVY)

# Gold accent horizontal bar (middle)
add_rect(slide,
         left=Inches(0.6), top=Inches(3.3),
         width=Inches(12.0), height=Inches(0.06),
         fill_color=GOLD)

# Main title
tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.0), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "A Comparative Study of Speech-to-Text Models"
run.font.name = "Calibri"; run.font.size = Pt(32); run.font.bold = True
run.font.color.rgb = WHITE

tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12.0), Inches(0.7))
tf2 = tb2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "for Noisy Radio Transmission"
run2.font.name = "Calibri"; run2.font.size = Pt(32); run2.font.bold = True
run2.font.color.rgb = WHITE

# Subtitle
tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(3.55), Inches(12.0), Inches(0.5))
tf3 = tb3.text_frame
p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "Literature Review Proposal  |  SOATE-44"
run3.font.name = "Calibri"; run3.font.size = Pt(20)
run3.font.color.rgb = GOLD

# Line 3
tb4 = slide.shapes.add_textbox(Inches(0.6), Inches(4.15), Inches(12.0), Inches(0.5))
tf4 = tb4.text_frame
p4 = tf4.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
run4 = p4.add_run()
run4.text = "VANI – Voice Analysis & Neural Intelligence"
run4.font.name = "Calibri"; run4.font.size = Pt(18); run4.font.italic = True
run4.font.color.rgb = WHITE

# Line 4
tb5 = slide.shapes.add_textbox(Inches(0.6), Inches(4.75), Inches(12.0), Inches(0.4))
tf5 = tb5.text_frame
p5 = tf5.paragraphs[0]; p5.alignment = PP_ALIGN.CENTER
run5 = p5.add_run()
run5.text = "March 2026"
run5.font.name = "Calibri"; run5.font.size = Pt(16)
run5.font.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Presentation Outline")
add_slide_number(slide, 2)

items = [
    "1.  Introduction & Background",
    "2.  Problem Statement",
    "3.  Research Objectives",
    "4.  Literature Overview (16 Papers)",
    "5.  Cluster 1: Foundational Architecture",
    "6.  Cluster 2: STT Architecture Evolution",
    "7.  Cluster 3: Noise Robustness & Enhancement",
    "8.  Cluster 4: Multilingual & Indic Language ASR",
    "9.  Cluster 5: Downstream NLP Pipeline",
    "10. Research Gaps",
    "11. Expected Findings",
    "12. Recommendations",
    "13. References",
]

col1 = items[:7]
col2 = items[7:]

txBox1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(5.8), Inches(6.0))
tf1 = txBox1.text_frame; tf1.word_wrap = True
for i, item in enumerate(col1):
    p = tf1.paragraphs[0] if i == 0 else tf1.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run(); run.text = item
    run.font.name = "Calibri"; run.font.size = Pt(15)
    run.font.color.rgb = NAVY if i == 0 else BLACK
    run.font.bold = (i == 0)

txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(0.9), Inches(5.8), Inches(6.0))
tf2 = txBox2.text_frame; tf2.word_wrap = True
for i, item in enumerate(col2):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run(); run.text = item
    run.font.name = "Calibri"; run.font.size = Pt(15)
    run.font.color.rgb = BLACK

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introduction: The SIGINT Challenge
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "The Problem: Radio Interception at Scale")
add_slide_number(slide, 3)

bullets = [
    "Military tactical communications still dominated by VHF/UHF radio — the primary SIGINT collection medium",
    "Indian subcontinent operational theatre: 22+ constitutionally scheduled languages, 5 active scripts (Devanagari, Gurmukhi, Arabic, Latin, Bengali)",
    "Radio transmissions suffer severe acoustic degradation: AWGN, frequency-selective fading, squelch artifacts, codec compression, burst interference (SNR: -5dB to +15dB)",
    "Human analyst transcription: slow, expensive, not scalable to operational SIGINT volumes",
    "Cloud ASR (Google, Azure): not viable for classified/offline environments",
    "No existing comparative benchmark for STT models on Indic-language noisy radio audio",
]

txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(12.5), Inches(5.6))
tf = txBox.text_frame; tf.word_wrap = True
tf.margin_left = Pt(6); tf.margin_top = Pt(4)
for i, b in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run(); run.text = f"• {b}"
    run.font.name = "Calibri"; run.font.size = Pt(13)
    run.font.color.rgb = BLACK

# Bottom note box
note = add_rounded_rect(slide,
    left=Inches(0.5), top=Inches(6.55),
    width=Inches(12.3), height=Inches(0.65),
    fill_color=LIGHT_BLUE, line_color=NAVY)
tb_note = slide.shapes.add_textbox(Inches(0.6), Inches(6.58), Inches(12.1), Inches(0.6))
tf_n = tb_note.text_frame; tf_n.word_wrap = True
p_n = tf_n.paragraphs[0]; p_n.alignment = PP_ALIGN.CENTER
r_n = p_n.add_run()
r_n.text = "This study conducts a systematic comparison of leading STT architectures under these specific conditions"
r_n.font.name = "Calibri"; r_n.font.size = Pt(12); r_n.font.italic = True
r_n.font.color.rgb = NAVY

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Introduction: VANI System Context
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "VANI: The Implementation Framework")
add_slide_number(slide, 4)

# Left column heading
add_textbox(slide, "System Constraints & Capabilities",
    Inches(0.4), Inches(0.9), Inches(5.8), Inches(0.4),
    font_size=13, bold=True, color=STEEL_BLUE)

left_bullets = [
    "Offline, CPU-only, 8 GB RAM constraint",
    "Fully air-gapped — no internet dependency",
    "19 operational languages supported",
    "Complete pipeline: VAD → ASR → LangID → MT → Keywords → ISUM",
    "8-tab Streamlit analyst interface",
]
add_bullets_box(slide, left_bullets,
    Inches(0.4), Inches(1.35), Inches(5.8), Inches(3.8),
    font_size=13, fill_color=LIGHT_BLUE, line_color=NAVY)

# Right column heading
add_textbox(slide, "Pipeline Flow",
    Inches(7.0), Inches(0.9), Inches(5.9), Inches(0.4),
    font_size=13, bold=True, color=STEEL_BLUE)

# Pipeline flow boxes
pipeline_steps = [
    ("Audio Input", NAVY, WHITE),
    ("VAD (Voice Activity Detection)", LIGHT_BLUE, NAVY),
    ("ASR (Whisper large-v3-turbo)", LIGHT_BLUE, NAVY),
    ("3-way Language Identification", LIGHT_BLUE, NAVY),
    ("Translation (NLLB-200 / IndicTrans2)", LIGHT_BLUE, NAVY),
    ("ISUM Report Generation", LIGHT_BLUE, NAVY),
    ("PDF / DOCX Export", NAVY, WHITE),
]
step_h = 0.52
for i, (txt, bg, fg) in enumerate(pipeline_steps):
    top = Inches(1.35 + i * (step_h + 0.06))
    rect = add_rounded_rect(slide,
        Inches(7.2), top, Inches(5.5), Inches(step_h),
        fill_color=bg, line_color=NAVY, radius_pct=8)
    tb = slide.shapes.add_textbox(Inches(7.25), top + Pt(2), Inches(5.4), Inches(step_h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.name = "Calibri"; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = fg
    if i < len(pipeline_steps) - 1:
        # Arrow: small triangle indicator
        arr = slide.shapes.add_textbox(Inches(9.5), top + Inches(step_h - 0.01),
                                        Inches(1.0), Inches(0.12))
        tf_a = arr.text_frame; p_a = tf_a.paragraphs[0]
        p_a.alignment = PP_ALIGN.CENTER
        r_a = p_a.add_run(); r_a.text = "▼"
        r_a.font.size = Pt(8); r_a.font.color.rgb = NAVY

# Bottom note
add_textbox(slide,
    "VANI serves as the test-bed for the comparative STT evaluation proposed in this study",
    Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
    font_size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Problem Statement")
add_slide_number(slide, 5)

# Quote box
quote_rect = add_rounded_rect(slide,
    Inches(0.5), Inches(0.95), Inches(12.3), Inches(2.6),
    fill_color=RGBColor(0xF2, 0xF7, 0xFF),
    line_color=NAVY, line_width_pt=2.0)

tb_q = slide.shapes.add_textbox(Inches(0.65), Inches(1.05), Inches(12.0), Inches(2.4))
tf_q = tb_q.text_frame; tf_q.word_wrap = True
p_q = tf_q.paragraphs[0]; p_q.alignment = PP_ALIGN.CENTER
r_q = p_q.add_run()
r_q.text = ('"Existing Speech-to-Text systems are optimised for clean audio and high-resource languages. '
            'Military radio transmissions are characterised by SNR values of -5dB to +15dB and are conducted '
            'in low-resource Indic languages. No systematic comparative evaluation exists for this specific '
            'combination of acoustic and linguistic challenges."')
r_q.font.name = "Calibri"; r_q.font.size = Pt(14); r_q.font.italic = True
r_q.font.color.rgb = NAVY

# 3 challenge boxes
boxes = [
    ("Acoustic Challenge", "Noise, fading, codec distortion", RED, LIGHT_RED),
    ("Linguistic Challenge", "19 languages, 5 scripts, code-switching", RGBColor(0xFF,0x8C,0x00), RGBColor(0xFF,0xF0,0xCC)),
    ("Operational Challenge", "Offline, 8 GB RAM, real-time", RGBColor(0x37,0x5C,0x23), LIGHT_GREEN),
]
box_w = Inches(3.8)
gap   = Inches(0.2)
for i, (title, desc, border, bg) in enumerate(boxes):
    left = Inches(0.5) + i * (box_w + gap)
    rect = add_rounded_rect(slide, left, Inches(3.85), box_w, Inches(1.8),
                             fill_color=bg, line_color=border, line_width_pt=2.0)
    tb = slide.shapes.add_textbox(left + Inches(0.1), Inches(3.95), box_w - Inches(0.2), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    p_t = tf.paragraphs[0]; p_t.alignment = PP_ALIGN.CENTER
    r_t = p_t.add_run(); r_t.text = title
    r_t.font.name = "Calibri"; r_t.font.size = Pt(14); r_t.font.bold = True
    r_t.font.color.rgb = border
    p_d = tf.add_paragraph(); p_d.alignment = PP_ALIGN.CENTER
    p_d.space_before = Pt(6)
    r_d = p_d.add_run(); r_d.text = desc
    r_d.font.name = "Calibri"; r_d.font.size = Pt(12)
    r_d.font.color.rgb = BLACK

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Research Objectives
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Research Objectives")
add_slide_number(slide, 6)

objectives_left = [
    ("1.", "Review leading STT architectures (CTC-RNN, self-supervised, encoder-decoder, convolution-augmented)"),
    ("2.", "Analyse noise robustness mechanisms and performance under degraded audio"),
    ("3.", "Evaluate multilingual capability for Indic and low-resource languages"),
]
objectives_right = [
    ("4.", "Compare computational requirements for CPU-only deployment"),
    ("5.", "Identify optimal architecture(s) for offline multilingual radio intercept pipeline"),
    ("6.", "Establish WER/CER baselines and evaluation framework for domain benchmarking"),
]

def draw_obj_box(slide, num, text, left, top, width):
    rect = add_rounded_rect(slide, left, top, width, Inches(1.1),
                             fill_color=LIGHT_BLUE, line_color=NAVY, line_width_pt=1.0)
    tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.05),
                                   width - Inches(0.2), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = f"{num}  {text}"
    r.font.name = "Calibri"; r.font.size = Pt(12)
    r.font.color.rgb = NAVY; r.font.bold = False

gap_v = Inches(0.15)
obj_h_start = Inches(0.95)
obj_w = Inches(6.1)
for i, (num, txt) in enumerate(objectives_left):
    draw_obj_box(slide, num, txt, Inches(0.4), obj_h_start + i*(Inches(1.1)+gap_v), obj_w)
for i, (num, txt) in enumerate(objectives_right):
    draw_obj_box(slide, num, txt, Inches(6.8), obj_h_start + i*(Inches(1.1)+gap_v), obj_w)

# Sub-objectives
add_textbox(slide, "Sub-objectives:",
    Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.35),
    font_size=12, bold=True, color=STEEL_BLUE)

sub_obj = [
    "Analyse SpecAugment and SEGAN as noise augmentation / enhancement strategies",
    "Evaluate 3-way LangID ensemble (Whisper + FastText + MMS-LID)",
    "Identify future directions: domain fine-tuning, cross-lingual transfer, speaker diarisation",
]
txBox = slide.shapes.add_textbox(Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.0))
tf = txBox.text_frame; tf.word_wrap = True
for i, s in enumerate(sub_obj):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run(); run.text = f"→  {s}"
    run.font.name = "Calibri"; run.font.size = Pt(12)
    run.font.color.rgb = BLACK

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Literature Overview: 16 Papers Thematic Map
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "16 Papers Reviewed — Thematic Map")
add_slide_number(slide, 7)

clusters = [
    {
        "name": "CLUSTER 1\nFoundation",
        "color": RGBColor(0xE2, 0xEF, 0xDA),
        "border": RGBColor(0x37, 0x5C, 0x23),
        "papers": ["[P6] Transformer\nVaswani et al. 2017"],
        "new": [],
    },
    {
        "name": "CLUSTER 2\nSTT Architectures",
        "color": LIGHT_BLUE,
        "border": NAVY,
        "papers": ["[P14] Deep Speech 2\nAmodei et al. 2016",
                   "[P11] wav2vec 2.0 ★\nBaevski et al. 2020",
                   "[P15] HuBERT ★\nHsu et al. 2021",
                   "[P12] Conformer ★\nGulati et al. 2020",
                   "[P1] Whisper\nRadford et al. 2022"],
        "new": [1, 2, 3],
    },
    {
        "name": "CLUSTER 3\nNoise Robustness",
        "color": RGBColor(0xFF, 0xF0, 0xCC),
        "border": RGBColor(0xFF, 0x8C, 0x00),
        "papers": ["[P13] SpecAugment ★\nPark et al. 2019",
                   "[P16] SEGAN ★\nPascual et al. 2017"],
        "new": [0, 1],
    },
    {
        "name": "CLUSTER 4\nMultilingual Indic",
        "color": RGBColor(0xFC, 0xE4, 0xD6),
        "border": RED,
        "papers": ["[P7] IndicWav2Vec\nJaved et al. 2022",
                   "[P4] MMS\nPratap et al. 2023"],
        "new": [],
    },
    {
        "name": "CLUSTER 5\nDownstream NLP",
        "color": RGBColor(0xED, 0xE7, 0xF6),
        "border": RGBColor(0x6A, 0x1B, 0x9A),
        "papers": ["[P5] FastText\nJoulin et al. 2017",
                   "[P2] NLLB-200\nCosta-jussà et al. 2022",
                   "[P3] IndicTrans2\nGala et al. 2023",
                   "[P8] XLM-R ★\nConneau et al. 2020",
                   "[P9] PEGASUS ★\nZhang et al. 2020",
                   "[P10] pyannote ★\nBredin et al. 2020"],
        "new": [3, 4, 5],
    },
]

cluster_w = Inches(2.45)
cluster_left_start = Inches(0.3)
cluster_gap = Inches(0.12)
cluster_top = Inches(0.9)
cluster_total_h = Inches(6.35)

for ci, cl in enumerate(clusters):
    cl_left = cluster_left_start + ci * (cluster_w + cluster_gap)
    # Cluster header
    hdr = add_rounded_rect(slide, cl_left, cluster_top, cluster_w, Inches(0.6),
                            fill_color=cl["border"], line_color=None, radius_pct=8)
    tb_h = slide.shapes.add_textbox(cl_left + Inches(0.05), cluster_top + Inches(0.05),
                                     cluster_w - Inches(0.1), Inches(0.55))
    tf_h = tb_h.text_frame; tf_h.word_wrap = True
    p_h = tf_h.paragraphs[0]; p_h.alignment = PP_ALIGN.CENTER
    r_h = p_h.add_run(); r_h.text = cl["name"].replace("\n", " ")
    r_h.font.name = "Calibri"; r_h.font.size = Pt(10); r_h.font.bold = True
    r_h.font.color.rgb = WHITE

    # Paper boxes
    n_papers = len(cl["papers"])
    paper_h = (cluster_total_h - Inches(0.7)) / max(n_papers, 1) - Inches(0.08)
    paper_h = min(paper_h, Inches(0.95))

    for pi, paper in enumerate(cl["papers"]):
        p_top = cluster_top + Inches(0.68) + pi * (paper_h + Inches(0.07))
        is_new = pi in cl["new"]
        p_fill = RGBColor(0xFF, 0xF8, 0xE7) if is_new else cl["color"]
        p_border = GOLD if is_new else cl["border"]
        prect = add_rounded_rect(slide, cl_left + Inches(0.05), p_top,
                                  cluster_w - Inches(0.1), paper_h,
                                  fill_color=p_fill, line_color=p_border,
                                  line_width_pt=1.0 if not is_new else 1.5,
                                  radius_pct=6)
        tb_p = slide.shapes.add_textbox(cl_left + Inches(0.1), p_top + Inches(0.04),
                                         cluster_w - Inches(0.2), paper_h - Inches(0.05))
        tf_p = tb_p.text_frame; tf_p.word_wrap = True
        p_p = tf_p.paragraphs[0]; p_p.alignment = PP_ALIGN.CENTER
        r_p = p_p.add_run(); r_p.text = paper
        r_p.font.name = "Calibri"; r_p.font.size = Pt(9)
        r_p.font.bold = is_new
        r_p.font.color.rgb = NAVY if not is_new else RGBColor(0x7B, 0x5C, 0x00)

# Legend
add_textbox(slide, "★ = New papers (P11-P16) added in this review",
    Inches(0.3), Inches(7.18), Inches(8.0), Inches(0.28),
    font_size=9, color=GREY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Cluster 1: Foundational Architecture (Transformer)
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Cluster 1: The Transformer — Foundation of All STT Models")
add_slide_number(slide, 8)

# Paper citation
add_textbox(slide,
    "P6  |  Vaswani et al. (2017)  |  NeurIPS  |  arXiv:1706.03762",
    Inches(0.4), Inches(0.85), Inches(12.0), Inches(0.35),
    font_size=11, color=STEEL_BLUE, italic=True)

# Left bullets
left_b = [
    "Multi-head self-attention replaces recurrence entirely",
    "Encoder-decoder architecture with positional encodings",
    "O(1) maximum path length → long-range dependency learning",
    "Parallelisable training: days (RNN) → hours (Transformer)",
    "Foundation for: Whisper (ASR), NLLB (MT), Conformer (ASR), HuBERT (SSL), Qwen2.5 (ISUM)",
]
add_bullets_box(slide, left_b,
    Inches(0.4), Inches(1.3), Inches(5.9), Inches(4.0),
    font_size=13, fill_color=LIGHT_BLUE, line_color=NAVY)

# Right: Architecture diagram
arch_left = Inches(6.6)
arch_top  = Inches(1.3)

# ENCODER box
enc_rect = add_rounded_rect(slide, arch_left, arch_top, Inches(2.9), Inches(2.8),
                              fill_color=LIGHT_BLUE, line_color=NAVY)
tb_enc = slide.shapes.add_textbox(arch_left + Inches(0.1), arch_top + Inches(0.05),
                                   Inches(2.7), Inches(2.7))
tf_enc = tb_enc.text_frame; tf_enc.word_wrap = True

enc_lines = [
    ("ENCODER", True, 12, NAVY),
    ("Multi-Head Self-Attention", False, 10, BLACK),
    ("Add & Norm", False, 10, BLACK),
    ("Feed Forward Network", False, 10, BLACK),
    ("Add & Norm", False, 10, BLACK),
    ("× N layers", True, 10, STEEL_BLUE),
]
for i, (txt, b, sz, clr) in enumerate(enc_lines):
    p = tf_enc.paragraphs[0] if i == 0 else tf_enc.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(2)
    r = p.add_run(); r.text = txt
    r.font.name = "Calibri"; r.font.size = Pt(sz)
    r.font.bold = b; r.font.color.rgb = clr

# DECODER box
dec_rect = add_rounded_rect(slide, arch_left + Inches(3.1), arch_top, Inches(2.9), Inches(2.8),
                              fill_color=RGBColor(0xFF, 0xF8, 0xE7), line_color=NAVY)
tb_dec = slide.shapes.add_textbox(arch_left + Inches(3.2), arch_top + Inches(0.05),
                                   Inches(2.7), Inches(2.7))
tf_dec = tb_dec.text_frame; tf_dec.word_wrap = True
dec_lines = [
    ("DECODER", True, 12, NAVY),
    ("Masked Multi-Head Attention", False, 10, BLACK),
    ("Add & Norm", False, 10, BLACK),
    ("Cross-Attention ← Encoder", False, 10, STEEL_BLUE),
    ("Feed Forward Network", False, 10, BLACK),
    ("Add & Norm", False, 10, BLACK),
    ("× N layers", True, 10, STEEL_BLUE),
]
for i, (txt, b, sz, clr) in enumerate(dec_lines):
    p = tf_dec.paragraphs[0] if i == 0 else tf_dec.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(2)
    r = p.add_run(); r.text = txt
    r.font.name = "Calibri"; r.font.size = Pt(sz)
    r.font.bold = b; r.font.color.rgb = clr

# Arrow between encoder & decoder
add_textbox(slide, "→",
    arch_left + Inches(2.9), arch_top + Inches(1.1),
    Inches(0.2), Inches(0.5),
    font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Params label
add_textbox(slide, "Base: 65M params  |  Large: 213M params",
    arch_left, arch_top + Inches(2.95), Inches(6.0), Inches(0.35),
    font_size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# Bottom note
note_r = add_rounded_rect(slide, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.55),
                            fill_color=RGBColor(0xFF, 0xF8, 0xE7), line_color=GOLD, line_width_pt=1.5)
add_textbox(slide, "Every neural model in VANI is a variant of this architecture",
    Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.45),
    font_size=12, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Cluster 2: STT Evolution Timeline
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Cluster 2: STT Architecture Evolution")
add_slide_number(slide, 9)

timeline_items = [
    {
        "year": "2016",
        "model": "Deep Speech 2",
        "paper": "P14",
        "arch": "RNN-CTC",
        "wer": "~8% / 21%",
        "multi": "No",
        "noise": "Low",
        "cpu": "Yes (slow)",
        "color": RGBColor(0xDE, 0xEA, 0xF1),
    },
    {
        "year": "2020",
        "model": "wav2vec 2.0",
        "paper": "P11",
        "arch": "CNN+Transformer\n(Self-supervised)",
        "wer": "1.8% / 3.3%",
        "multi": "Per-lang FT",
        "noise": "High",
        "cpu": "Limited",
        "color": RGBColor(0xE2, 0xEF, 0xDA),
    },
    {
        "year": "2020",
        "model": "Conformer",
        "paper": "P12",
        "arch": "Conv+Attention\n(Supervised)",
        "wer": "1.9% / 3.9%",
        "multi": "Per-lang FT",
        "noise": "Medium-High",
        "cpu": "Limited",
        "color": RGBColor(0xFF, 0xF0, 0xCC),
    },
    {
        "year": "2021",
        "model": "HuBERT",
        "paper": "P15",
        "arch": "Masked Pred.\n(Self-supervised)",
        "wer": "2.0% / 4.0%",
        "multi": "Per-lang FT",
        "noise": "High",
        "cpu": "Limited",
        "color": RGBColor(0xED, 0xE7, 0xF6),
    },
    {
        "year": "2022",
        "model": "Whisper",
        "paper": "P1",
        "arch": "Enc-Dec Transformer\n(Weak Supervised)",
        "wer": "2.7% / 5.2%",
        "multi": "Yes (99 langs)",
        "noise": "Medium-High",
        "cpu": "Yes (int8)",
        "color": RGBColor(0xFF, 0xF8, 0xE7),
        "highlight": True,
    },
]

node_w = Inches(2.35)
node_h = Inches(4.8)
node_top = Inches(1.1)
gap_x   = Inches(0.18)

for i, item in enumerate(timeline_items):
    n_left = Inches(0.35) + i * (node_w + gap_x)
    is_hl = item.get("highlight", False)
    border = GOLD if is_hl else NAVY
    bw = 2.0 if is_hl else 1.0
    rect = add_rounded_rect(slide, n_left, node_top, node_w, node_h,
                             fill_color=item["color"], line_color=border, line_width_pt=bw)

    tb = slide.shapes.add_textbox(n_left + Inches(0.08), node_top + Inches(0.06),
                                   node_w - Inches(0.16), node_h - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = True

    rows = [
        (item["year"], True, 20, NAVY),
        (item["model"], True, 13, NAVY if not is_hl else RGBColor(0x7B,0x5C,0x00)),
        (f"[{item['paper']}]", False, 10, GREY),
        ("", False, 4, BLACK),
        (item["arch"], False, 10, BLACK),
        ("", False, 4, BLACK),
        (f"WER (clean/other):", True, 9, STEEL_BLUE),
        (item["wer"], False, 10, BLACK),
        ("", False, 4, BLACK),
        (f"Multilingual: {item['multi']}", False, 9, BLACK),
        (f"Noise: {item['noise']}", False, 9, BLACK),
        (f"CPU: {item['cpu']}", False, 9, BLACK),
    ]
    for j, (txt, b, sz, clr) in enumerate(rows):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        r.font.name = "Calibri"; r.font.size = Pt(sz)
        r.font.bold = b; r.font.color.rgb = clr

    if i < len(timeline_items) - 1:
        add_textbox(slide, "→",
            n_left + node_w, node_top + Inches(2.2),
            Inches(0.18), Inches(0.4),
            font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

if is_hl:
    add_textbox(slide, "VANI Selected ▲",
        Inches(0.35) + 4*(node_w + gap_x), node_top + node_h + Inches(0.05),
        node_w, Inches(0.3),
        font_size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_textbox(slide,
    "Trajectory: From labelled RNN → self-supervised transformers → weakly supervised multilingual",
    Inches(0.35), Inches(6.1), Inches(12.5), Inches(0.35),
    font_size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — STT Model Comparison Table
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "STT Model Comparison: Architecture & Performance")
add_slide_number(slide, 10)

headers = ["Model", "Architecture", "Training Paradigm",
           "LibriSpeech WER\n(clean/other)", "Multilingual",
           "Noise Robustness", "CPU Feasible"]

rows_data = [
    ["Deep Speech 2",  "Bi-RNN + CTC",           "Supervised\n(12K hrs)",          "~8% / 21%",   "No",             "Low",         "Yes (slow)"],
    ["wav2vec 2.0",    "CNN + Transformer",       "Self-supervised",                "1.8% / 3.3%", "No\n(per-lang FT)", "High",     "Limited"],
    ["Conformer",      "Conv + Attention",        "Supervised",                     "1.9% / 3.9%", "No\n(per-lang FT)", "Medium-High","Limited"],
    ["HuBERT",         "Transformer\n(masked)",   "Self-supervised",                "2.0% / 4.0%", "No\n(per-lang FT)", "High",     "Limited"],
    ["Whisper\n(large-v3-turbo)", "Enc-Dec Transformer", "Weak Supervised\n(680K hrs)", "2.7% / 5.2%", "Yes\n(99 langs)", "Medium-High", "Yes (int8)"],
    ["IndicWav2Vec",   "Wav2Vec + CTC",           "SSL + Supervised FT",            "Pa: 22.3%\nHi: 18.7%", "9 Indic\nlangs", "Medium", "Yes"],
]

col_widths = [Inches(1.6), Inches(1.7), Inches(1.9), Inches(1.5), Inches(1.3), Inches(1.5), Inches(1.3)]
tbl_left   = Inches(0.35)
tbl_top    = Inches(0.9)
tbl_width  = sum(col_widths)
tbl_height = Inches(6.2)

table = slide.shapes.add_table(
    len(rows_data) + 1, len(headers),
    tbl_left, tbl_top, tbl_width, tbl_height
).table

# Set column widths
for ci, cw in enumerate(col_widths):
    table.columns[ci].width = cw

# Header row
for ci, hdr in enumerate(headers):
    cell = table.cell(0, ci)
    set_cell_bg(cell, NAVY)
    set_cell_text(cell, hdr, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Data rows
for ri, row in enumerate(rows_data):
    for ci, val in enumerate(row):
        cell = table.cell(ri + 1, ci)
        is_whisper = (ri == 4)
        if is_whisper:
            set_cell_bg(cell, LIGHT_GOLD)
            set_cell_text(cell, val, font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        else:
            bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
            set_cell_bg(cell, bg)
            set_cell_text(cell, val, font_size=10, bold=False, color=BLACK, align=PP_ALIGN.CENTER)

# Whisper label
add_textbox(slide, "◄ VANI Selected",
    tbl_left + tbl_width + Inches(0.05),
    tbl_top + Inches(0.9) + 4 * (tbl_height / 7),
    Inches(1.2), Inches(0.35),
    font_size=9, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Cluster 3: Noise Robustness
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Cluster 3: Noise Robustness — Augmentation & Enhancement")
add_slide_number(slide, 11)

# LEFT — SpecAugment
add_rect(slide, Inches(0.3), Inches(0.9), Inches(5.9), Inches(5.3),
         fill_color=LIGHT_BLUE, line_color=NAVY)

add_textbox(slide, "SpecAugment — P13, Park et al. 2019",
    Inches(0.4), Inches(0.95), Inches(5.7), Inches(0.45),
    font_size=13, bold=True, color=NAVY)

add_textbox(slide, '"Train the ASR to be noise-robust"',
    Inches(0.4), Inches(1.42), Inches(5.7), Inches(0.35),
    font_size=11, italic=True, color=STEEL_BLUE)

spec_bullets = [
    "Frequency Masking: mask F consecutive mel-frequency channels → simulates channel dropout/fading",
    "Time Masking: mask T consecutive time steps → simulates squelch, burst noise, dropout",
    "Time Warping: non-linear time deformation → simulates Doppler, timing errors",
    "Result: WER improvement LibriSpeech 6.8% → 5.8% (with LM)",
    "Applied in: Whisper training | Domain fine-tuning for VANI",
]
txb = slide.shapes.add_textbox(Inches(0.45), Inches(1.85), Inches(5.7), Inches(4.0))
tf = txb.text_frame; tf.word_wrap = True
for i, b in enumerate(spec_bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    r = p.add_run(); r.text = f"• {b}"
    r.font.name = "Calibri"; r.font.size = Pt(11.5)
    r.font.color.rgb = BLACK

# RIGHT — SEGAN
add_rect(slide, Inches(6.6), Inches(0.9), Inches(6.35), Inches(5.3),
         fill_color=RGBColor(0xFF, 0xF0, 0xCC), line_color=RGBColor(0xFF, 0x8C, 0x00))

add_textbox(slide, "SEGAN — P16, Pascual et al. 2017",
    Inches(6.7), Inches(0.95), Inches(6.1), Inches(0.45),
    font_size=13, bold=True, color=RGBColor(0xAA, 0x50, 0x00))

add_textbox(slide, '"Enhance the audio before ASR"',
    Inches(6.7), Inches(1.42), Inches(6.1), Inches(0.35),
    font_size=11, italic=True, color=RGBColor(0xAA, 0x50, 0x00))

segan_bullets = [
    "GAN-based end-to-end waveform enhancement",
    "Generator: encoder-decoder with skip connections",
    "Discriminator: distinguishes enhanced vs. real clean speech",
    "PESQ: 1.97 (noisy) → 2.16 (enhanced)",
    "Pipeline: DeepFilterNet (~7 MB) + Whisper → estimated 10-20% WER gain at SNR <5dB",
]
txb2 = slide.shapes.add_textbox(Inches(6.7), Inches(1.85), Inches(6.1), Inches(4.0))
tf2 = txb2.text_frame; tf2.word_wrap = True
for i, b in enumerate(segan_bullets):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_before = Pt(5)
    r = p.add_run(); r.text = f"• {b}"
    r.font.name = "Calibri"; r.font.size = Pt(11.5)
    r.font.color.rgb = BLACK

# Bottom decision box
dec_rect = add_rounded_rect(slide, Inches(0.3), Inches(6.4), Inches(12.65), Inches(0.8),
                              fill_color=LIGHT_BLUE, line_color=NAVY, line_width_pt=1.5)
add_textbox(slide,
    "Architecture Decision: Noise-robust end-to-end ASR  vs.  Enhance-then-Recognise  vs.  Joint Training",
    Inches(0.4), Inches(6.48), Inches(12.5), Inches(0.65),
    font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Cluster 3: Radio Noise SNR Analysis
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Radio Noise Profile & STT Model Performance")
add_slide_number(slide, 12)

snr_headers = ["SNR Level", "Noise Type", "Whisper WER\n(estimated)",
               "With Enhancement", "Recommended Approach"]
snr_rows = [
    ("> 15dB",     "Mild AWGN",           "< 15%",      "No benefit",          "Whisper direct"),
    ("10 – 15dB",  "Moderate",            "15 – 25%",   "Marginal",            "Whisper direct"),
    ("5 – 10dB",   "Significant fading",  "25 – 40%",   "+10% improvement",    "Whisper + SpecAugment"),
    ("0 – 5dB",    "Severe",              "40 – 60%",   "+15 – 20%",           "DeepFilterNet + Whisper"),
    ("< 0dB",      "Extreme squelch",     "> 60%",      "Limited",             "VAD gating + enhance"),
]
snr_colors = [
    RGBColor(0xE2, 0xEF, 0xDA),
    RGBColor(0xFF, 0xFF, 0xCC),
    RGBColor(0xFF, 0xE0, 0xB2),
    RGBColor(0xFF, 0xCC, 0xBC),
    RGBColor(0xFF, 0xCC, 0xCC),
]

snr_col_w = [Inches(1.4), Inches(2.2), Inches(1.8), Inches(2.0), Inches(3.6)]
snr_left  = Inches(0.5)
snr_top   = Inches(0.95)
snr_h     = Inches(5.5)

snr_table = slide.shapes.add_table(
    len(snr_rows) + 1, len(snr_headers),
    snr_left, snr_top, sum(snr_col_w), snr_h
).table

for ci, cw in enumerate(snr_col_w):
    snr_table.columns[ci].width = cw

for ci, hdr in enumerate(snr_headers):
    cell = snr_table.cell(0, ci)
    set_cell_bg(cell, NAVY)
    set_cell_text(cell, hdr, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(snr_rows):
    for ci, val in enumerate(row):
        cell = snr_table.cell(ri + 1, ci)
        set_cell_bg(cell, snr_colors[ri])
        bold = (ci == 0)
        set_cell_text(cell, val, font_size=12, bold=bold, color=BLACK, align=PP_ALIGN.CENTER)

# Bottom note
note_r = add_rounded_rect(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.65),
                            fill_color=LIGHT_RED, line_color=RED, line_width_pt=1.5)
add_textbox(slide,
    "No published benchmark exists for Indic STT at these SNR levels — a key research gap this study addresses",
    Inches(0.6), Inches(6.67), Inches(12.1), Inches(0.55),
    font_size=11, italic=True, color=RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Cluster 4: Multilingual & Indic ASR
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Cluster 4: Multilingual & Indic Language ASR")
add_slide_number(slide, 13)

# LEFT — IndicWav2Vec
add_rect(slide, Inches(0.3), Inches(0.9), Inches(5.6), Inches(4.8),
         fill_color=LIGHT_BLUE, line_color=NAVY)
add_textbox(slide, "P7 — IndicWav2Vec (Javed et al. 2022)",
    Inches(0.45), Inches(0.97), Inches(5.3), Inches(0.45),
    font_size=12, bold=True, color=NAVY)
indic_b = [
    "9 Indic languages: Hi, Mr, Gu, Te, Ta, Kn, Or, Pa, Bn",
    "17,000 hrs unlabelled Indian speech pre-training",
    "WER baselines: Punjabi 22.3%, Hindi 18.7%",
    "Specialist fine-tuned → domain-specific advantage",
    "Limitation: requires labelled data, no zero-shot",
]
txb = slide.shapes.add_textbox(Inches(0.45), Inches(1.48), Inches(5.3), Inches(4.0))
tf = txb.text_frame; tf.word_wrap = True
for i, b in enumerate(indic_b):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    r = p.add_run(); r.text = f"• {b}"
    r.font.name = "Calibri"; r.font.size = Pt(12)
    r.font.color.rgb = BLACK

# RIGHT — MMS
add_rect(slide, Inches(7.4), Inches(0.9), Inches(5.6), Inches(4.8),
         fill_color=RGBColor(0xE2, 0xEF, 0xDA), line_color=RGBColor(0x37, 0x5C, 0x23))
add_textbox(slide, "P4 — MMS (Pratap et al. 2023)",
    Inches(7.55), Inches(0.97), Inches(5.3), Inches(0.45),
    font_size=12, bold=True, color=RGBColor(0x37, 0x5C, 0x23))
mms_b = [
    "1,000+ languages  |  256-language audio LangID",
    "Built on wav2vec 2.0  |  150 MB compact model",
    "> 90% LangID accuracy from raw audio",
    "Solves Whisper's Punjabi/Hindi confusion",
    "In VANI: Vote 3 in 3-way LangID ensemble",
]
txb2 = slide.shapes.add_textbox(Inches(7.55), Inches(1.48), Inches(5.3), Inches(4.0))
tf2 = txb2.text_frame; tf2.word_wrap = True
for i, b in enumerate(mms_b):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_before = Pt(5)
    r = p.add_run(); r.text = f"• {b}"
    r.font.name = "Calibri"; r.font.size = Pt(12)
    r.font.color.rgb = BLACK

# Centre comparison box
add_rect(slide, Inches(0.3), Inches(5.85), Inches(12.7), Inches(1.35),
         fill_color=RGBColor(0xFF, 0xF8, 0xE7), line_color=GOLD, )
add_textbox(slide,
    "VANI Decision: Whisper (zero-shot multilingual) + MMS-LID (audio LangID) over IndicWav2Vec (fine-tuned specialist)\n"
    "— justified by absence of labelled military-domain training data",
    Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.15),
    font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — 3-Way LangID Ensemble
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "VANI's 3-Way Language Identification Ensemble")
add_slide_number(slide, 14)

# Diagram using text boxes and shapes
def flow_box(slide, text, left, top, w, h, fill, border, font_sz=11, bold=False, clr=NAVY):
    add_rounded_rect(slide, left, top, w, h, fill_color=fill, line_color=border, line_width_pt=1.5)
    tb = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.04),
                                   w - Inches(0.16), h - Inches(0.06))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = "Calibri"; r.font.size = Pt(font_sz)
    r.font.bold = bold; r.font.color.rgb = clr

def arrow_down(slide, left, top, h=0.25):
    add_textbox(slide, "▼", left, Inches(top), Inches(0.3), Inches(h),
                font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

def arrow_right_txt(slide, left, top):
    add_textbox(slide, "→", Inches(left), Inches(top), Inches(0.35), Inches(0.35),
                font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Audio input
flow_box(slide, "Audio Input", Inches(5.5), Inches(0.95), Inches(2.3), Inches(0.5),
         NAVY, NAVY, font_sz=13, bold=True, clr=WHITE)
arrow_down(slide, Inches(6.3), 1.5)

# Three vote branches
branch_tops_y = 1.85
branch_data = [
    (Inches(1.0), "Whisper ASR", LIGHT_BLUE, NAVY,
     "Language Probability\n(Vote 1 — weight: high)"),
    (Inches(5.2), "FastText lid.176.bin\n(from ASR transcript)", LIGHT_BLUE, NAVY,
     "Text LangID\n(Vote 2 — weight: medium)"),
    (Inches(9.4), "MMS-LID-256\n(from raw audio)", LIGHT_BLUE, NAVY,
     "Audio LangID\n(Vote 3 — weight: medium)"),
]
for bleft, bt, bf, bb, label in branch_data:
    flow_box(slide, bt, bleft, Inches(branch_tops_y), Inches(3.1), Inches(0.65),
             bf, bb, font_sz=11, bold=True)
    flow_box(slide, label, bleft, Inches(branch_tops_y + 0.72), Inches(3.1), Inches(0.65),
             WHITE, STEEL_BLUE, font_sz=10)

# Voting engine
arrow_down(slide, Inches(6.3), branch_tops_y + 1.45)
flow_box(slide, "Confidence-Weighted Voting Engine",
         Inches(3.2), Inches(branch_tops_y + 1.72), Inches(6.5), Inches(0.55),
         NAVY, NAVY, font_sz=12, bold=True, clr=WHITE)

# Voting rules
vote_rules = [
    "Unanimous → boost ×1.10",
    "Majority → average of agreers",
    "All disagree → best single model",
]
txb_r = slide.shapes.add_textbox(Inches(3.2), Inches(branch_tops_y + 2.35), Inches(6.5), Inches(0.9))
tf_r = txb_r.text_frame; tf_r.word_wrap = True
for i, v in enumerate(vote_rules):
    p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"{'  ' if i else ''}• {v}"
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.color.rgb = BLACK

# Final output
arrow_down(slide, Inches(6.3), branch_tops_y + 3.3)
flow_box(slide, "Final Language + Confidence Score",
         Inches(3.5), Inches(branch_tops_y + 3.58), Inches(6.0), Inches(0.55),
         LIGHT_BLUE, NAVY, font_sz=12, bold=True)

# Flags
flags_txb = slide.shapes.add_textbox(Inches(0.4), Inches(branch_tops_y + 4.3), Inches(12.5), Inches(0.55))
tf_f = flags_txb.text_frame; tf_f.word_wrap = True
p_f = tf_f.paragraphs[0]; p_f.alignment = PP_ALIGN.CENTER
r_f = p_f.add_run()
r_f.text = ("< 0.60 confidence → flagged UNCERTAIN  |  "
            "Punjabi override: FastText+MMS pa + Whisper hi → force pa via NLLB")
r_f.font.name = "Calibri"; r_f.font.size = Pt(10.5)
r_f.font.italic = True; r_f.font.color.rgb = STEEL_BLUE

# Papers cited
add_textbox(slide,
    "Papers cited: P1 (Whisper), P5 (FastText), P4 (MMS-LID)   |   "
    '"This ensemble resolves >90% of language ambiguity cases on tested intercepts"',
    Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.28),
    font_size=9, color=GREY, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Cluster 5: Downstream NLP Pipeline
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Cluster 5: Downstream Intelligence Extraction")
add_slide_number(slide, 15)

stages = [
    {
        "num": "Stage 1",
        "label": "TRANSLATION",
        "body": "P2 NLLB-200 (19 languages → English)\nP3 IndicTrans2 (Dogri fallback)\nCoverage: 19 operational languages  |  +44% BLEU over prior best",
        "fill": LIGHT_BLUE,
        "border": NAVY,
    },
    {
        "num": "Stage 2",
        "label": "KEYWORD DETECTION (current)",
        "body": "P5 FastText character n-grams  |  Regex-based threat keyword matching\nLimitation: language-specific dictionaries, misses paraphrase",
        "fill": RGBColor(0xFF, 0xF0, 0xCC),
        "border": RGBColor(0xFF, 0x8C, 0x00),
    },
    {
        "num": "Stage 3",
        "label": "KEYWORD DETECTION (planned — P8 XLM-R)",
        "body": "Zero-shot cross-lingual NER trained on English military annotations\nZero-shot transfer to Hindi, Punjabi, Urdu, Pashto via cross-lingual representations",
        "fill": RGBColor(0xED, 0xE7, 0xF6),
        "border": RGBColor(0x6A, 0x1B, 0x9A),
    },
    {
        "num": "Stage 4",
        "label": "ISUM GENERATION",
        "body": "P9 PEGASUS: Gap Sentence Generation → 5W structured summary\nCurrent: Qwen2.5 rule-based  |  Planned: Fine-tuned on 1000+ annotated intercepts",
        "fill": LIGHT_BLUE,
        "border": NAVY,
    },
    {
        "num": "Stage 5",
        "label": "SPEAKER DIARISATION — P10 pyannote (Phase 5)",
        "body": '"Who spoke when" → per-speaker ISUM  |  Voiceprint tracking across intercepts',
        "fill": LIGHT_GREEN,
        "border": RGBColor(0x37, 0x5C, 0x23),
    },
]

stage_h = Inches(1.12)
stage_top_start = Inches(0.95)
stage_gap = Inches(0.07)

for i, st in enumerate(stages):
    top = stage_top_start + i * (stage_h + stage_gap)
    add_rounded_rect(slide, Inches(0.3), top, Inches(12.7), stage_h,
                     fill_color=st["fill"], line_color=st["border"], line_width_pt=1.5)
    tb = slide.shapes.add_textbox(Inches(0.45), top + Inches(0.06),
                                   Inches(12.4), stage_h - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = True
    # Stage header line
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r_num = p1.add_run(); r_num.text = f"{st['num']} — {st['label']}"
    r_num.font.name = "Calibri"; r_num.font.size = Pt(12); r_num.font.bold = True
    r_num.font.color.rgb = st["border"]
    # Body
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT; p2.space_before = Pt(2)
    r2 = p2.add_run(); r2.text = st["body"]
    r2.font.name = "Calibri"; r2.font.size = Pt(11); r2.font.color.rgb = BLACK

    if i < len(stages) - 1:
        arrow_down(slide, Inches(6.5), (top + stage_h) / 914400)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Research Gaps
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Key Research Gaps Identified")
add_slide_number(slide, 16)

gaps = [
    "No public benchmark exists comparing STT models specifically on tactical radio noise (multi-path fading, squelch, codec compression) for Indic languages at controlled SNR levels",
    "IndicWav2Vec and Conformer have not been evaluated on military radio-domain audio — existing benchmarks use broadcast/read speech",
    "Speaker diarisation for multi-party Indic radio intercepts is absent from all reviewed systems — who-said-what is operationally critical",
    "Regex-based keyword detection (VANI current) is brittle to paraphrase and cross-script code-switching — XLM-R NER offers a research-backed solution but requires domain annotation data",
    "LLM-based ISUM generation (Qwen2.5) is unconstrained without domain fine-tuning — PEGASUS fine-tuning requires 1000+ annotated transcript→5W pairs not yet collected",
]

gap_h = Inches(1.1)
gap_top_start = Inches(0.95)
gap_gap = Inches(0.1)

for i, gap_text in enumerate(gaps):
    top = gap_top_start + i * (gap_h + gap_gap)
    add_rounded_rect(slide, Inches(0.3), top, Inches(12.7), gap_h,
                     fill_color=LIGHT_RED, line_color=RED, line_width_pt=1.5)
    tb = slide.shapes.add_textbox(Inches(0.45), top + Inches(0.1),
                                   Inches(12.4), gap_h - Inches(0.15))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r_n = p.add_run(); r_n.text = f"Gap {i+1}:  "
    r_n.font.name = "Calibri"; r_n.font.size = Pt(12); r_n.font.bold = True
    r_n.font.color.rgb = RED
    r_t = p.add_run(); r_t.text = gap_text
    r_t.font.name = "Calibri"; r_t.font.size = Pt(11.5)
    r_t.font.color.rgb = BLACK

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Expected Findings
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Expected Findings")
add_slide_number(slide, 17)

findings = [
    "Whisper large-v3-turbo is the most practical multilingual noisy radio STT for CPU-constrained deployment — zero-shot, int8 quantised, RTF <3× on 8 GB RAM",
    "Conformer achieves lowest WER on moderately noisy audio but requires per-language fine-tuning — unsuitable for zero-shot multilingual deployment",
    "HuBERT/wav2vec 2.0 show strongest robustness on unseen noise types due to self-supervised noise-invariant representations",
    "Speech enhancement (DeepFilterNet) as preprocessing improves Whisper WER by 10-20% at SNR <5dB at cost of ~0.3× additional RTF",
    "3-way LangID ensemble (Whisper + FastText + MMS-LID) achieves >90% accuracy — significantly better than any single model alone",
    "SpecAugment with radio-specific parameters (frequency masking F=20-40 bins, time masking T=50-100ms) is the recommended augmentation policy for domain fine-tuning",
]

box_w = Inches(6.1)
box_h = Inches(1.15)
gap_h = Inches(0.1)
cols = [Inches(0.3), Inches(6.55)]
col_items = [findings[:3], findings[3:]]

for ci, (col_left, items) in enumerate(zip(cols, col_items)):
    for ri, finding in enumerate(items):
        top = Inches(0.95) + ri * (box_h + gap_h)
        add_rounded_rect(slide, col_left, top, box_w, box_h,
                         fill_color=LIGHT_BLUE, line_color=NAVY, line_width_pt=1.0)
        tb = slide.shapes.add_textbox(col_left + Inches(0.1), top + Inches(0.07),
                                       box_w - Inches(0.2), box_h - Inches(0.1))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r_n = p.add_run(); r_n.text = f"{ci*3 + ri + 1}.  "
        r_n.font.name = "Calibri"; r_n.font.size = Pt(12); r_n.font.bold = True
        r_n.font.color.rgb = NAVY
        r_t = p.add_run(); r_t.text = finding
        r_t.font.name = "Calibri"; r_t.font.size = Pt(11)
        r_t.font.color.rgb = BLACK

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Recommendations
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "Recommendations")
add_slide_number(slide, 18)

recs = [
    ("[1]  Establish Radio-Noise ASR Benchmark",
     "Collect real VHF/UHF recordings at controlled SNR levels (-5 to +20dB) across 5-7 Indic languages. "
     "No such public benchmark currently exists."),
    ("[2]  Evaluate DeepFilterNet as Preprocessing",
     "Lightweight (<10MB, real-time CPU) enhancement before Whisper for SNR <10dB conditions. "
     "Expected 10-20% WER improvement on severely degraded intercepts."),
    ("[3]  Integrate pyannote.audio (Phase 5)",
     'Speaker diarisation for per-operator ISUM generation. Enable voiceprint tracking across multiple intercepts '
     '— high-value intelligence capability.'),
    ("[4]  Domain Fine-tuning Path",
     "Collect 500+ annotated intercept samples per language via VANI annotation system → "
     "IndicWav2Vec domain fine-tuning → Whisper LoRA fine-tuning"),
    ("[5]  Ablation Study Design",
     "Evaluate: Whisper alone vs. DeepFilterNet+Whisper vs. Conformer (fine-tuned) on held-out radio eval set "
     "across SNR: {-5, 0, 5, 10, 15, 20} dB × 5 Indic languages"),
]

rec_h = Inches(1.12)
rec_gap = Inches(0.09)
rec_top_start = Inches(0.92)

for i, (title, body) in enumerate(recs):
    top = rec_top_start + i * (rec_h + rec_gap)
    add_rounded_rect(slide, Inches(0.3), top, Inches(12.7), rec_h,
                     fill_color=LIGHT_GREEN, line_color=RGBColor(0x37, 0x5C, 0x23), line_width_pt=1.5)
    tb = slide.shapes.add_textbox(Inches(0.45), top + Inches(0.08),
                                   Inches(12.4), rec_h - Inches(0.12))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r_t = p.add_run(); r_t.text = f"→  {title}  "
    r_t.font.name = "Calibri"; r_t.font.size = Pt(12); r_t.font.bold = True
    r_t.font.color.rgb = GREEN_DARK
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT; p2.space_before = Pt(1)
    r2 = p2.add_run(); r2.text = f"     {body}"
    r2.font.name = "Calibri"; r2.font.size = Pt(11)
    r2.font.color.rgb = BLACK

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — References
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_line(slide, "References")
add_slide_number(slide, 19)

refs_col1 = [
    "P1:  Radford, A., et al. (2022). Robust speech recognition via large-scale weak supervision. arXiv:2212.04356",
    "P2:  Costa-jussà, M. R., et al. (2022). No language left behind. arXiv:2207.04672",
    "P3:  Gala, J., et al. (2023). IndicTrans2. arXiv:2305.16307",
    "P4:  Pratap, V., et al. (2023). Scaling speech technology to 1,000+ languages. arXiv:2305.13516",
    "P5:  Joulin, A., et al. (2017). Bag of tricks for efficient text classification. arXiv:1607.01759",
    "P6:  Vaswani, A., et al. (2017). Attention is all you need. arXiv:1706.03762",
    "P7:  Javed, T., et al. (2022). IndicWav2Vec. arXiv:2111.03945",
    "P8:  Conneau, A., et al. (2020). Unsupervised cross-lingual representation learning at scale. arXiv:1911.02116",
]
refs_col2 = [
    "P9:  Zhang, J., et al. (2020). PEGASUS. arXiv:1912.08777",
    "P10: Bredin, H., et al. (2020). pyannote.audio. arXiv:2001.01980",
    "P11: Baevski, A., et al. (2020). wav2vec 2.0. arXiv:2006.11477  ★",
    "P12: Gulati, A., et al. (2020). Conformer. arXiv:2005.08100  ★",
    "P13: Park, D. S., et al. (2019). SpecAugment. arXiv:1904.08779  ★",
    "P14: Amodei, D., et al. (2016). Deep Speech 2. arXiv:1512.02595  ★",
    "P15: Hsu, W. N., et al. (2021). HuBERT. arXiv:2106.07447  ★",
    "P16: Pascual, S., et al. (2017). SEGAN. arXiv:1703.09452  ★",
]

def make_ref_col(slide, refs, left, top):
    txb = slide.shapes.add_textbox(left, top, Inches(6.2), Inches(6.2))
    tf = txb.text_frame; tf.word_wrap = True
    for i, ref in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run(); r.text = ref
        r.font.name = "Calibri"; r.font.size = Pt(9.5)
        is_new = "★" in ref
        r.font.color.rgb = RGBColor(0x7B, 0x5C, 0x00) if is_new else BLACK
        r.font.bold = is_new

make_ref_col(slide, refs_col1, Inches(0.35), Inches(0.9))
make_ref_col(slide, refs_col2, Inches(6.8), Inches(0.9))

add_textbox(slide, "★ = New papers added in this review (P11-P16)",
    Inches(0.35), Inches(7.1), Inches(6.0), Inches(0.28),
    font_size=9, color=GREY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Thank You / Closing
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()

# Navy background
add_rect(slide, 0, 0, W, H, NAVY)

# Gold accent bar
add_rect(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.06), GOLD)

# Summary title
add_textbox(slide, "Summary",
    Inches(0.6), Inches(0.85), Inches(12.1), Inches(0.65),
    font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Three key takeaways
takeaways = [
    ("1.", "Whisper + 3-way LangID ensemble is the current best-practice for offline multilingual radio STT on CPU-constrained hardware"),
    ("2.", "Six new papers (P11-P16) establish the comparative framework for evaluating STT models on noisy radio conditions — a benchmark that does not yet exist"),
    ("3.", "VANI implements this pipeline and provides the annotation infrastructure for future domain fine-tuning"),
]

for i, (num, txt) in enumerate(takeaways):
    top = Inches(1.8) + i * Inches(1.35)
    add_rounded_rect(slide, Inches(0.8), top, Inches(11.7), Inches(1.15),
                     fill_color=RGBColor(0x28, 0x44, 0x72),
                     line_color=GOLD, line_width_pt=1.5)
    tb = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.12),
                                   Inches(11.3), Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = f"{num}  "
    r1.font.name = "Calibri"; r1.font.size = Pt(14); r1.font.bold = True
    r1.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = txt
    r2.font.name = "Calibri"; r2.font.size = Pt(13)
    r2.font.color.rgb = WHITE

# Bottom line
add_rect(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.04), GOLD)

add_textbox(slide,
    "VANI – Voice Analysis & Neural Intelligence  |  SOATE-44  |  March 2026",
    Inches(0.6), Inches(6.72), Inches(12.1), Inches(0.4),
    font_size=12, color=GOLD, align=PP_ALIGN.CENTER)

add_textbox(slide, "Questions?",
    Inches(0.6), Inches(7.08), Inches(12.1), Inches(0.35),
    font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out_path = "/Users/vik/offline_ai_system_v2/VANI_LRP_Comparative_STT.pptx"
prs.save(out_path)
print(f"SUCCESS: Saved to {out_path}")
