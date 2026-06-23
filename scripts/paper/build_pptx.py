"""
build_pptx.py — Generate VANI_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x2a, 0x4a)   # dark navy
GREEN  = RGBColor(0x00, 0x7a, 0x3d)   # military green
GOLD   = RGBColor(0xd4, 0xaf, 0x37)   # gold accent
WHITE  = RGBColor(0xff, 0xff, 0xff)
LGRAY  = RGBColor(0xf0, 0xf4, 0xf8)
RED    = RGBColor(0xcc, 0x00, 0x00)
AMBER  = RGBColor(0xff, 0x99, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ─────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color):
    """Fill slide background with solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height,
        text="", font_size=18, bold=False, color=WHITE,
        bg_color=None, align=PP_ALIGN.LEFT, italic=False,
        wrap=True, font_name="Calibri"):
    """Add a text box and return the text frame."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return tf

def rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a filled rectangle (no text)."""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_title_bar(slide, title, subtitle=None):
    """Standard navy title bar across top."""
    rect(slide, 0, 0, 13.33, 1.1, NAVY)
    box(slide, 0.3, 0.1, 10, 0.6, title,
        font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        box(slide, 0.3, 0.68, 12, 0.35, subtitle,
            font_size=14, color=GOLD, align=PP_ALIGN.LEFT)

def add_bullet_block(slide, left, top, width, height, items,
                     bullet="▸", font_size=16, color=WHITE, line_gap=0.36):
    """Add a list of bullet items."""
    y = top
    for item in items:
        text = f"{bullet}  {item}" if bullet else item
        box(slide, left, y, width, line_gap + 0.05,
            text, font_size=font_size, color=color)
        y += line_gap

def slide_number(slide, n, total=15):
    box(slide, 12.3, 7.1, 1.0, 0.3,
        f"{n} / {total}", font_size=11, color=GOLD, align=PP_ALIGN.RIGHT)

def green_bar(slide, top=1.1, height=0.06):
    rect(slide, 0, top, 13.33, height, GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)

# Gold diagonal accent
rect(s, 0, 0, 0.45, 7.5, GREEN)

box(s, 1.2, 0.8, 11, 1.2, "V A N I",
    font_size=72, bold=True, color=GOLD, align=PP_ALIGN.LEFT, font_name="Calibri")
box(s, 1.2, 2.0, 11, 0.7, "Voice Analysis & Neural Intelligence",
    font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
box(s, 1.2, 2.7, 11, 0.55,
    "Military-Grade Offline Radio Intercept Analysis System",
    font_size=20, italic=True, color=LGRAY, align=PP_ALIGN.LEFT)
rect(s, 1.2, 3.35, 9.5, 0.05, GOLD)
box(s, 1.2, 3.5, 11, 0.45,
    "Interim Presentation  |  March 2026",
    font_size=18, color=GOLD, align=PP_ALIGN.LEFT)

box(s, 1.2, 6.5, 11, 0.4,
    "CONFIDENTIAL — FOR INTERNAL USE ONLY",
    font_size=12, italic=True, color=RGBColor(0xaa, 0xaa, 0xaa),
    align=PP_ALIGN.LEFT)
slide_number(s, 1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AIM / SCOPE
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "AIM & SCOPE", "What VANI is designed to do")
green_bar(s)

box(s, 0.4, 1.3, 12.5, 0.45,
    "AIM", font_size=16, bold=True, color=GOLD)
box(s, 0.4, 1.7, 12.5, 0.9,
    "Develop a fully offline, CPU-deployable intelligence system that automatically transcribes, "
    "translates, and analyses military radio intercepts in multiple languages — producing structured "
    "intelligence summaries (ISUM) in real time, with zero internet dependency.",
    font_size=15, color=WHITE)

box(s, 0.4, 2.65, 12.5, 0.4, "SCOPE", font_size=16, bold=True, color=GOLD)

scope = [
    ("Languages",       "Hindi · Punjabi · Urdu · Pashto · Nepali · Dogri · Kashmiri · Bengali · Maithili · Sindhi · Sinhala · Mandarin · Burmese · Tibetan · English"),
    ("Input Formats",   "WAV · MP3 · OGG · FLAC · M4A"),
    ("Output",          "5W ISUM report · PDF/DOCX export · SQLite history"),
    ("Deployment",      "Standalone Windows/Linux · Air-gapped · No internet required"),
    ("Hardware",        "8 GB RAM · CPU-only (no GPU required)"),
    ("Interface",       "Web-based GUI (Streamlit) accessible via browser"),
]
y = 3.1
for label, val in scope:
    rect(s, 0.4, y, 2.5, 0.33, RGBColor(0x00, 0x4a, 0x2a))
    box(s, 0.45, y+0.03, 2.45, 0.3, label, font_size=12, bold=True, color=GOLD)
    box(s, 3.0, y+0.03, 10.0, 0.35, val, font_size=12, color=WHITE)
    y += 0.39

slide_number(s, 2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM BEING SOLVED
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "PROBLEM BEING SOLVED", "Why VANI is needed")
green_bar(s)

problems = [
    "Manual transcription of radio intercepts is TIME-CONSUMING and ERROR-PRONE",
    "Existing tools require CLOUD CONNECTIVITY — unacceptable in field deployments",
    "Most ASR/NLP tools do NOT support Indian border-area languages (Punjabi, Dogri, Pashto, Kashmiri)",
    "Raw intercepts provide NO STRUCTURED OUTPUT — analysts must extract 5W fields manually",
    "No existing system combines ASR + Language ID + Translation + Keywords + ISUM offline",
]
y = 1.4
for i, prob in enumerate(problems):
    rect(s, 0.4, y, 0.55, 0.5, RED)
    box(s, 0.42, y+0.06, 0.5, 0.38, f"P{i+1}", font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.1, y+0.05, 11.8, 0.45, prob, font_size=15, color=WHITE)
    y += 0.62

rect(s, 0.4, 4.7, 12.5, 0.9, RGBColor(0x00, 0x4a, 0x2a))
box(s, 0.6, 4.75, 12.2, 0.8,
    "VANI addresses all five gaps in a single integrated, offline, multi-language system "
    "deployable on standard 8 GB RAM hardware with no specialist IT infrastructure.",
    font_size=15, color=WHITE)

slide_number(s, 3)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "SYSTEM ARCHITECTURE", "Three-layer design: Pipeline → Models → Output")
green_bar(s)

# Layer 1: Input
rect(s, 0.3, 1.3, 2.0, 0.7, RGBColor(0x1a, 0x4a, 0x7a))
box(s, 0.3, 1.3, 2.0, 0.7, "AUDIO INPUT\nWAV / MP3 / FLAC",
    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow
box(s, 2.35, 1.55, 0.5, 0.3, "▶", font_size=18, color=GOLD, align=PP_ALIGN.CENTER)

# Layer 2: Pipeline bar
rect(s, 2.8, 1.3, 10.2, 0.7, RGBColor(0x00, 0x50, 0x30))
box(s, 2.8, 1.3, 10.2, 0.7,
    "PROCESSING PIPELINE:  VAD → Preprocess → Chunk → ASR → LangID → Translate → Keywords → ISUM",
    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Models
box(s, 0.3, 2.2, 12.7, 0.35, "MODEL LAYER  (all fully offline)", font_size=14, bold=True, color=GOLD)
models = [
    ("Whisper\nLarge-v3-Turbo\n(ASR)", NAVY),
    ("FastText\n176-lang\n(LangID)", RGBColor(0x1a, 0x4a, 0x7a)),
    ("MMS-LID\n256-lang\n(LangID)", RGBColor(0x1a, 0x4a, 0x7a)),
    ("NLLB-200\n600M\n(Translation)", RGBColor(0x00, 0x4a, 0x2a)),
    ("IndicTrans2\n1B\n(Dogri)", RGBColor(0x00, 0x4a, 0x2a)),
    ("Qwen2.5\n1.5B ISUM\n(GPU opt.)", RGBColor(0x4a, 0x1a, 0x1a)),
]
x = 0.3
for label, col in models:
    rect(s, x, 2.6, 2.05, 1.1, col)
    box(s, x, 2.6, 2.05, 1.1, label, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += 2.18

# Output
box(s, 0.3, 3.9, 12.7, 0.35, "OUTPUT LAYER", font_size=14, bold=True, color=GOLD)
outputs = [
    ("SQLite Database\n(History & Search)", NAVY),
    ("JSON Result\n(Pipeline output)", RGBColor(0x1a, 0x4a, 0x7a)),
    ("PDF / DOCX\nReport Export", RGBColor(0x00, 0x4a, 0x2a)),
]
x = 0.3
for label, col in outputs:
    rect(s, x, 4.25, 4.0, 0.85, col)
    box(s, x, 4.25, 4.0, 0.85, label, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += 4.2

# UI Tabs
box(s, 0.3, 5.25, 12.7, 0.35, "STREAMLIT WEB INTERFACE — 8 Tabs", font_size=14, bold=True, color=GOLD)
tabs = ["[P] PROCESS", "[I] ISUM REPORT", "[S] SEARCH", "[D] DASHBOARD",
        "[H] HISTORY", "[X] EXPORT", "[M] METRICS", "[A] ANNOTATE"]
x = 0.3
for tab in tabs:
    rect(s, x, 5.6, 1.55, 0.55, RGBColor(0x1a, 0x2a, 0x5a))
    box(s, x, 5.6, 1.55, 0.55, tab, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += 1.6

slide_number(s, 4)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PIPELINE WORKFLOW
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "PIPELINE WORKFLOW", "10-stage processing — audio in, intelligence out")
green_bar(s)

stages = [
    ("1\nVAD", "Silero VAD 4.0\nRemove silence & noise", RGBColor(0x1a, 0x4a, 0x7a)),
    ("2\nPREPROC", "librosa 16kHz\nNoise reduction", RGBColor(0x1a, 0x4a, 0x7a)),
    ("3\nCHUNK", "VAD-aware split\nmax 29s chunks", RGBColor(0x1a, 0x4a, 0x7a)),
    ("4\nASR", "Whisper Large-v3\nCTranslate2 int8", RGBColor(0x00, 0x3a, 0x7a)),
    ("5\nLANG ID", "3-way vote\nWhisper+FT+MMS", RGBColor(0x00, 0x3a, 0x7a)),
    ("6\nTRANSLATE", "NLLB-600M\n+ chrF scoring", RGBColor(0x00, 0x4a, 0x2a)),
    ("7\nKEYWORDS", "804 keywords\n12 categories", RGBColor(0x00, 0x4a, 0x2a)),
    ("8\nISUM", "WHO/WHERE/WHEN\n+ Assessment", RGBColor(0x4a, 0x2a, 0x00)),
]

x = 0.2
for num_lbl, detail, col in stages:
    rect(s, x, 1.35, 1.55, 0.75, col)
    box(s, x, 1.35, 1.55, 0.75, num_lbl, font_size=13, bold=True,
        color=GOLD, align=PP_ALIGN.CENTER)
    box(s, x, 2.12, 1.55, 0.55, detail, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
    if x < 11.8:
        box(s, x+1.57, 1.6, 0.25, 0.4, "▶", font_size=16, color=GOLD, align=PP_ALIGN.CENTER)
    x += 1.63

# Output branch
rect(s, 0.4, 2.85, 12.5, 0.06, GREEN)
box(s, 0.3, 2.95, 12.7, 0.35,
    "OUTPUT BRANCH", font_size=13, bold=True, color=GOLD)

outs = [
    ("SQLite DB\nHistory + Search", RGBColor(0x1a, 0x2a, 0x5a)),
    ("PDF / DOCX\nReport", RGBColor(0x00, 0x4a, 0x2a)),
    ("JSON\nPipeline result", RGBColor(0x1a, 0x4a, 0x7a)),
    ("Metrics\n(RTF, WER, BLEU)", RGBColor(0x3a, 0x3a, 0x00)),
    ("Annotation\n(Training data)", RGBColor(0x4a, 0x1a, 0x1a)),
]
x = 0.4
for lbl, col in outs:
    rect(s, x, 3.35, 2.4, 0.85, col)
    box(s, x, 3.35, 2.4, 0.85, lbl, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += 2.55

# Performance note
rect(s, 0.4, 4.4, 12.5, 0.9, RGBColor(0x00, 0x30, 0x1a))
box(s, 0.6, 4.45, 12.2, 0.8,
    "Performance:  First run (cold) ~60–70 s for 5 s audio (model loading).  "
    "Subsequent runs (models cached via @st.cache_resource):  ~15–25 s.",
    font_size=14, color=WHITE)

slide_number(s, 5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LANGUAGE IDENTIFICATION VOTING
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "3-WAY LANGUAGE ID VOTING", "No single model is trusted alone")
green_bar(s)

# Three model boxes
models_info = [
    ("Whisper", "ASR Vote\nlang_prob", RGBColor(0x00, 0x3a, 0x7a)),
    ("FastText", "Text Vote\n176 languages", RGBColor(0x00, 0x4a, 0x2a)),
    ("MMS-LID", "Audio Vote\n256 languages", RGBColor(0x4a, 0x2a, 0x00)),
]
x = 1.0
for name, detail, col in models_info:
    rect(s, x, 1.35, 3.3, 1.0, col)
    box(s, x, 1.35, 3.3, 0.55, name, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, x, 1.88, 3.3, 0.45, detail, font_size=13, color=LGRAY, align=PP_ALIGN.CENTER)
    x += 3.65

# Converge arrow
box(s, 6.2, 2.4, 1.0, 0.5, "▼", font_size=24, color=GOLD, align=PP_ALIGN.CENTER)

# Voting box
rect(s, 2.5, 2.9, 8.3, 0.8, RGBColor(0x1a, 0x2a, 0x5a))
box(s, 2.5, 2.9, 8.3, 0.8,
    "CONFIDENCE-WEIGHTED VOTE\nUnanimous (3/3) → boost ×1.10  |  Majority (2/3) → avg conf  |  No consensus → penalty ×0.85",
    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

# Three outcomes
outcomes = [
    ("UNANIMOUS\n3/3 agree\nconf × 1.10", RGBColor(0x00, 0x4a, 0x2a)),
    ("MAJORITY\n2/3 agree\navg conf", RGBColor(0x00, 0x3a, 0x7a)),
    ("NO CONSENSUS\n1/3 or 0/3\nconf × 0.85", RGBColor(0x4a, 0x1a, 0x00)),
]
x = 0.8
for lbl, col in outcomes:
    rect(s, x, 3.85, 3.6, 0.95, col)
    box(s, x, 3.85, 3.6, 0.95, lbl, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += 3.9

# Punjabi fix
rect(s, 0.8, 4.95, 11.7, 0.75, RGBColor(0x4a, 0x3a, 0x00))
box(s, 0.9, 4.98, 11.5, 0.68,
    "PUNJABI DISAMBIGUATION:  Gurmukhi script detected  OR  (FastText=pa OR MMS=pa)  AND  Whisper=hi  "
    "→  FORCE language = Punjabi (pa)  →  route via NLLB-200",
    font_size=13, color=WHITE)

# Final output
rect(s, 2.5, 5.85, 8.3, 0.65, RGBColor(0x00, 0x4a, 0x2a))
box(s, 2.5, 5.85, 8.3, 0.65,
    "FINAL: language code + conf + route (nllb / indictrans2 / none)  |  conf < 0.60 → UNCERTAIN flag",
    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(s, 6)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DELIVERABLES
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "DELIVERABLES", "Seven defined outputs against the project mandate")
green_bar(s)

deliverables = [
    ("D1", "Core Processing Pipeline",          "End-to-end offline pipeline — VAD→ASR→LangID→Translate→Keywords→ISUM. 15 languages. No internet.",    "✅ COMPLETE",  GREEN),
    ("D2", "Web Intelligence Interface",         "8-tab Streamlit UI: Process · ISUM · Search · Dashboard · History · Export · Metrics · Annotate",    "✅ COMPLETE",  GREEN),
    ("D3", "ISUM Report Generation",             "Automated 5W extraction · 12-category threat classification · PDF/DOCX export",                        "✅ COMPLETE",  GREEN),
    ("D4", "Quality Metrics Framework",          "Tier 1 auto (RTF/confidence/5W) · Tier 2 reference (WER/CER/BLEU/chrF/TER)",                          "✅ COMPLETE",  GREEN),
    ("D5", "Training Data Collection",           "Analyst annotation tab · corrected ASR/translation/ISUM export · structured JSON for fine-tuning",     "✅ COMPLETE",  GREEN),
    ("D6", "Windows Server Deployment Package",  "Installation guide · dependency pinning · model download scripts · air-gapped validation",             "🔄 IN PROGRESS", AMBER),
    ("D7", "Fine-Tuned Domain Models",           "Whisper/NLLB/Qwen fine-tuned on actual military intercepts on offline Windows Server",                "📅 FUTURE",    RGBColor(0x1a, 0x4a, 0x7a)),
]

y = 1.3
for tag, title, desc, status, col in deliverables:
    rect(s, 0.3, y, 0.7, 0.5, col)
    box(s, 0.3, y, 0.7, 0.5, tag, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.1, y+0.02, 4.0, 0.25, title, font_size=13, bold=True, color=GOLD)
    box(s, 1.1, y+0.24, 8.8, 0.24, desc, font_size=11, color=LGRAY)
    rect(s, 10.05, y+0.05, 3.0, 0.38, col)
    box(s, 10.05, y+0.05, 3.0, 0.38, status, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y += 0.58

slide_number(s, 7)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PHASE-WISE IMPLEMENTATION
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "PHASE-WISE IMPLEMENTATION", "Six phases from foundation to scaled deployment")
green_bar(s)

phases = [
    ("Ph 1\nFOUNDATION", "System architecture · Base pipeline (VAD→ASR→Translate) · Streamlit skeleton · SQLite DB schema", "Nov–Dec 2025", "✅", GREEN),
    ("Ph 2\nLANG ID",    "3-way LangID voting · NLLB-200 integration · IndicTrans2 patches · Punjabi disambiguation · MMS-LID",  "Jan 2026",     "✅", GREEN),
    ("Ph 3\nINTEL",      "804-keyword dictionary (12 cats) · Rule-based ISUM (5W) · Qwen LLM ISUM (GPU opt.) · Threat classification", "Feb 2026", "✅", GREEN),
    ("Ph 4\nOUTPUT",     "PDF/DOCX export · Metrics framework (RTF/WER/BLEU) · Annotation tab · Model caching (10× speedup)", "Feb–Mar 2026", "✅", GREEN),
    ("Ph 5\nTRAIN",      "Windows Server deployment · Offline fine-tuning on actual intercepts · Domain-adaptive Whisper", "Apr–Jun 2026", "🔄", AMBER),
    ("Ph 6\nSCALE",      "GPU-accelerated deployment · Real-time streaming · Multi-intercept batch · SIGINT integration",  "Jul–Sep 2026", "📅", RGBColor(0x1a, 0x4a, 0x7a)),
]

y = 1.3
for phase, detail, timeline, icon, col in phases:
    rect(s, 0.3, y, 1.55, 0.72, col)
    box(s, 0.3, y, 1.55, 0.72, phase, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 2.0, y+0.04, 8.9, 0.3, detail, font_size=11, color=WHITE)
    rect(s, 10.9, y+0.1, 1.5, 0.35, RGBColor(0x0d, 0x1b, 0x2e))
    box(s, 10.9, y+0.08, 1.5, 0.35, timeline, font_size=10, color=GOLD, align=PP_ALIGN.CENTER)
    rect(s, 12.55, y+0.1, 0.55, 0.35, col)
    box(s, 12.55, y+0.1, 0.55, 0.35, icon, font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
    y += 0.82

slide_number(s, 8)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CURRENT STATUS
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "CURRENT STATUS — MARCH 2026", "What is operational today")
green_bar(s)

# Stage table
cols = ["Stage", "Technology", "Performance"]
widths = [1.4, 4.8, 5.8]
headers = [("Stage", NAVY), ("Technology", RGBColor(0x00, 0x3a, 0x7a)), ("Performance", RGBColor(0x00, 0x4a, 0x2a))]

stages_data = [
    ("VAD",          "Silero VAD 4.0",                       "~0.6 s for 5 min audio"),
    ("Preprocess",   "librosa + noisereduce",                 "16 kHz normalised"),
    ("Chunking",     "VAD-aware, max 29 s",                   "~0.2 s"),
    ("ASR",          "Whisper Large-v3-Turbo · CTranslate2 int8", "~40 s cold / ~5 s cached (CPU)"),
    ("Language ID",  "3-way vote: Whisper + FastText + MMS",  "< 1 s (cached)"),
    ("Translation",  "NLLB-200-distilled-600M",               "~2–5 s"),
    ("Keywords",     "804 entries · 12 categories",           "< 0.1 s"),
    ("ISUM",         "Rule-based (default) / Qwen LLM (GPU)", "< 0.1 s (rule-based)"),
]

x_starts = [0.3, 1.75, 6.6]
y = 1.3
# header row
for i, (h, col) in enumerate(headers):
    rect(s, x_starts[i], y, widths[i]-0.05, 0.38, col)
    box(s, x_starts[i], y, widths[i]-0.05, 0.38, h, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
y += 0.38
for row_i, (stage, tech, perf) in enumerate(stages_data):
    row_col = RGBColor(0x12, 0x25, 0x40) if row_i % 2 == 0 else RGBColor(0x0d, 0x1b, 0x2e)
    for i, val in enumerate([stage, tech, perf]):
        rect(s, x_starts[i], y, widths[i]-0.05, 0.38, row_col)
        box(s, x_starts[i]+0.05, y+0.05, widths[i]-0.15, 0.3, val, font_size=11, color=WHITE)
    y += 0.38

# Models installed
box(s, 0.3, 4.6, 12.7, 0.3, "MODELS INSTALLED (ALL OFFLINE)", font_size=13, bold=True, color=GOLD)
models_data = [
    ("Whisper Large-v3-Turbo", "~1.6 GB"),
    ("NLLB-200-distilled-600M", "~2.2 GB"),
    ("IndicTrans2-1B", "~3.5 GB"),
    ("FastText lid.176.bin", "~900 MB"),
    ("MMS-LID-256", "~150 MB"),
    ("Qwen2.5-1.5B", "~3.1 GB"),
    ("TOTAL", "~11.5 GB"),
]
x = 0.3
for name, size in models_data:
    col = RGBColor(0x4a, 0x1a, 0x00) if name == "TOTAL" else RGBColor(0x00, 0x30, 0x1a)
    rect(s, x, 4.95, 1.8, 0.65, col)
    box(s, x, 4.95, 1.8, 0.35, name, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, x, 5.3, 1.8, 0.28, size, font_size=11, color=GOLD, align=PP_ALIGN.CENTER)
    x += 1.88

slide_number(s, 9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — KEY FEATURES
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "KEY FEATURES", "Eight distinguishing capabilities")
green_bar(s)

features = [
    ("1", "Air-Gapped Operation",         "All models run locally. HF_HUB_OFFLINE=1 enforced. Zero internet dependency. Field deployable."),
    ("2", "Multi-Language Radio Optimised","Custom ASR prompts with military terminology. VAD threshold tuned for radio static (0.45)."),
    ("3", "3-Way Language Voting",         "Whisper + FastText + MMS-LID ensemble with confidence weighting and vote_note in every result."),
    ("4", "Tactical ISUM Extraction",     "WHO: callsigns/ranks/units · WHERE: MGRS/terrain · WHEN: H-hour/military time · Category-aware assessment."),
    ("5", "12-Category Threat Intelligence","804 keywords · CRITICAL: attack/explosives/CBRN · HIGH: weapons/movement/casualties · 8 languages."),
    ("6", "Quality Metrics Framework",    "RTF · ASR confidence · Model agreement · 5W completeness · Back-translation chrF · WER/CER/BLEU."),
    ("7", "Analyst Annotation System",    "Correct ASR/translation/ISUM per intercept. Export structured JSON for Whisper/NLLB/Qwen fine-tuning."),
    ("8", "Professional Report Export",   "ReportLab PDF + python-docx DOCX, military-format, threat-colour-coded, metrics tables included."),
]

y = 1.3
for i, (num, title, desc) in enumerate(features):
    col = GREEN if i < 4 else RGBColor(0x1a, 0x4a, 0x7a)
    rect(s, 0.3, y, 0.55, 0.55, col)
    box(s, 0.3, y, 0.55, 0.55, num, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.0, y+0.02, 3.2, 0.25, title, font_size=13, bold=True, color=GOLD)
    box(s, 1.0, y+0.26, 11.9, 0.26, desc, font_size=11, color=WHITE)
    if i == 3:
        y += 0.65
        rect(s, 0, y-0.05, 13.33, 0.04, GREEN)
        y += 0.1
    else:
        y += 0.64

slide_number(s, 10)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — LITERATURE REVIEW
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "LITERATURE REVIEW", "Five foundational papers underpinning VANI")
green_bar(s)

papers = [
    ("P1", "Robust Speech Recognition via Large-Scale Weak Supervision",
     "Radford et al. (OpenAI) — arXiv:2212.04356, 2022",
     "Whisper architecture — encoder-decoder transformer on 680K hrs multilingual data. Informs ASR backbone, initial prompt injection, and hallucination prevention."),
    ("P2", "No Language Left Behind: Scaling Human-Centered MT",
     "Costa-jussà et al. (Meta AI) — arXiv:2207.04672, 2022",
     "NLLB-200 — 200-language single-model MT. NLLB-200-distilled-600M fits 8 GB RAM and handles all Indic languages in VANI."),
    ("P3", "IndicTrans2: Towards High-Quality MT for all 22 Scheduled Indian Languages",
     "Gala et al. (AI4Bharat) — TMLR 2023",
     "SotA Indic translation. Used for Dogri (not in NLLB-200). Script-detection logic from this paper informs Punjabi disambiguation."),
    ("P4", "Scaling Speech Technology to 1,000+ Languages",
     "Pratap et al. (Meta AI) — arXiv:2305.13516, 2023",
     "MMS-LID-256 audio-based language ID. Third vote in VANI ensemble — critical for romanised Indic scripts where text models fail."),
    ("P5", "Bag of Tricks for Efficient Text Classification",
     "Joulin et al. (Facebook AI) — EACL 2017 | arXiv:1607.01759",
     "FastText sub-word n-gram classification. lid.176.bin provides text-based language ID on Whisper transcript — second vote in ensemble."),
]

y = 1.3
for tag, title, ref, relevance in papers:
    rect(s, 0.3, y, 0.65, 0.9, RGBColor(0x00, 0x4a, 0x2a))
    box(s, 0.3, y+0.2, 0.65, 0.5, tag, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.1, y+0.02, 11.9, 0.28, title, font_size=12, bold=True, color=GOLD)
    box(s, 1.1, y+0.3, 11.9, 0.22, ref, font_size=10, italic=True, color=LGRAY)
    box(s, 1.1, y+0.53, 11.9, 0.32, f"Relevance: {relevance}", font_size=10, color=WHITE)
    y += 1.02

slide_number(s, 11)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PAPER PUBLICATION
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "PAPER PUBLICATION DETAILS", "Planned academic dissemination")
green_bar(s)

box(s, 0.4, 1.35, 12.5, 0.38,
    "Planned Title:  \"VANI: An Offline Multilingual Spoken Intelligence Analysis System for Low-Resource Indian Border Languages\"",
    font_size=13, italic=True, color=GOLD)

venues = [
    ("PRIMARY", "IEEE Transactions on Information Forensics and Security",
     "Scope: Intelligence systems, signal processing, security applications\nIF: ~7.2 | Covers SIGINT/NLP intersection directly"),
    ("ALT 1",   "INTERSPEECH 2026 / ICASSP 2026",
     "Scope: Speech processing, multilingual ASR, low-resource languages\nStrong match for multilingual ASR + Language ID contribution"),
    ("ALT 2",   "ACL 2026 / EMNLP 2026",
     "Scope: Natural Language Processing\nMatch: multilingual NLP pipeline, ISUM generation, annotation framework"),
]
y = 1.9
for tag, venue, detail in venues:
    col = RGBColor(0x00, 0x4a, 0x2a) if tag == "PRIMARY" else RGBColor(0x00, 0x3a, 0x7a)
    rect(s, 0.4, y, 0.9, 0.8, col)
    box(s, 0.4, y, 0.9, 0.8, tag, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.4, y+0.02, 11.5, 0.3, venue, font_size=13, bold=True, color=WHITE)
    box(s, 1.4, y+0.34, 11.5, 0.44, detail, font_size=11, color=LGRAY)
    y += 0.98

box(s, 0.4, 5.05, 12.5, 0.3, "KEY CONTRIBUTIONS FOR PUBLICATION", font_size=13, bold=True, color=GOLD)
contribs = [
    "3-way confidence-weighted Language ID voting combining ASR + text + audio models with Punjabi script disambiguation",
    "Rule-based military ISUM generation with domain-specific tactical pattern extraction (WHO / WHERE / WHEN)",
    "Indic-language translation routing solving IndicTrans2 / NLLB-200 complementarity for 15 languages",
    "End-to-end offline deployment of 11.5 GB model stack within 8 GB RAM CPU constraint",
    "Military annotation framework for building labelled intercept training data",
]
y = 5.4
for i, c in enumerate(contribs):
    box(s, 0.6, y, 12.5, 0.26, f"C{i+1}.  {c}", font_size=11, color=WHITE)
    y += 0.28

slide_number(s, 12)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — HARDWARE & SOFTWARE
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "HARDWARE & SOFTWARE PROCUREMENT", "Current stack and target deployment specifications")
green_bar(s)

# Software table
box(s, 0.3, 1.3, 6.3, 0.32, "CURRENT SOFTWARE STACK", font_size=13, bold=True, color=GOLD)
sw = [
    ("Deep Learning", "PyTorch", "2.2.2 CPU"),
    ("ASR",           "faster-whisper + CTranslate2", "Latest"),
    ("Translation",   "HuggingFace Transformers", "4.40.0"),
    ("Language ID",   "FastText", "0.9.2"),
    ("UI",            "Streamlit", "1.34.0"),
    ("Database",      "SQLite3", "Native Python"),
    ("PDF Export",    "ReportLab", "4.0.7"),
    ("Audio",         "librosa + soundfile + noisereduce", "Pinned"),
]
sw_headers = [("Category", 1.8), ("Package", 3.2), ("Version", 1.3)]
x_sw = [0.3, 2.15, 5.4]
y = 1.65
for i, (h, w) in enumerate(sw_headers):
    rect(s, x_sw[i], y, w, 0.3, NAVY)
    box(s, x_sw[i], y, w, 0.3, h, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
y += 0.3
for ri, (cat, pkg, ver) in enumerate(sw):
    rc = RGBColor(0x12, 0x25, 0x40) if ri % 2 == 0 else RGBColor(0x0d, 0x1b, 0x2e)
    for i, val in enumerate([cat, pkg, ver]):
        widths_sw = [1.8, 3.2, 1.3]
        rect(s, x_sw[i], y, widths_sw[i], 0.3, rc)
        box(s, x_sw[i]+0.05, y+0.04, widths_sw[i]-0.1, 0.23, val, font_size=10, color=WHITE)
    y += 0.3

# Target hardware
box(s, 7.1, 1.3, 6.0, 0.32, "TARGET DEPLOYMENT (Windows Server)", font_size=13, bold=True, color=GOLD)
hw = [
    ("CPU",      "8-core+ x86-64",          "Parallel model loading"),
    ("RAM",      "16–32 GB recommended",     "Multiple models cached"),
    ("Storage",  "100 GB SSD",              "Models + DB + audio archive"),
    ("GPU",      "NVIDIA 16 GB+ VRAM",       "Phase 5: Qwen + GPU ASR"),
    ("OS",       "Windows Server 2019/2022", "Standalone deployment"),
    ("Network",  "Air-gapped",              "Operational security"),
]
hw_headers = [("Component", 1.5), ("Spec", 2.3), ("Reason", 2.3)]
x_hw = [7.1, 8.65, 11.0]
y = 1.65
for i, (h, w) in enumerate(hw_headers):
    rect(s, x_hw[i], y, w, 0.3, NAVY)
    box(s, x_hw[i], y, w, 0.3, h, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
y += 0.3
for ri, (comp, spec, reason) in enumerate(hw):
    rc = RGBColor(0x12, 0x25, 0x40) if ri % 2 == 0 else RGBColor(0x0d, 0x1b, 0x2e)
    for i, val in enumerate([comp, spec, reason]):
        widths_hw = [1.5, 2.3, 2.3]
        rect(s, x_hw[i], y, widths_hw[i], 0.3, rc)
        box(s, x_hw[i]+0.05, y+0.04, widths_hw[i]-0.1, 0.23, val, font_size=10, color=WHITE)
    y += 0.3

# Procurement status
rect(s, 0.3, 5.0, 12.7, 0.75, RGBColor(0x00, 0x30, 0x1a))
proc = [
    ("✅", "All software open-source — no licensing cost"),
    ("✅", "All models downloaded and operational offline (~11.5 GB)"),
    ("🔄", "Windows Server hardware — procurement in progress"),
    ("📅", "GPU (Phase 5 training) — to be requested after server delivery"),
]
box(s, 0.5, 5.05, 12.3, 0.24, "PROCUREMENT STATUS:", font_size=12, bold=True, color=GOLD)
y2 = 5.3
for icon, text in proc:
    box(s, 0.5, y2, 12.3, 0.2, f"{icon}  {text}", font_size=11, color=WHITE)
    y2 += 0.22

slide_number(s, 13)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — FUTURE IMPROVEMENTS
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x0d, 0x1b, 0x2e))
add_title_bar(s, "FUTURE IMPROVEMENTS", "Roadmap beyond the current operational system")
green_bar(s)

# Three columns: Short / Medium / Long term
terms = [
    ("SHORT-TERM\n0–3 months", RGBColor(0x00, 0x4a, 0x2a), [
        "Offline fine-tuning on actual military intercepts",
        "Whisper domain adaptation (WER < 10% on tactical vocab)",
        "NLLB fine-tune for Punjabi/Urdu (chrF > 55)",
        "Qwen SFT for structured ISUM JSON generation",
        "Real-time streaming mode (< 10 s latency)",
        "GPU-accelerated deployment (5× speedup)",
    ]),
    ("MEDIUM-TERM\n3–6 months", RGBColor(0x1a, 0x4a, 0x7a), [
        "Speaker diarisation (multi-party intercepts)",
        "Auto-escalation: CRITICAL + conf > 0.90 → alert",
        "Expanded language coverage (Balochi, Tibetan dialects)",
        "Voice biometrics — identify known speaker patterns",
        "Confidence-based priority queue for analyst review",
        "Batch multi-intercept processing",
    ]),
    ("LONG-TERM\n6–12 months", RGBColor(0x4a, 0x1a, 0x00), [
        "Multi-intercept entity correlation (callsigns, locations)",
        "Timeline reconstruction across intercept corpus",
        "Network graph of communication participants",
        "LLM ISUM with full operational context window",
        "NATO STANAG output format",
        "GIS mapping integration for location visualisation",
    ]),
]

x = 0.3
for header, col, items in terms:
    rect(s, x, 1.3, 4.2, 0.65, col)
    box(s, x, 1.3, 4.2, 0.65, header, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y = 2.05
    for item in items:
        box(s, x+0.1, y, 4.0, 0.38, f"▸  {item}", font_size=11, color=WHITE)
        y += 0.42
    x += 4.37

slide_number(s, 14)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ACHIEVEMENT SUMMARY
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 0.45, 7.5, GREEN)

box(s, 0.7, 0.2, 12.0, 0.55, "WHAT HAS BEEN ACHIEVED",
    font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
box(s, 0.7, 0.75, 12.0, 0.35, "VANI — Interim Project Review · March 2026",
    font_size=16, color=GOLD, align=PP_ALIGN.LEFT)
rect(s, 0.7, 1.15, 12.0, 0.05, GREEN)

done = [
    "10-stage end-to-end offline pipeline — fully operational",
    "15 languages supported (all major Indian border-area languages)",
    "3-way Language ID voting: Whisper + FastText + MMS-LID",
    "Punjabi/Hindi disambiguation via Gurmukhi script detection + model voting",
    "NLLB-200 + IndicTrans2 dual translation routing (15 languages)",
    "804-keyword multilingual threat detection — 12 categories",
    "Rule-based ISUM: WHO / WHERE / WHEN / WHAT / ASSESSMENT",
    "Qwen2.5 LLM ISUM integration (GPU deployment ready)",
    "8-tab professional Streamlit web interface",
    "Professional PDF + DOCX military report export",
    "Comprehensive metrics framework: Tier 1 (auto) + Tier 2 (reference)",
    "Analyst annotation system for training data collection",
    "SQLite persistent database with full-text search",
    "Model caching — 10× speedup on subsequent pipeline runs",
    "All models installed and operational offline (~11.5 GB total)",
]
pending = [
    "Windows Server deployment package (in progress)",
    "Offline fine-tuning on actual military intercepts (planned)",
    "Real-time streaming mode (planned)",
    "GPU-accelerated deployment (pending hardware)",
]

y = 1.3
x = 0.7
for i, item in enumerate(done):
    if i == 8:
        x = 6.9
        y = 1.3
    rect(s, x, y, 0.4, 0.32, GREEN)
    box(s, x, y, 0.4, 0.32, "✓", font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, x+0.5, y+0.04, 5.8, 0.26, item, font_size=11, color=WHITE)
    y += 0.37

y = max(1.3 + len(done[:8])*0.37, 1.3 + len(done[8:])*0.37) + 1.3
rect(s, 0.7, 5.55, 12.0, 0.04, GOLD)
for i, item in enumerate(pending):
    icon = "🔄" if i == 0 else "📅"
    box(s, 0.7 + i*3.15, 5.65, 3.1, 0.38, f"{icon}  {item}", font_size=10, color=LGRAY)

slide_number(s, 15)

# ── Save ────────────────────────────────────────────────────────────────────
prs.save("VANI_Presentation.pptx")
print("Saved: VANI_Presentation.pptx")
