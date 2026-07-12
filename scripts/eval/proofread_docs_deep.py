"""Deep proofread: FINETUNE_REPORT.md tables, built PPTX, built PDF vs source data."""
import io, sys, json, csv, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pathlib import Path

REPO = Path(r"C:\Users\vis15\offline_ai_system_v2")
sys.path.insert(0, str(REPO))

comp = {r["lang"]: r for r in json.load(open(REPO/"docs/model_comparison_results.json", encoding="utf-8"))}
smft = {r["lang"]: r for r in json.load(open(REPO/"docs/seamless_ft_results.json", encoding="utf-8"))}
pre_comp = {r["lang"]: r for r in json.load(open(REPO/"docs/model_comparison_results_PRE_FIX_2026-07-10.json", encoding="utf-8"))}
pre_smft = {r["lang"]: r for r in json.load(open(REPO/"docs/seamless_ft_results_PRE_FIX_2026-07-10.json", encoding="utf-8"))}
sweep = {}
for r in csv.DictReader(open(REPO/"eval_data/wer_robustness_results.csv", encoding="utf-8")):
    sweep[(r["system"], r["lang"], r["condition"])] = float(r["wer"])

from generate_report_pdf import LANG_META, EVAL_RESULTS

issues, ok = [], []
def check(label, claim, truth, tol=0.051):
    if claim is None:
        issues.append(f"{label}: could not parse claimed value")
    elif truth is None:
        issues.append(f"{label}: claimed {claim} but source has no value")
    elif abs(claim - truth) > tol:
        issues.append(f"{label}: claimed {claim} but source says {truth}")
    else:
        ok.append(label)

def num(cell):
    """Extract first numeric value from a md cell ('**57.39%**', '−20.2 pp', '+39.0')."""
    c = cell.replace("−", "-").replace("**", "")
    m = re.search(r"[-+]?\d+(?:\.\d+)?", c)
    return float(m.group()) if m else None

def num2(cell):
    """Second numeric value in a cell — the CER in 'WER (CER)' cells; None if absent."""
    c = cell.replace("−", "-").replace("**", "")
    m = re.findall(r"[-+]?\d+(?:\.\d+)?", c)
    return float(m[1]) if len(m) > 1 else None

md = (REPO/"docs/FINETUNE_REPORT.md").read_text(encoding="utf-8")

def table_rows(anchor, n_langs=7):
    """Return dict lang -> list of cells for the FIRST md table after `anchor`.
    A row belongs to the table if the line starts with '|'; the table ends at the
    first non-'|' line after data rows started."""
    i = md.index(anchor)
    rows, in_table = {}, False
    for line in md[i:].splitlines():
        is_row = line.lstrip().startswith("|")
        if in_table and not is_row:
            break
        if is_row:
            m = re.search(r"\((pa|ps|ur|ne|zh|hi|ks)\)", line, re.I)
            if m:
                in_table = True
                cells = [c.strip() for c in re.split(r"(?<!\\)\|", line.strip().strip("|"))]
                # drop a leading pure row-number cell ('1'..'7')
                if re.fullmatch(r"\d", cells[0]):
                    cells = cells[1:]
                rows[m.group(1).lower()] = cells
        if len(rows) == n_langs:
            break
    return rows

# ── 1. Cross-model ASR WER table ─────────────────────────────────────────────
rows = table_rows("### ASR Word Error Rate", 7)
for l in ["pa","ps","ur","ne","zh","hi"]:
    c = rows[l]
    check(f"md ASR {l}.baseline", num(c[1]), comp[l]["whisper_baseline_wer"])
    check(f"md ASR {l}.ft",       num(c[2]), comp[l]["whisper_ft_wer"])
    check(f"md ASR {l}.seamless", num(c[3]), comp[l]["seamless_asr_wer"])
    check(f"md ASR {l}.baseline-cer", num2(c[1]), comp[l]["whisper_baseline_cer"])
    check(f"md ASR {l}.ft-cer",       num2(c[2]), comp[l]["whisper_ft_cer"])
    check(f"md ASR {l}.seamless-cer", num2(c[3]), comp[l]["seamless_asr_cer"])
    gain = num(c[4])
    check(f"md ASR {l}.gain-arith", gain, round(comp[l]["whisper_ft_wer"] - comp[l]["whisper_baseline_wer"], 1), tol=0.06)
ksc = rows["ks"]
check("md ASR ks.baseline", num(ksc[1]), 96.87)
check("md ASR ks.ft", num(ksc[2]), 74.02)
check("md ASR ks.gain-arith", num(ksc[4]), round(74.02-96.87, 2), tol=0.06)

# deployed-backend column matches EVAL_RESULTS["best"]
for l in ["pa","ps","ur","ne","zh","hi","ks"]:
    want = "SeamlessM4T" if EVAL_RESULTS[l]["best"] == "seamless" else "FT Whisper"
    got = rows[l][5]
    (ok if want in got else issues).append(
        f"md ASR {l}.backend" if want in got else f"md ASR {l}.backend: says '{got}', EVAL_RESULTS says {want}")

# ── 2. chrF table ─────────────────────────────────────────────────────────────
rows = table_rows("### Translation Quality", 7)
for l in ["pa","ps","ur","ne","zh","hi"]:
    c = rows[l]
    n_, s_ = comp[l]["whisper_nllb_chrf"], comp[l]["seamless_s2tt_chrf"]
    check(f"md chrF {l}.nllb", num(c[1]), n_)
    check(f"md chrF {l}.sm",   num(c[2]), s_)
    want = "SeamlessM4T" if s_ > n_ else "Whisper+NLLB"
    (ok if want in c[3] else issues).append(
        f"md chrF {l}.winner" if want in c[3] else f"md chrF {l}.winner: says '{c[3]}', data says {want}")

# ── 3. Degradation table (ft − sm per condition) ─────────────────────────────
rows = table_rows("### ASR WER under Radio-Channel Degradation", 6)
conds = ["clean", "bandpass", "awgn_10", "awgn_0", "codec_mp3"]
for l in ["pa","ne","hi","ur","zh","ps"]:
    c = rows[l]
    for j, cnd in enumerate(conds):
        true_d = round(sweep[("whisper_ft", l, cnd)] - sweep[("seamless_zs", l, cnd)], 1)
        check(f"md degr {l}/{cnd}", num(c[1+j]), true_d, tol=0.06)

# ── 4. SM4T FT ASR table ──────────────────────────────────────────────────────
rows = table_rows("### ASR WER — Fine-tuned SM4T", 6)
for l in ["pa","ps","ur","ne","zh","hi"]:
    c = rows[l]
    zs, ft = comp[l]["seamless_asr_wer"], smft[l]["sm_ft_asr_wer"]
    check(f"md smft {l}.zs", num(c[1]), zs)
    check(f"md smft {l}.ft", num(c[2]), ft)
    check(f"md smft {l}.delta-arith", num(c[3]), round(ft - zs, 1), tol=0.06)
    check(f"md smft {l}.ftwhisper", num(c[4]), comp[l]["whisper_ft_wer"])

# ── 5. SM4T FT chrF table ────────────────────────────────────────────────────
rows = table_rows("### Translation (S2TT chrF)", 6)
for l in ["pa","ps","ur","ne","zh","hi"]:
    c = rows[l]
    check(f"md smft-chrf {l}.zs",   num(c[1]), comp[l]["seamless_s2tt_chrf"])
    check(f"md smft-chrf {l}.ft",   num(c[2]), smft[l]["sm_ft_s2tt_chrf"])
    check(f"md smft-chrf {l}.nllb", num(c[3]), comp[l]["whisper_nllb_chrf"])

# ── 6. Results Summary table vs LANG_META + JSON ─────────────────────────────
rows = table_rows("## Results Summary", 7)
for l in ["pa","ps","ur","ne","zh","hi","ks"]:
    c = rows[l]
    m = LANG_META[l]
    # cells (row number dropped): Language, Script, Base Model, Dataset, Train Samples, Best Train-Val, Held-Out, Backend
    tv = c[5]; ho = c[6]; bk = c[7]
    check(f"md summary {l}.trainval", num(tv), m["best_wer"])
    step = re.search(r"step (\d+)", tv)
    if step and int(step.group(1)) != m["best_step"]:
        issues.append(f"md summary {l}.beststep: says {step.group(1)}, LANG_META {m['best_step']}")
    else:
        ok.append(f"md summary {l}.beststep")
    check(f"md summary {l}.heldout", num(ho), m["eval_wer"])
    if EVAL_RESULTS[l]["best"] == "seamless":
        mm = re.search(r"\((\d+\.?\d*)%\)", bk)
        check(f"md summary {l}.backend-wer", float(mm.group(1)) if mm else None,
              comp[l]["seamless_asr_wer"])

# ── 7. Corrections section quotes vs PRE_FIX archives ────────────────────────
corr = md[md.index("## Corrections"):md.index("## Fine-Tuning Configuration")]
quotes = {
    "zh old baseline 100.03": (100.03, pre_comp["zh"]["whisper_baseline_wer"]),
    "zh old ft 16.03":        (16.03,  pre_comp["zh"]["whisper_ft_wer"]),
    "zh old smft 60.53":      (60.53,  pre_smft["zh"]["sm_ft_asr_wer"]),
    "zh old sm 100.0":        (100.0,  pre_comp["zh"]["seamless_asr_wer"]),
}
for label, (quoted, archived) in quotes.items():
    (ok if (f"{quoted}" in corr and abs(quoted-archived) < 0.01) else issues).append(
        f"corr {label}" if (f"{quoted}" in corr and abs(quoted-archived) < 0.01)
        else f"corr {label}: quote/archive mismatch (archive {archived})")
# old pa gain −50 pp / new −20 pp
old_pa = pre_comp["pa"]["whisper_ft_wer"] - pre_comp["pa"]["whisper_baseline_wer"]
new_pa = comp["pa"]["whisper_ft_wer"] - comp["pa"]["whisper_baseline_wer"]
(ok if round(old_pa) == -50 and round(new_pa) == -20 else issues).append(
    "corr pa gain shrink −50→−20" if round(old_pa) == -50 and round(new_pa) == -20
    else f"corr pa gain: archive {old_pa:.1f}, corrected {new_pa:.1f}")
# ur −1.4
(ok if round(comp['ur']['whisper_ft_wer']-comp['ur']['whisper_baseline_wer'],1) == -1.4 else issues).append("corr ur −1.4")
# −84 pp headline
(ok if round(pre_comp["zh"]["whisper_ft_wer"] - pre_comp["zh"]["whisper_baseline_wer"], 1) == -84.0 else issues).append("corr −84 pp")

# ── 8. Prose derived claims in md ────────────────────────────────────────────
ne_lead = comp["ne"]["whisper_ft_wer"] - comp["ne"]["seamless_asr_wer"]
(ok if "22.5 pp" in md and abs(ne_lead - 22.46) < 0.01 else issues).append(
    "md ne 22.5pp lead" if "22.5 pp" in md else f"md ne lead: {ne_lead}")
sm_wins = [l for l in ["pa","ps","ur","ne","zh","hi"] if comp[l]["seamless_asr_wer"] < comp[l]["whisper_ft_wer"]]
(ok if len(sm_wins) == 5 and "5 of 6" in md else issues).append(
    "md 5-of-6 claim" if len(sm_wins) == 5 else f"md 5-of-6: actually {sm_wins}")
chrf_sm = [l for l in ["pa","ps","ur","ne","zh","hi"] if comp[l]["seamless_s2tt_chrf"] > comp[l]["whisper_nllb_chrf"]]
(ok if sorted(chrf_sm) == ["ne","pa","zh"] else issues).append(
    "md chrF 3-3 split" if sorted(chrf_sm) == ["ne","pa","zh"] else f"chrF SM4T wins: {chrf_sm}")

# ── 9. Stale-number scan of md outside correction contexts ───────────────────
STALE = ["105.79","94.55","94.23","98.64","30.29","24.44","55.67","49.24",
         "58.72","56.02","54.91","53.42","47.72","41.54","42.85","43.92","2.41"]
for s in STALE:
    for mline in [ln for ln in md.splitlines() if s in ln]:
        issues.append(f"md stale {s}: {mline.strip()[:90]}")
# allowed-in-context stale numbers: report context lines for eyeball
print("Context of intentionally-quoted old figures in md:")
for s in ["100.03","16.03","60.53","103.58"]:
    for mline in [ln for ln in md.splitlines() if s in ln]:
        print(f"  [{s}] {mline.strip()[:110]}")

# ── 10. Built PPTX text vs sources ───────────────────────────────────────────
from pptx import Presentation
prs = Presentation(str(REPO/"docs/VANI_Finetune_Presentation_v6.pptx"))
slides_txt = []
for idx, slide in enumerate(prs.slides, 1):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
    slides_txt.append((idx, "\n".join(parts)))
pptx_all = "\n".join(t for _, t in slides_txt)

for s in STALE:
    for idx, t in slides_txt:
        if s in t:
            line = next(ln for ln in t.splitlines() if s in ln)
            issues.append(f"PPTX slide {idx} stale {s}: {line.strip()[:90]}")
print("\nPPTX occurrences of old zh figures (should be correction notes only):")
for s in ["100.03","16.03","60.53","103.58"]:
    for idx, t in slides_txt:
        for ln in t.splitlines():
            if s in ln:
                print(f"  slide {idx} [{s}] {ln.strip()[:110]}")

EXPECT = {"77.6":"pa base","57.39":"pa ft","19.77":"pa sm","89.76":"ps base","38.55":"ps ft",
          "44.4":"ps sm","21.23":"ur base","19.82":"ur ft","16.9":"ur sm","88.85":"ne base",
          "50.92":"ne ft","28.46":"ne sm","10.99":"zh base","14.22":"zh ft","11.69":"zh sm",
          "26.34":"hi base","19.78":"hi ft","15.44":"hi sm","96.87":"ks base","74.02":"ks ft"}
for s, what in EXPECT.items():
    (ok if s in pptx_all else issues).append(
        f"PPTX has {what}" if s in pptx_all else f"PPTX missing corrected {what} ({s})")

# CER presence: PPTX labels format CER as "(x.y)" (one decimal), PDF/md as "(x.yz)".
CER_KEYS = [("whisper_ft_cer", "ft"), ("seamless_asr_cer", "sm")]
for l in ["pa","ps","ur","ne","zh","hi"]:
    for key, tag in CER_KEYS:
        v = comp[l][key]
        s1 = f"({v:.1f})"
        (ok if s1 in pptx_all else issues).append(
            f"PPTX has {l}.{tag} CER" if s1 in pptx_all
            else f"PPTX missing {l}.{tag} CER {s1}")

# ── 11. Built PDF text vs sources ────────────────────────────────────────────
import pypdf
rd = pypdf.PdfReader(str(REPO/"docs/VANI_Finetune_Report.pdf"))
pdf_all = " ".join((p.extract_text() or "") for p in rd.pages)
for s, what in EXPECT.items():
    (ok if s in pdf_all else issues).append(
        f"PDF has {what}" if s in pdf_all else f"PDF missing corrected {what} ({s})")

for l in ["pa","ps","ur","ne","zh","hi"]:
    for key, tag in [("whisper_baseline_cer","base"), ("whisper_ft_cer","ft"),
                     ("seamless_asr_cer","sm")]:
        v = comp[l][key]
        s2 = f"({v:.2f})"
        (ok if s2 in pdf_all else issues).append(
            f"PDF has {l}.{tag} CER" if s2 in pdf_all
            else f"PDF missing {l}.{tag} CER {s2}")
print("\nPDF occurrences of old zh figures (should be correction notes only):")
for s in ["100.03","16.03","60.53","103.58","105.79"]:
    for i, p in enumerate(rd.pages, 1):
        t = p.extract_text() or ""
        if s in t:
            j = t.index(s)
            print(f"  page {i} [{s}] ...{t[max(0,j-70):j+50].replace(chr(10),' ')}...")

for s in ["94.55","94.23","98.64","30.29","24.44","55.67","49.24","58.72","54.91","53.42","47.72","41.54","42.85","43.92"]:
    if s in pdf_all:
        j = pdf_all.index(s)
        issues.append(f"PDF stale {s}: ...{pdf_all[max(0,j-60):j+40]}...")

# ── 12. Diagram data: LANG_META eval/baseline vs JSON (charts derive from these) ──
for l in ["pa","ps","ur","ne","zh","hi"]:
    check(f"chart-src {l}.baseline", LANG_META[l]["baseline_wer"], comp[l]["whisper_baseline_wer"])
    check(f"chart-src {l}.eval",     LANG_META[l]["eval_wer"],     comp[l]["whisper_ft_wer"])

print(f"\n{'='*70}\nPASSED: {len(ok)} checks")
print(f"ISSUES: {len(issues)}")
for i in issues:
    print(f"  X {i}")
